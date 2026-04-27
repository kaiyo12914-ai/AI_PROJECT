from __future__ import annotations

import io
import os
import re
import json
import copy
import posixpath
import zipfile
import mimetypes
import logging
import platform
import unicodedata
from datetime import datetime
from xml.etree import ElementTree as ET
from urllib.parse import quote
from typing import Optional, List, Dict, Any, Set

from django.conf import settings
from django.contrib import messages
from django.http import HttpResponse, JsonResponse
from django.shortcuts import render, redirect

from webapps.portal.decorators import require_node
from webapps.llm.llm_factory import get_chat_model
from webapps.text2pptx.image_service import ImageGenError, generate_image
from webapps.text2pptx.pptx2schema.extractors.template_structure_extractor import (
    extract_template_structure,
)


# ---------------------------
# Paths / constants
# ---------------------------
BASE_DIR = getattr(settings, "BASE_DIR", None) or os.getcwd()
PPTX_TEMPLATE_DIR = os.path.join(BASE_DIR, "webapps", "text2pptx", "pptx_templates")
GENERATED_IMAGE_DIR = str(
    getattr(
        settings,
        "TEXT2PPTX_IMAGE_DIR",
        os.path.join(settings.MEDIA_ROOT, "generated_images", "text2pptx"),
    )
)

MAX_BULLETS_PER_SLIDE = 7
TITLE_FONT_PT = 34
BODY_FONT_PT = 20
SUBTITLE_FONT_PT = 18
MAX_INPUT_CHARS = int(getattr(settings, "TEXT2PPTX_MAX_CHARS", 20000))
MAX_TEMPLATE_UPLOAD_MB = int(getattr(settings, "TEXT2PPTX_TEMPLATE_MAX_MB", 20))
MAX_SAMPLE_UPLOAD_MB = int(getattr(settings, "TEXT2PPTX_SAMPLE_MAX_MB", 50))
MAX_SAMPLE_SLIDES = int(getattr(settings, "TEXT2PPTX_SAMPLE_MAX_SLIDES", 80))
MAX_SAMPLE_CHARS = int(getattr(settings, "TEXT2PPTX_SAMPLE_MAX_CHARS", 50000))
TEXT2PPTX_IMAGE_MODE = str(getattr(settings, "TEXT2PPTX_IMAGE_MODE", "mock")).strip().lower() or "mock"
TEXT2PPTX_IMAGE_TIMEOUT_SEC = int(getattr(settings, "TEXT2PPTX_IMAGE_TIMEOUT_SEC", 30))
TEXT2PPTX_IMAGE_RETRY = max(0, int(getattr(settings, "TEXT2PPTX_IMAGE_RETRY", 2)))
DEFAULT_TEMPLATE_NAME = "預設範本.pptx"
DEFAULT_MAIN_TITLE = "簡報"
DEFAULT_MAIN_SUBTITLE = "文字內容自動生成簡報"
VALID_SLIDE_TYPES = {"content", "section", "two_content", "comparison"}
VALID_IMAGE_INTENTS = {"concept", "data", "process", "hero"}
MARKER_TO_SLIDE_TYPE = {
    "cover": "cover",
    "section": "section",
    "content": "content",
    "two_content": "two_content",
    "comparison": "comparison",
    "標題投影片": "cover",
    "章節標題": "section",
    "內容頁": "content",
    "雙欄必較頁": "two_content",
}
LAYOUT_NAME_KEYS = {
    "cover": ["title slide", "標題投影片"],
    "content": ["title and content", "內容頁"],
    "section": ["section header", "章節標題"],
    "two_content": ["two content", "雙欄必較頁"],
}
TITLE_PLACEHOLDER_TYPES = {1, 3}
SUBTITLE_PLACEHOLDER_TYPES = {4}
CONTENT_PLACEHOLDER_TYPES = {2, 7}
IMAGE_PLACEHOLDER_TYPES = {11, 14}
PPTX_MAIN_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
POTX_MAIN_CONTENT_TYPE = PPTX_MAIN_CONTENT_TYPE

logger = logging.getLogger(__name__)


def _select_font_names() -> tuple[str, str]:
    forced_ea = (getattr(settings, "TEXT2PPTX_FONT_EAST_ASIAN", "") or "").strip()
    forced_ascii = (getattr(settings, "TEXT2PPTX_FONT_ASCII", "") or "").strip()
    if forced_ea and forced_ascii:
        return forced_ea, forced_ascii

    system = platform.system().lower()
    if "windows" in system:
        return forced_ea or "Microsoft JhengHei", forced_ascii or "Times New Roman"
    if "darwin" in system:
        return forced_ea or "PingFang TC", forced_ascii or "Helvetica Neue"
    return forced_ea or "Noto Sans CJK TC", forced_ascii or "DejaVu Sans"


FONT_NAME_EAST_ASIAN, FONT_NAME_ASCII = _select_font_names()


def _placeholder_type_id(ph) -> int | None:
    try:
        return int(ph.placeholder_format.type)
    except Exception:
        return None


def _is_title_placeholder_type(v: int | None) -> bool:
    return v in TITLE_PLACEHOLDER_TYPES


def _is_subtitle_placeholder_type(v: int | None) -> bool:
    return v in SUBTITLE_PLACEHOLDER_TYPES


def _is_content_placeholder_type(v: int | None) -> bool:
    return v in CONTENT_PLACEHOLDER_TYPES

def _is_image_placeholder_type(v: int | None) -> bool:
    return v in IMAGE_PLACEHOLDER_TYPES


def _list_pptx_templates() -> List[str]:
    if not os.path.isdir(PPTX_TEMPLATE_DIR):
        return []
    files = []
    for fn in os.listdir(PPTX_TEMPLATE_DIR):
        if fn.lower().endswith(".pptx"):
            full = os.path.join(PPTX_TEMPLATE_DIR, fn)
            if os.path.isfile(full):
                files.append(fn)
    files.sort()
    return files


def _normalize_template_name(name: str) -> str:
    return unicodedata.normalize("NFC", str(name or "")).strip()


def _template_name_key(name: str) -> str:
    return _normalize_template_name(name).casefold()


def _list_ignored_template_files() -> List[str]:
    if not os.path.isdir(PPTX_TEMPLATE_DIR):
        return []
    ignored: List[str] = []
    for fn in os.listdir(PPTX_TEMPLATE_DIR):
        full = os.path.join(PPTX_TEMPLATE_DIR, fn)
        if not os.path.isfile(full):
            continue
        lower = fn.lower()
        if lower.endswith(".pptx"):
            continue
        if lower.endswith(".potx") or lower.startswith("~$"):
            ignored.append(fn)
    ignored.sort()
    return ignored


def _safe_select_template(tpl_name: str) -> Optional[str]:
    tpl_name = _normalize_template_name(tpl_name)
    if not tpl_name:
        return None
    wanted_key = _template_name_key(tpl_name)
    for fn in _list_pptx_templates():
        if _template_name_key(fn) == wanted_key:
            return os.path.join(PPTX_TEMPLATE_DIR, fn)
    return None


def _pick_layout_adaptive(prs, layout_type: str):
    # PowerPoint 2016 class mapping:
    # comparison and two_content both map to "Two Content（雙欄必較頁）".
    if layout_type == "comparison":
        layout_type = "two_content"

    keys = LAYOUT_NAME_KEYS.get(layout_type, [])
    
    for layout in prs.slide_layouts:
        name = (getattr(layout, "name", "") or "").lower()
        for k in keys:
            if k in name:
                return layout

    if layout_type == "cover":
        return prs.slide_layouts[0]
    
    if layout_type == "content":
        for layout in prs.slide_layouts:
            ph_types = [_placeholder_type_id(p) for p in layout.placeholders]
            has_title = any(_is_title_placeholder_type(t) for t in ph_types)
            has_content = any(_is_content_placeholder_type(t) for t in ph_types)
            if has_title and has_content:
                return layout
        return prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    if layout_type in ("two_content", "comparison"):
        for layout in prs.slide_layouts:
            ph_types = [_placeholder_type_id(p) for p in layout.placeholders]
            content_count = sum(1 for t in ph_types if _is_content_placeholder_type(t))
            has_title = any(_is_title_placeholder_type(t) for t in ph_types)
            if has_title and content_count >= 2:
                return layout
        return _pick_layout_adaptive(prs, "content")

    if layout_type == "section":
        return _pick_layout_adaptive(prs, "cover")

    return prs.slide_layouts[0]


def _layout_name_key(name: str) -> str:
    text = unicodedata.normalize("NFC", str(name or "")).strip().casefold()
    return re.sub(r"\s+", "", text)


def _pick_layout_by_display_name(prs, layout_name: str):
    wanted = _layout_name_key(layout_name)
    if not wanted:
        return None
    for layout in prs.slide_layouts:
        try:
            name = str(getattr(layout, "name", "") or "").strip()
        except Exception:
            name = ""
        if _layout_name_key(name) == wanted:
            return layout
    return None


def _pick_layout_for_slide_data(prs, slide_data: Dict[str, Any], *, default_type: str = "content"):
    layout_name = str(slide_data.get("layout_name") or "").strip()
    if layout_name:
        layout = _pick_layout_by_display_name(prs, layout_name)
        if layout is not None:
            return layout
    return _pick_layout_adaptive(prs, default_type)


def _clear_all_slides(prs):
    """
    Keep masters/layouts from template, but remove existing content slides.
    """
    try:
        sld_id_list = prs.slides._sldIdLst  # pylint: disable=protected-access
        for sld_id in list(sld_id_list):
            r_id = sld_id.rId
            prs.part.drop_rel(r_id)
            sld_id_list.remove(sld_id)
    except Exception as e:
        logger.warning("Failed to clear template slides: %s", e)


def _keep_only_main_slide(prs):
    """
    Keep only the first (main) slide while preserving masters/layouts/theme.
    """
    try:
        sld_id_list = prs.slides._sldIdLst  # pylint: disable=protected-access
        if len(sld_id_list) <= 1:
            return
        for idx, sld_id in enumerate(list(sld_id_list)):
            if idx == 0:
                continue
            r_id = sld_id.rId
            prs.part.drop_rel(r_id)
            sld_id_list.remove(sld_id)
    except Exception as e:
        logger.warning("Failed to keep only main slide: %s", e)

    # Clear any text content on the retained main slide.
    try:
        if len(prs.slides) > 0:
            _clear_all_text_on_slide(prs.slides[0])
    except Exception as e:
        logger.warning("Failed to clear text on main slide: %s", e)


def _retain_representative_slides_and_clear_text(prs) -> None:
    """
    Keep total slide count unchanged and clear text on all slides.
    """
    try:
        total = len(prs.slides)
    except Exception:
        total = 0
    if total <= 0:
        return

    for slide in prs.slides:
        _clear_all_text_on_slide(slide)


def _slide_layout_identity(slide, fallback_index: int) -> str:
    try:
        layout = getattr(slide, "slide_layout", None)
        layout_name = str(getattr(layout, "name", "") or "").strip() or f"layout_{fallback_index}"
        master = getattr(layout, "slide_master", None)
        master_name = str(getattr(master, "name", "") or "").strip() or "master"
        return f"{master_name}::{layout_name}"
    except Exception:
        return f"layout_{fallback_index}"


def _collect_layout_occurrences(prs) -> Dict[str, List[Dict[str, Any]]]:
    """
    Collect all slide occurrences by layout identity, preserving order.
    """
    by_layout: Dict[str, List[Dict[str, Any]]] = {}
    for idx, slide in enumerate(prs.slides):
        try:
            layout_part_name = str(slide.slide_layout.part.partname or "").strip().lstrip("/")
        except Exception:
            layout_part_name = ""
        key = layout_part_name or _slide_layout_identity(slide, idx)
        try:
            part_name = str(slide.part.partname or "").lstrip("/")
        except Exception:
            part_name = ""
        if not part_name:
            continue
        occurrence = {
            "slide_index": idx,
            "source_slide_part": part_name,
        }
        if layout_part_name:
            occurrence["layout_part"] = layout_part_name
        by_layout.setdefault(key, []).append(occurrence)
    return by_layout


def _layout_identity_to_partname_map(prs) -> Dict[str, str]:
    out: Dict[str, str] = {}
    for _midx, master in enumerate(prs.slide_masters):
        master_name = str(getattr(master, "name", "") or "").strip() or "master"
        for lidx, layout in enumerate(master.slide_layouts):
            layout_name = str(getattr(layout, "name", "") or "").strip() or f"layout_{lidx}"
            key = f"{master_name}::{layout_name}"
            try:
                partname = str(layout.part.partname or "").strip()
            except Exception:
                partname = ""
            if partname:
                out[key] = partname
    return out


def _collect_referenced_layout_parts_from_presentation(prs) -> Set[str]:
    referenced: Set[str] = set()
    for slide in getattr(prs, "slides", []):
        try:
            slide_layout = getattr(slide, "slide_layout", None)
            partname = str(getattr(getattr(slide_layout, "part", None), "partname", "") or "").strip().lstrip("/")
        except Exception:
            partname = ""
        if partname:
            referenced.add(partname)
    return referenced


def _layout_part_display_name_map(prs, keep_layout_parts: Optional[Set[str]] = None) -> Dict[str, str]:
    """
    Assign sequential display names to layouts in the order they appear in slide masters.

    The numbering is presentation-wide and stable for the current PPTX order.
    """
    out: Dict[str, str] = {}
    seq = 1
    for master in getattr(prs, "slide_masters", []):
        for layout in getattr(master, "slide_layouts", []):
            try:
                partname = str(layout.part.partname or "").strip().lstrip("/")
            except Exception:
                partname = ""
            if keep_layout_parts is not None and (not partname or partname not in keep_layout_parts):
                continue
            if not partname or partname in out:
                continue
            out[partname] = f"版型{seq}"
            seq += 1
    return out


def _collect_referenced_layout_parts_from_file_map(file_map: Dict[str, bytes]) -> Set[str]:
    rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    layout_rel_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
    referenced: Set[str] = set()
    for rels_path, raw in file_map.items():
        if not re.match(r"^ppt/slides/_rels/slide\d+\.xml\.rels$", str(rels_path)):
            continue
        slide_part = str(rels_path).replace("/_rels/", "/").replace(".rels", "")
        try:
            rels_root = ET.fromstring(raw)
        except Exception:
            continue
        for rel in rels_root.findall(f"{{{rel_ns}}}Relationship"):
            if str(rel.get("Type") or "") != layout_rel_type:
                continue
            target = str(rel.get("Target") or "").strip()
            if not target:
                continue
            resolved = _resolve_target_part_path(slide_part, target)
            if resolved:
                referenced.add(resolved)
    return referenced


def _remove_content_type_override_for_part(content_types_root: ET.Element, part_path: str) -> None:
    ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    part_name = f"/{str(part_path or '').lstrip('/')}"
    for node in list(content_types_root.findall(f"{{{ct_ns}}}Override")):
        if str(node.get("PartName") or "") == part_name:
            content_types_root.remove(node)


def _part_sort_key(part_path: str) -> tuple[int, str]:
    text = str(part_path or "")
    m = re.search(r"(\d+)(?=\.xml(?:\.rels)?$)", text)
    if m:
        return int(m.group(1)), text
    return 10**9, text


def _rename_kept_layout_display_names(file_map: Dict[str, bytes]) -> None:
    rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    slide_layout_rel_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
    seq = 1
    master_rels_paths = sorted(
        [
            str(name)
            for name in file_map.keys()
            if re.match(r"^ppt/slideMasters/_rels/slideMaster\d+\.xml\.rels$", str(name))
        ],
        key=_part_sort_key,
    )
    for master_rels_path in master_rels_paths:
        master_xml = str(master_rels_path).replace("/_rels/", "/").replace(".rels", "")
        if master_xml not in file_map:
            continue
        try:
            rels_root = ET.fromstring(file_map[master_rels_path])
        except Exception:
            continue
        for rel in rels_root.findall(f"{{{rel_ns}}}Relationship"):
            if str(rel.get("Type") or "") != slide_layout_rel_type:
                continue
            target = str(rel.get("Target") or "").strip()
            if not target:
                continue
            layout_part = _resolve_target_part_path(master_xml, target)
            if layout_part not in file_map:
                continue
            _set_layout_display_name(file_map, layout_part, f"版型{seq}")
            seq += 1


def _prune_unused_master_layout_parts(file_map: Dict[str, bytes]) -> None:
    content_types_name = "[Content_Types].xml"
    presentation_name = "ppt/presentation.xml"
    presentation_rels_name = "ppt/_rels/presentation.xml.rels"
    rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    slide_layout_rel_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
    slide_master_rel_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster"

    if content_types_name not in file_map or presentation_name not in file_map or presentation_rels_name not in file_map:
        return

    referenced_layout_parts = _collect_referenced_layout_parts_from_file_map(file_map)
    content_types_root = ET.fromstring(file_map[content_types_name])
    presentation_root = ET.fromstring(file_map[presentation_name])
    presentation_rels_root = ET.fromstring(file_map[presentation_rels_name])

    deleted_master_parts: Set[str] = set()
    deleted_master_rel_parts: Set[str] = set()

    master_rels_paths = [
        str(name)
        for name in file_map.keys()
        if re.match(r"^ppt/slideMasters/_rels/slideMaster\d+\.xml\.rels$", str(name))
    ]
    for master_rels_path in master_rels_paths:
        master_xml = str(master_rels_path).replace("/_rels/", "/").replace(".rels", "")
        if master_xml not in file_map:
            continue
        try:
            rels_root = ET.fromstring(file_map[master_rels_path])
        except Exception:
            continue

        kept_rids: Set[str] = set()
        removed_any = False
        for rel in list(rels_root.findall(f"{{{rel_ns}}}Relationship")):
            if str(rel.get("Type") or "") != slide_layout_rel_type:
                continue
            target = str(rel.get("Target") or "").strip()
            resolved = _resolve_target_part_path(master_xml, target) if target else ""
            if resolved and resolved in referenced_layout_parts and resolved in file_map:
                rid = str(rel.get("Id") or "").strip()
                if rid:
                    kept_rids.add(rid)
                continue
            rels_root.remove(rel)
            removed_any = True

        if not kept_rids:
            deleted_master_parts.add(master_xml)
            deleted_master_rel_parts.add(master_rels_path)
            continue

        if removed_any:
            file_map[master_rels_path] = ET.tostring(rels_root, encoding="utf-8", xml_declaration=True)

        try:
            master_root = ET.fromstring(file_map[master_xml])
        except Exception:
            continue
        layout_id_list = master_root.find(f"{{{p_ns}}}sldLayoutIdLst")
        if layout_id_list is not None:
            for node in list(layout_id_list.findall(f"{{{p_ns}}}sldLayoutId")):
                rid = str(node.get(f"{{{r_ns}}}id") or "").strip()
                if rid and rid in kept_rids:
                    continue
                layout_id_list.remove(node)
            if len(list(layout_id_list)) == 0:
                master_root.remove(layout_id_list)
        file_map[master_xml] = ET.tostring(master_root, encoding="utf-8", xml_declaration=True)

    for layout_part in sorted(
        [
            str(name)
            for name in file_map.keys()
            if re.match(r"^ppt/slideLayouts/slideLayout\d+\.xml$", str(name))
        ]
    ):
        if layout_part in referenced_layout_parts:
            continue
        file_map.pop(layout_part, None)
        file_map.pop(_rel_path_for_part(layout_part), None)
        _remove_content_type_override_for_part(content_types_root, layout_part)

    presentation_id_list = presentation_root.find(f"{{{p_ns}}}sldMasterIdLst")
    if presentation_id_list is not None:
        for node in list(presentation_id_list.findall(f"{{{p_ns}}}sldMasterId")):
            rid = str(node.get(f"{{{r_ns}}}id") or "").strip()
            target_part = ""
            if rid:
                for rel in presentation_rels_root.findall(f"{{{rel_ns}}}Relationship"):
                    if str(rel.get("Id") or "").strip() != rid:
                        continue
                    if str(rel.get("Type") or "") != slide_master_rel_type:
                        continue
                    target_part = _resolve_target_part_path(
                        presentation_name,
                        str(rel.get("Target") or "").strip(),
                    )
                    break
            if target_part and target_part in deleted_master_parts:
                presentation_id_list.remove(node)
                if rid:
                    for rel in list(presentation_rels_root.findall(f"{{{rel_ns}}}Relationship")):
                        if str(rel.get("Id") or "").strip() == rid:
                            presentation_rels_root.remove(rel)
                            break

    for master_part in sorted(deleted_master_parts):
        deleted_master_rel_parts.add(_rel_path_for_part(master_part))
        file_map.pop(master_part, None)
        file_map.pop(_rel_path_for_part(master_part), None)
        _remove_content_type_override_for_part(content_types_root, master_part)

    for rels_part in sorted(deleted_master_rel_parts):
        file_map.pop(rels_part, None)

    _rename_kept_layout_display_names(file_map)

    for override in content_types_root.findall("{http://schemas.openxmlformats.org/package/2006/content-types}Override"):
        if str(override.get("PartName") or "") == "/ppt/presentation.xml":
            override.set("ContentType", PPTX_MAIN_CONTENT_TYPE)
            break

    file_map[content_types_name] = ET.tostring(content_types_root, encoding="utf-8", xml_declaration=True)
    file_map[presentation_name] = ET.tostring(presentation_root, encoding="utf-8", xml_declaration=True)
    file_map[presentation_rels_name] = ET.tostring(presentation_rels_root, encoding="utf-8", xml_declaration=True)


def _rel_path_for_part(part_path: str) -> str:
    # ppt/slideLayouts/slideLayout1.xml -> ppt/slideLayouts/_rels/slideLayout1.xml.rels
    base_dir, filename = os.path.split(part_path)
    return "/".join([base_dir, "_rels", f"{filename}.rels"])


def _next_rid(rels_root: ET.Element) -> str:
    used: List[int] = []
    for rel in rels_root.findall("{http://schemas.openxmlformats.org/package/2006/relationships}Relationship"):
        rid = str(rel.get("Id") or "")
        m = re.match(r"^rId(\d+)$", rid)
        if m:
            used.append(int(m.group(1)))
    return f"rId{(max(used) + 1) if used else 1}"


def _ensure_content_type_default_for_ext(content_types_root: ET.Element, ext: str, content_type: str) -> None:
    ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    wanted_ext = (ext or "png").lower()
    for node in content_types_root.findall(f"{{{ct_ns}}}Default"):
        if str(node.get("Extension") or "").lower() == wanted_ext:
            return
    default_node = ET.SubElement(content_types_root, f"{{{ct_ns}}}Default")
    default_node.set("Extension", wanted_ext)
    default_node.set("ContentType", content_type)


def _resolve_target_part_path(base_part_path: str, target: str) -> str:
    # base: ppt/slides/slide1.xml, target: ../media/image1.png -> ppt/media/image1.png
    base_dir = posixpath.dirname(base_part_path)
    return posixpath.normpath(posixpath.join(base_dir, target))


def _local_name(tag: str) -> str:
    return str(tag or "").split("}")[-1]


def _shape_placeholder_type_id(node: ET.Element, ns: Dict[str, str]) -> int | None:
    checks = [
        "p:nvSpPr/p:nvPr/p:ph",
        "p:nvPicPr/p:nvPr/p:ph",
        "p:nvGrpSpPr/p:nvPr/p:ph",
        "p:nvGraphicFramePr/p:nvPr/p:ph",
        "p:nvCxnSpPr/p:nvPr/p:ph",
        "p:nvContentPartPr/p:nvPr/p:ph",
    ]
    for c in checks:
        ph = node.find(c, ns)
        if ph is not None:
            return _placeholder_type_id(ph)
    return None


def _remove_placeholder_metadata(node: ET.Element, ns: Dict[str, str]) -> None:
    for path in [
        "p:nvSpPr/p:nvPr",
        "p:nvPicPr/p:nvPr",
        "p:nvGrpSpPr/p:nvPr",
        "p:nvGraphicFramePr/p:nvPr",
        "p:nvCxnSpPr/p:nvPr",
        "p:nvContentPartPr/p:nvPr",
    ]:
        nv_pr = node.find(path, ns)
        if nv_pr is None:
            continue
        for ph in list(nv_pr.findall("p:ph", ns)):
            nv_pr.remove(ph)


def _next_shape_id_from_tree(root: ET.Element, ns: Dict[str, str]) -> int:
    max_id = 0
    for c_nv_pr in root.findall(".//p:cNvPr", ns):
        try:
            max_id = max(max_id, int(str(c_nv_pr.get("id") or "0")))
        except Exception:
            continue
    return max_id + 1


def _assign_shape_ids(node: ET.Element, start_id: int, ns: Dict[str, str]) -> int:
    next_id = start_id
    for c_nv_pr in node.findall(".//p:cNvPr", ns):
        c_nv_pr.set("id", str(next_id))
        next_id += 1
    return next_id


def _clear_text_in_shape_xml(node: ET.Element) -> None:
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for text_node in node.findall(f".//{{{a_ns}}}t"):
        text_node.text = ""


def _node_is_text_shape(node: ET.Element, ns: Dict[str, str]) -> bool:
    """
    Preserve real textbox / placeholder shapes on the slide.

    We only treat `p:sp` nodes with a local `p:txBody` as text objects.
    This keeps title/body text boxes while still pruning images, charts, and
    other non-text objects from the slide copy.
    """
    local_name = _local_name(node.tag)
    if local_name == "sp":
        return node.find("p:txBody", ns) is not None
    if local_name == "grpSp":
        return any(_node_is_text_shape(child, ns) for child in list(node))
    return False


def _copy_content_type_for_part(
    *,
    part_path: str,
    src_content_types_root: ET.Element,
    dst_content_types_root: ET.Element,
) -> None:
    ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    ext = posixpath.splitext(part_path)[1].lstrip(".").lower()

    # Prefer Override for exact part.
    part_name = f"/{part_path.lstrip('/')}"
    for node in src_content_types_root.findall(f"{{{ct_ns}}}Override"):
        if str(node.get("PartName") or "") == part_name:
            exists = any(
                str(x.get("PartName") or "") == part_name
                for x in dst_content_types_root.findall(f"{{{ct_ns}}}Override")
            )
            if not exists:
                new_node = ET.SubElement(dst_content_types_root, f"{{{ct_ns}}}Override")
                new_node.set("PartName", part_name)
                new_node.set("ContentType", str(node.get("ContentType") or "application/octet-stream"))
            return

    # Fallback to Default by extension.
    if not ext:
        return
    for node in src_content_types_root.findall(f"{{{ct_ns}}}Default"):
        if str(node.get("Extension") or "").lower() == ext:
            _ensure_content_type_default_for_ext(
                dst_content_types_root,
                ext,
                str(node.get("ContentType") or "application/octet-stream"),
            )
            return


def _copy_source_part_tree(
    *,
    source_map: Dict[str, bytes],
    file_map: Dict[str, bytes],
    src_content_types_root: ET.Element,
    dst_content_types_root: ET.Element,
    part_path: str,
    visited: Optional[Set[str]] = None,
) -> None:
    rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"

    part_path = str(part_path or "").strip().lstrip("/")
    if not part_path:
        return
    if visited is None:
        visited = set()
    if part_path in visited:
        return
    visited.add(part_path)

    if part_path in source_map and part_path not in file_map:
        file_map[part_path] = source_map[part_path]
        _copy_content_type_for_part(
            part_path=part_path,
            src_content_types_root=src_content_types_root,
            dst_content_types_root=dst_content_types_root,
        )

    rels_path = _rel_path_for_part(part_path)
    if rels_path not in source_map:
        return
    if rels_path not in file_map:
        file_map[rels_path] = source_map[rels_path]

    try:
        rels_root = ET.fromstring(source_map[rels_path])
    except Exception:
        return

    for rel in rels_root.findall(f"{{{rel_ns}}}Relationship"):
        if str(rel.get("TargetMode") or "").lower() == "external":
            continue
        target = str(rel.get("Target") or "").strip()
        if not target:
            continue
        target_part = _resolve_target_part_path(part_path, target)
        if not target_part:
            continue
        _copy_source_part_tree(
            source_map=source_map,
            file_map=file_map,
            src_content_types_root=src_content_types_root,
            dst_content_types_root=dst_content_types_root,
            part_path=target_part,
            visited=visited,
        )


def _strip_nonplaceholder_shapes_from_slide_parts(file_map: Dict[str, bytes]) -> None:
    rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    slide_rel_types_to_keep = {
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideComments",
    }
    ns = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }

    slide_part_names = sorted(
        [
            str(name)
            for name in file_map.keys()
            if re.match(r"^ppt/slides/slide\d+\.xml$", str(name))
        ],
        key=_part_sort_key,
    )
    for slide_part in slide_part_names:
        if slide_part not in file_map:
            continue
        slide_root = ET.fromstring(file_map[slide_part])
        sp_tree = slide_root.find(".//p:spTree", ns)
        if sp_tree is None:
            continue

        for node in list(sp_tree):
            if _local_name(node.tag) in {"nvGrpSpPr", "grpSpPr", "extLst"}:
                continue
            if _shape_placeholder_type_id(node, ns) is not None:
                continue
            if _node_is_text_shape(node, ns):
                continue
            sp_tree.remove(node)

        remaining_rids: Set[str] = set()
        for el in slide_root.iter():
            for attr_name, attr_val in el.attrib.items():
                if not attr_name.startswith("{") or not attr_val:
                    continue
                if attr_name[1:].split("}", 1)[0] != "http://schemas.openxmlformats.org/officeDocument/2006/relationships":
                    continue
                remaining_rids.add(str(attr_val))

        rels_path = _rel_path_for_part(slide_part)
        if rels_path not in file_map:
            file_map[slide_part] = ET.tostring(slide_root, encoding="utf-8", xml_declaration=True)
            continue

        rels_root = ET.fromstring(file_map[rels_path])
        for rel in list(rels_root.findall(f"{{{rel_ns}}}Relationship")):
            rid = str(rel.get("Id") or "").strip()
            rel_type = str(rel.get("Type") or "").strip()
            if rid in remaining_rids or rel_type in slide_rel_types_to_keep:
                continue
            rels_root.remove(rel)

        file_map[slide_part] = ET.tostring(slide_root, encoding="utf-8", xml_declaration=True)
        file_map[rels_path] = ET.tostring(rels_root, encoding="utf-8", xml_declaration=True)


def _layout_base_name_from_key(layout_key: str) -> str:
    text = str(layout_key or "")
    base = text.split("::", 1)[1] if "::" in text else text
    base = re.sub(r"_\d+$", "", base).strip()
    return base or "layout"


def _next_slide_layout_part_path(file_map: Dict[str, bytes]) -> str:
    used: List[int] = []
    for name in file_map.keys():
        m = re.match(r"^ppt/slideLayouts/slideLayout(\d+)\.xml$", str(name))
        if m:
            used.append(int(m.group(1)))
    nxt = (max(used) + 1) if used else 1
    return f"ppt/slideLayouts/slideLayout{nxt}.xml"


def _copy_content_type_override_for_new_part(content_types_root: ET.Element, source_part: str, new_part: str) -> None:
    ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    src_name = f"/{source_part.lstrip('/')}"
    dst_name = f"/{new_part.lstrip('/')}"
    exists = any(str(x.get("PartName") or "") == dst_name for x in content_types_root.findall(f"{{{ct_ns}}}Override"))
    if exists:
        return
    for node in content_types_root.findall(f"{{{ct_ns}}}Override"):
        if str(node.get("PartName") or "") == src_name:
            new_node = ET.SubElement(content_types_root, f"{{{ct_ns}}}Override")
            new_node.set("PartName", dst_name)
            new_node.set("ContentType", str(node.get("ContentType") or "application/xml"))
            return


def _set_layout_display_name(file_map: Dict[str, bytes], layout_part: str, layout_name: str) -> None:
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    ns = {"p": p_ns}
    if layout_part not in file_map:
        return
    root = ET.fromstring(file_map[layout_part])
    c_sld = root.find("p:cSld", ns)
    if c_sld is None:
        c_sld = root.find(".//p:cSld", ns)
    if c_sld is not None:
        c_sld.set("name", layout_name)
    file_map[layout_part] = ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _find_master_binding_for_layout(file_map: Dict[str, bytes], layout_part: str) -> Optional[Dict[str, str]]:
    rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    layout_rel_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
    for rels_path, raw in file_map.items():
        if not re.match(r"^ppt/slideMasters/_rels/slideMaster\d+\.xml\.rels$", str(rels_path)):
            continue
        master_xml = str(rels_path).replace("/_rels/", "/").replace(".rels", "")
        if master_xml not in file_map:
            continue
        rels_root = ET.fromstring(raw)
        for rel in rels_root.findall(f"{{{rel_ns}}}Relationship"):
            if str(rel.get("Type") or "") != layout_rel_type:
                continue
            target = str(rel.get("Target") or "")
            if not target:
                continue
            resolved = _resolve_target_part_path(master_xml, target)
            if resolved == layout_part:
                return {"master_xml": master_xml, "master_rels": rels_path}
    return None


def _attach_layout_to_master(file_map: Dict[str, bytes], master_xml: str, master_rels: str, new_layout_part: str) -> None:
    rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    ns = {"p": p_ns}

    rels_root = ET.fromstring(file_map[master_rels])
    new_rid = _next_rid(rels_root)
    rel_node = ET.SubElement(rels_root, f"{{{rel_ns}}}Relationship")
    rel_node.set("Id", new_rid)
    rel_node.set("Type", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout")
    rel_target = posixpath.relpath(new_layout_part, posixpath.dirname(master_xml))
    rel_node.set("Target", rel_target)
    file_map[master_rels] = ET.tostring(rels_root, encoding="utf-8", xml_declaration=True)

    master_root = ET.fromstring(file_map[master_xml])
    lst = master_root.find("p:sldLayoutIdLst", ns)
    if lst is None:
        lst = ET.SubElement(master_root, f"{{{p_ns}}}sldLayoutIdLst")
    max_id = 255
    for node in lst.findall("p:sldLayoutId", ns):
        try:
            max_id = max(max_id, int(node.get("id") or 0))
        except Exception:
            continue
    entry = ET.SubElement(lst, f"{{{p_ns}}}sldLayoutId")
    entry.set("id", str(max_id + 1))
    entry.set(f"{{{r_ns}}}id", new_rid)
    file_map[master_xml] = ET.tostring(master_root, encoding="utf-8", xml_declaration=True)


def _set_slide_layout_relationship_target(file_map: Dict[str, bytes], slide_part: str, layout_part: str) -> None:
    rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    layout_rel_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
    rels_path = _rel_path_for_part(slide_part)
    if rels_path not in file_map:
        return
    rels_root = ET.fromstring(file_map[rels_path])
    changed = False
    for rel in rels_root.findall(f"{{{rel_ns}}}Relationship"):
        if str(rel.get("Type") or "") != layout_rel_type:
            continue
        target = posixpath.relpath(layout_part, posixpath.dirname(slide_part))
        rel.set("Target", target)
        changed = True
    if changed:
        file_map[rels_path] = ET.tostring(rels_root, encoding="utf-8", xml_declaration=True)


def _replace_layout_content_from_source_slide(
    *,
    file_map: Dict[str, bytes],
    source_map: Dict[str, bytes],
    src_content_types_root: ET.Element,
    dst_content_types_root: ET.Element,
    source_slide_part: str,
    target_layout_part: str,
) -> None:
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    rel_attr_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    ns = {"p": p_ns}

    if source_slide_part not in source_map or target_layout_part not in file_map:
        return

    source_slide_root = ET.fromstring(source_map[source_slide_part])
    source_sp_tree = source_slide_root.find(".//p:spTree", ns)
    if source_sp_tree is None:
        return

    source_shape_nodes: List[ET.Element] = []
    for node in list(source_sp_tree):
        if _local_name(node.tag) in {"nvGrpSpPr", "grpSpPr"}:
            continue
        source_shape_nodes.append(node)

    layout_root = ET.fromstring(file_map[target_layout_part])
    sp_tree = layout_root.find(".//p:spTree", ns)
    if sp_tree is None:
        return

    for node in list(sp_tree):
        if _local_name(node.tag) in {"nvGrpSpPr", "grpSpPr"}:
            continue
        sp_tree.remove(node)

    rels_path = _rel_path_for_part(target_layout_part)
    if rels_path in file_map:
        rels_root = ET.fromstring(file_map[rels_path])
    else:
        rels_root = ET.Element(f"{{{rel_ns}}}Relationships")

    src_slide_rels_path = _rel_path_for_part(source_slide_part)
    if src_slide_rels_path in source_map:
        src_slide_rels_root = ET.fromstring(source_map[src_slide_rels_path])
    else:
        src_slide_rels_root = ET.Element(f"{{{rel_ns}}}Relationships")
    src_rel_by_id: Dict[str, ET.Element] = {
        str(rel.get("Id") or ""): rel
        for rel in src_slide_rels_root.findall(f"{{{rel_ns}}}Relationship")
    }
    next_shape_id = _next_shape_id_from_tree(sp_tree, ns)

    for node in source_shape_nodes:
        cloned = copy.deepcopy(node)
        _clear_text_in_shape_xml(cloned)
        next_shape_id = _assign_shape_ids(cloned, next_shape_id, ns)

        rid_map: Dict[str, str] = {}
        for el in cloned.iter():
            for attr_name, attr_val in list(el.attrib.items()):
                if not attr_name.startswith("{") or not attr_val:
                    continue
                ns_uri, _local = attr_name[1:].split("}", 1)
                if ns_uri != rel_attr_ns:
                    continue
                old_rid = str(attr_val)
                if old_rid not in rid_map:
                    src_rel = src_rel_by_id.get(old_rid)
                    if src_rel is None:
                        continue
                    new_rid = _next_rid(rels_root)
                    rid_map[old_rid] = new_rid

                    new_rel = ET.SubElement(rels_root, f"{{{rel_ns}}}Relationship")
                    new_rel.set("Id", new_rid)
                    new_rel.set("Type", str(src_rel.get("Type") or ""))
                    new_rel.set("Target", str(src_rel.get("Target") or ""))
                    target_mode = str(src_rel.get("TargetMode") or "")
                    if target_mode:
                        new_rel.set("TargetMode", target_mode)

                    if target_mode.lower() != "external":
                        target = str(src_rel.get("Target") or "")
                        if target:
                            src_target_part = _resolve_target_part_path(source_slide_part, target)
                            _copy_source_part_tree(
                                source_map=source_map,
                                file_map=file_map,
                                src_content_types_root=src_content_types_root,
                                dst_content_types_root=dst_content_types_root,
                                part_path=src_target_part,
                            )
                mapped = rid_map.get(old_rid)
                if mapped:
                    el.set(attr_name, mapped)

        sp_tree.append(cloned)

    file_map[target_layout_part] = ET.tostring(layout_root, encoding="utf-8", xml_declaration=True)
    file_map[rels_path] = ET.tostring(rels_root, encoding="utf-8", xml_declaration=True)


def _inject_layout_representative_shapes(
    source_raw: bytes,
    normalized_raw: bytes,
    *,
    layout_occurrences: Dict[str, List[Dict[str, Any]]],
    layout_part_map: Dict[str, str],
) -> bytes:
    if not layout_occurrences:
        return normalized_raw

    content_types_name = "[Content_Types].xml"
    source_map: Dict[str, bytes] = {}
    with zipfile.ZipFile(io.BytesIO(source_raw), "r") as zf_source:
        for name in zf_source.namelist():
            source_map[name] = zf_source.read(name)

    with zipfile.ZipFile(io.BytesIO(normalized_raw), "r") as zin:
        file_map: Dict[str, bytes] = {name: zin.read(name) for name in zin.namelist()}

    src_content_types_root = ET.fromstring(source_map[content_types_name])
    content_types_root = ET.fromstring(file_map[content_types_name])

    for key, occs in layout_occurrences.items():
        if not occs:
            continue
        explicit_layout = str(occs[0].get("layout_part") or "").strip().lstrip("/")
        original_layout = explicit_layout or str(layout_part_map.get(key) or "").lstrip("/")
        if not original_layout or original_layout not in file_map:
            continue

        target_layouts = [original_layout]
        if len(occs) > 1:
            binding = _find_master_binding_for_layout(file_map, original_layout)
            if binding:
                for _ in range(1, len(occs)):
                    new_layout = _next_slide_layout_part_path(file_map)
                    file_map[new_layout] = file_map[original_layout]
                    original_layout_rels = _rel_path_for_part(original_layout)
                    new_layout_rels = _rel_path_for_part(new_layout)
                    if original_layout_rels in file_map:
                        file_map[new_layout_rels] = file_map[original_layout_rels]
                    _copy_content_type_override_for_new_part(content_types_root, original_layout, new_layout)
                    _attach_layout_to_master(
                        file_map,
                        binding["master_xml"],
                        binding["master_rels"],
                        new_layout,
                    )
                    target_layouts.append(new_layout)

        if len(occs) > 1:
            base = _layout_base_name_from_key(key)
            for idx, layout_part in enumerate(target_layouts, start=1):
                _set_layout_display_name(file_map, layout_part, f"{base}_{idx}")

        for idx, occ in enumerate(occs):
            target_layout = target_layouts[min(idx, len(target_layouts) - 1)]
            src_slide_part = str(occ.get("source_slide_part") or "").strip()
            if src_slide_part:
                _replace_layout_content_from_source_slide(
                    file_map=file_map,
                    source_map=source_map,
                    src_content_types_root=src_content_types_root,
                    dst_content_types_root=content_types_root,
                    source_slide_part=src_slide_part,
                    target_layout_part=target_layout,
                )
            normalized_slide_part = str(occ.get("normalized_slide_part") or "").strip()
            if normalized_slide_part:
                _set_slide_layout_relationship_target(file_map, normalized_slide_part, target_layout)

    file_map[content_types_name] = ET.tostring(content_types_root, encoding="utf-8", xml_declaration=True)
    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", compression=zipfile.ZIP_DEFLATED) as zout:
        for name, raw_data in file_map.items():
            zout.writestr(name, raw_data)
    return out.getvalue()


def _clear_all_text_on_slide(slide):
    for shape in list(getattr(slide, "shapes", [])):
        _clear_all_text_on_shape(shape)


def _clear_all_text_on_shape(shape):
    element = getattr(shape, "_element", None)
    if element is not None:
        try:
            _clear_text_in_shape_xml(element)
        except Exception:
            pass
    elif getattr(shape, "has_text_frame", False):
        try:
            shape.text_frame.clear()
        except Exception:
            pass

    # Some shapes may raise while accessing has_table/table.
    try:
        has_table = bool(getattr(shape, "has_table", False))
    except Exception:
        has_table = False
    if has_table:
        try:
            for row in shape.table.rows:
                for cell in row.cells:
                    cell_element = getattr(cell, "_tc", None)
                    if cell_element is not None:
                        _clear_text_in_shape_xml(cell_element)
                    else:
                        cell.text = ""
        except Exception:
            pass

    # Group shape recursion.
    child_shapes = getattr(shape, "shapes", None)
    if child_shapes is not None:
        try:
            for child in list(child_shapes):
                _clear_all_text_on_shape(child)
        except Exception:
            pass


def _layout_has_title(layout) -> bool:
    ph_types = [_placeholder_type_id(p) for p in layout.placeholders]
    return any(_is_title_placeholder_type(t) for t in ph_types)


def _layout_content_count(layout) -> int:
    ph_types = [_placeholder_type_id(p) for p in layout.placeholders]
    return sum(1 for t in ph_types if _is_content_placeholder_type(t))


def _layout_has_subtitle_or_textbox(layout) -> bool:
    ph_types = [_placeholder_type_id(p) for p in layout.placeholders]
    if any(_is_subtitle_placeholder_type(t) for t in ph_types):
        return True
    if any(_is_content_placeholder_type(t) for t in ph_types):
        return True
    for shape in layout.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if getattr(shape, "is_placeholder", False):
            t = _placeholder_type_id(shape)
            if _is_title_placeholder_type(t):
                continue
        return True
    return False


def _validate_layout_for_type(layout, layout_type: str) -> bool:
    if layout_type == "cover":
        return _layout_has_title(layout) and _layout_has_subtitle_or_textbox(layout)
    if layout_type == "content":
        return _layout_has_title(layout) and _layout_content_count(layout) >= 1
    if layout_type in ("two_content", "comparison"):
        return _layout_has_title(layout) and _layout_content_count(layout) >= 2
    if layout_type == "section":
        return _layout_has_title(layout) and _layout_has_subtitle_or_textbox(layout)
    return False


def _find_layout_for_import_validation(prs, layout_type: str):
    # PowerPoint 2016 class mapping:
    # comparison and two_content both map to "Two Content（雙欄必較頁）".
    if layout_type == "comparison":
        layout_type = "two_content"

    keys = LAYOUT_NAME_KEYS.get(layout_type, [])
    for layout in prs.slide_layouts:
        name = (getattr(layout, "name", "") or "").lower()
        if keys and not any(k in name for k in keys):
            continue
        if _validate_layout_for_type(layout, layout_type):
            return layout
    # Fallback: allow non-default localized layout names if structure is valid.
    for layout in prs.slide_layouts:
        if _validate_layout_for_type(layout, layout_type):
            return layout
    return None


def _audit_template_bytes(raw: bytes) -> Dict[str, Any]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(raw))
    # Required layouts follow PowerPoint 2016 default categories only.
    required = ["cover", "section", "content", "two_content"]
    labels = {
        "cover": "標題投影片（Title Slide）",
        "section": "章節標題（Section Header）",
        "content": "內容頁（Title and Content）",
        "two_content": "雙欄必較頁（Two Content）",
    }
    found: Dict[str, str] = {}
    missing: List[str] = []
    for t in required:
        layout = _find_layout_for_import_validation(prs, t)
        if layout is None:
            missing.append(labels[t])
        else:
            found[t] = getattr(layout, "name", "") or "(unnamed)"
    return {"ok": len(missing) == 0, "missing": missing, "found": found}


def _save_template_bytes(filename: str, raw: bytes) -> str:
    os.makedirs(PPTX_TEMPLATE_DIR, exist_ok=True)
    filename = _normalize_template_name(filename)
    base, ext = os.path.splitext(filename)
    base = _safe_filename(base, default="template")
    ext = ext if ext else ".pptx"
    candidate = f"{base}{ext}"
    existing_keys = set()
    for fn in os.listdir(PPTX_TEMPLATE_DIR):
        full = os.path.join(PPTX_TEMPLATE_DIR, fn)
        if os.path.isfile(full):
            existing_keys.add(_template_name_key(fn))
    idx = 1
    while (
        _template_name_key(candidate) in existing_keys
        or os.path.exists(os.path.join(PPTX_TEMPLATE_DIR, candidate))
    ):
        candidate = f"{base}_{idx}{ext}"
        idx += 1
    with open(os.path.join(PPTX_TEMPLATE_DIR, candidate), "wb") as f:
        f.write(raw)
    return candidate


def _find_largest_text_frame(slide, exclude_shape=None):
    best_shape = None
    best_area = -1
    
    for shape in slide.shapes:
        if exclude_shape is not None and shape == exclude_shape:
            continue
        if getattr(shape, "is_placeholder", False):
            if _is_content_placeholder_type(_placeholder_type_id(shape)):
                return shape.text_frame

    for shape in slide.shapes:
        if exclude_shape is not None and shape == exclude_shape:
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            area = int(shape.width) * int(shape.height)
        except Exception:
            area = 0
        if area > best_area:
            best_area = area
            best_shape = shape
            
    return best_shape.text_frame if best_shape else None


def _find_title_shape(slide):
    title_shape = slide.shapes.title if slide.shapes.title else None
    if title_shape:
        return title_shape
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if getattr(shape, "is_placeholder", False) and _is_title_placeholder_type(_placeholder_type_id(shape)):
            return shape
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        name = (getattr(shape, "name", "") or "").lower()
        if "title" in name or "標題" in name:
            return shape
    return None


def _find_subtitle_text_frame(slide, exclude_shape=None):
    for shape in slide.shapes:
        if exclude_shape is not None and shape == exclude_shape:
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        if getattr(shape, "is_placeholder", False) and _is_subtitle_placeholder_type(_placeholder_type_id(shape)):
            return shape.text_frame
    return _find_largest_text_frame(slide, exclude_shape=exclude_shape)


def _chunk_list(items: List[str], n: int):
    for i in range(0, len(items), n):
        yield items[i : i + n]


def _split_evenly(items: List[str], parts: int) -> List[List[str]]:
    if parts <= 0:
        return []
    seq = list(items or [])
    if not seq:
        return [[] for _ in range(parts)]
    base, extra = divmod(len(seq), parts)
    out: List[List[str]] = []
    start = 0
    for idx in range(parts):
        take = base + (1 if idx < extra else 0)
        out.append(seq[start : start + take] if take > 0 else [])
        start += take
    return out


def _fill_text_frames_with_lines(text_frames: List[Any], lines: List[str]) -> None:
    frames = [frame for frame in (text_frames or []) if frame is not None]
    if not frames:
        return

    items = [str(line or "").strip() for line in (lines or []) if str(line or "").strip()]
    if len(frames) == 1:
        _set_text_frame_lines_preserving_style(frames[0], items)
        return

    chunks = _split_evenly(items, len(frames))
    for text_frame, chunk in zip(frames, chunks):
        _set_text_frame_lines_preserving_style(text_frame, chunk)


def _resolve_slide_type(raw: Any) -> str:
    value = str(raw or "").strip().lower()
    if value in VALID_SLIDE_TYPES:
        return value
    resolved = _resolve_marker(str(raw or "").strip())
    if resolved in VALID_SLIDE_TYPES:
        return resolved
    return "content"


def _coerce_bool(value: Any, *, default: bool = False) -> bool:
    if isinstance(value, bool):
        return value
    if value is None:
        return default
    text = str(value).strip().lower()
    if text in {"1", "true", "yes", "y", "on"}:
        return True
    if text in {"0", "false", "no", "n", "off"}:
        return False
    return default


def _normalize_aspect_ratio(value: Any, default: str = "16:9") -> str:
    text = str(value or "").strip()
    if text in {"1:1", "4:3", "3:2", "16:9", "9:16"}:
        return text
    return default


def _normalize_image_intent(value: Any, default: str = "concept") -> str:
    intent = str(value or "").strip().lower()
    if intent in VALID_IMAGE_INTENTS:
        return intent
    return default


def _clean_prompt_fragment(value: Any, *, max_len: int = 48) -> str:
    text = str(value or "").strip()
    if not text:
        return ""
    text = re.sub(r"^[\-\*\u2022\d\.\)\(]+", "", text).strip()
    if re.match(r"^(required|intent|aspect_ratio|prompt)\s*:", text, flags=re.IGNORECASE):
        return ""
    if re.match(r"^\[\d+\]\s+", text):
        return ""
    text = re.sub(r"\s+", " ", text).strip(" ,;")
    if len(text) > max_len:
        text = text[:max_len].rstrip() + "..."
    return text


def _infer_image_intent_from_slide(slide: Dict[str, Any]) -> str:
    slide_type = _resolve_slide_type(slide.get("slide_type"))
    if slide_type == "section":
        return "hero"
    if slide_type == "comparison":
        return "data"

    title = str(slide.get("title") or "")
    bullets = slide.get("bullets") or []
    if not isinstance(bullets, list):
        bullets = [bullets]
    combined = f"{title} " + " ".join(str(x) for x in bullets)
    low = combined.lower()

    process_keys = ("process", "workflow", "timeline", "roadmap", "phase", "steps")
    data_keys = ("compare", "comparison", "trend", "kpi", "roi", "metric", "chart", "data", "vs")
    if any(k in low for k in process_keys):
        return "process"
    if any(k in low for k in data_keys):
        return "data"
    return "concept"


def _synthesize_image_prompt(slide: Dict[str, Any], *, intent: str) -> str:
    title = str(slide.get("title") or "").strip()
    title = re.sub(r"\s*\((section|content|two_content|comparison)\)\s*$", "", title, flags=re.IGNORECASE).strip()
    title = title or "presentation key visual"

    bullets = slide.get("bullets") or []
    if not isinstance(bullets, list):
        bullets = [bullets]
    fragments: List[str] = []
    for item in bullets:
        frag = _clean_prompt_fragment(item)
        if not frag:
            continue
        fragments.append(frag)
        if len(fragments) >= 3:
            break
    focus = f"; focus: {'; '.join(fragments)}" if fragments else ""

    style_map = {
        "hero": "corporate hero visual, wide scene, clear subject, professional lighting",
        "data": "business infographic style, clean hierarchy, data-focused visual language",
        "process": "workflow illustration style, clear step nodes, process-focused composition",
        "concept": "modern business illustration, clean composition, key point emphasis",
    }
    style = style_map.get(intent, style_map["concept"])
    return f"{title}{focus}. {style}, 16:9, no text, no watermark."


def _autofill_image_prompts_if_all_empty(analysis_result: Dict[str, Any]) -> Dict[str, Any]:
    if not isinstance(analysis_result, dict):
        return analysis_result
    slides = analysis_result.get("slides")
    if not isinstance(slides, list) or not slides:
        return analysis_result

    has_any_prompt = False
    for slide in slides:
        if not isinstance(slide, dict):
            continue
        if str(slide.get("image_prompt") or "").strip():
            has_any_prompt = True
            break
    if has_any_prompt:
        return analysis_result

    for slide in slides:
        if not isinstance(slide, dict):
            continue
        inferred_intent = _infer_image_intent_from_slide(slide)
        intent = _normalize_image_intent(slide.get("image_intent"), default=inferred_intent)
        slide["image_intent"] = intent
        slide["aspect_ratio"] = _normalize_aspect_ratio(slide.get("aspect_ratio"))
        slide["image_prompt"] = _synthesize_image_prompt(slide, intent=intent)
        slide["image_required"] = True
    return analysis_result


def _resolve_marker(raw_marker: str) -> str | None:
    marker = (raw_marker or "").strip()
    if not marker:
        return None
    lower = marker.lower()
    if lower in MARKER_TO_SLIDE_TYPE:
        return MARKER_TO_SLIDE_TYPE[lower]
    return MARKER_TO_SLIDE_TYPE.get(marker)


def _is_layout_name_marker(raw_marker: str) -> bool:
    marker = (raw_marker or "").strip()
    if not marker:
        return False
    return bool(re.match(r"^(?:版型|layout)\s*\d+$", marker, flags=re.IGNORECASE))


def _is_layout_meta_line(line: str) -> bool:
    text = str(line or "").strip()
    if not text:
        return False
    return bool(
        re.match(
            r"^(?:版面名稱|版型名稱|layout(?:\s*name)?)\s*[:：]",
            text,
            flags=re.IGNORECASE,
        )
    )


def _is_extraction_meta_line(line: str) -> bool:
    text = str(line or "").strip()
    if not text:
        return False
    if _is_layout_meta_line(text):
        return True
    return bool(
        re.match(
            r"^(?:@schema|@template_schema|@slide_schema|@master_schema|@layout_schema|@shape|@textbox|@textframe|@paragraph|@run)\b",
            text,
            flags=re.IGNORECASE,
        )
    )


def _parse_marked_text_structure(text: str) -> Dict[str, Any] | None:
    blocks = [s.strip() for s in text.split("===") if s.strip()]
    if not blocks:
        return None

    parsed = []
    has_any_marker = False
    for block in blocks:
        lines = [ln.strip() for ln in block.splitlines() if ln.strip()]
        if not lines:
            continue
        first_line = lines[0].lstrip("\ufeff\u200b\u200e\u200f").strip()
        if first_line.startswith("[") and "]" in first_line:
            marker_part, title_part = first_line[1:].split("]", 1)
            marker_part = marker_part.strip()
            title_part = title_part.strip()
            if marker_part:
                has_any_marker = True
                parsed.append(
                    {
                        "marker": marker_part,
                        "layout_name": marker_part if _is_layout_name_marker(marker_part) else "",
                        "title": title_part,
                        "lines": lines[1:],
                    }
                )
                continue
        parsed.append({"marker": "", "layout_name": "", "title": lines[0], "lines": lines[1:]})

    if not has_any_marker:
        return None

    out = {
        "main_title": DEFAULT_MAIN_TITLE,
        "main_subtitle": DEFAULT_MAIN_SUBTITLE,
        "main_layout_name": "",
        "slides": [],
    }

    for idx, item in enumerate(parsed, start=1):
        marker_type = _resolve_marker(item.get("marker") or "")
        title = (item.get("title") or "").strip() or f"投影片 {idx}"
        layout_name = str(item.get("layout_name") or "").strip()
        lines = [
            str(x).strip()
            for x in (item.get("lines") or [])
            if str(x).strip() and not _is_extraction_meta_line(str(x))
        ]
        if marker_type is None:
            fallback_index = idx if _is_layout_name_marker(item.get("marker") or "") else max(2, idx)
            marker_type = _pick_marker_for_extracted_slide(
                extracted_index=fallback_index,
                title=title,
                body_lines=lines,
            )

        if marker_type == "cover":
            out["main_title"] = title or DEFAULT_MAIN_TITLE
            if lines:
                out["main_subtitle"] = "｜".join(lines)
            if layout_name:
                out["main_layout_name"] = layout_name
            continue

        if marker_type == "two_content":
            lower_title = title.lower()
            if (
                "comparison" in lower_title
                or "compare" in lower_title
                or "比較" in title
                or "對照" in title
            ):
                marker_type = "comparison"

        slide_type = marker_type if marker_type in VALID_SLIDE_TYPES else "content"
        out["slides"].append(
            {
                "title": title,
                "layout_name": layout_name,
                "slide_type": slide_type,
                "left_title": "",
                "right_title": "",
                "bullets": lines,
                "image_required": False,
                "image_prompt": "",
                "image_intent": "",
                "aspect_ratio": "16:9",
            }
        )

    if not out["slides"]:
        out = _fallback_from_raw_text(text)
    elif "main_layout_name" not in out:
        out["main_layout_name"] = ""
    return out


def _find_content_text_frames(slide, exclude_shape=None) -> List[Any]:
    candidates = []
    for order, shape in enumerate(slide.shapes):
        if exclude_shape is not None and shape == exclude_shape:
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        priority = 3
        if getattr(shape, "is_placeholder", False):
            ph_type = _placeholder_type_id(shape)
            if _is_title_placeholder_type(ph_type):
                continue
            if _is_content_placeholder_type(ph_type):
                priority = 0
            elif _is_subtitle_placeholder_type(ph_type):
                priority = 1
            else:
                # Skip non-text placeholders such as footer/date/slide number.
                continue
        candidates.append(
            (
                priority,
                int(getattr(shape, "top", 0)),
                int(getattr(shape, "left", 0)),
                order,
                shape.text_frame,
            )
        )
    if candidates:
        candidates.sort(key=lambda x: (x[0], x[1], x[2], x[3]))
        return [x[4] for x in candidates]
    ordered: List[tuple[int, int, Any]] = []
    for shape in slide.shapes:
        if exclude_shape is not None and shape == exclude_shape:
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        ordered.append((int(getattr(shape, "left", 0)), int(getattr(shape, "top", 0)), shape.text_frame))
    if ordered:
        ordered.sort(key=lambda x: (x[0], x[1]))
        return [x[2] for x in ordered]
    tf = _find_largest_text_frame(slide, exclude_shape=exclude_shape)
    return [tf] if tf else []


def _shape_area(shape) -> int:
    try:
        return int(shape.width) * int(shape.height)
    except Exception:
        return 0


def _shape_has_text(shape) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    text = (getattr(shape.text_frame, "text", "") or "").strip()
    return bool(text)


def _slide_size_emu(slide) -> tuple[int, int]:
    emu = 914400
    width = 10 * emu
    height = int(7.5 * emu)
    try:
        presentation = slide.part.package.presentation_part.presentation
        width = int(getattr(presentation, "slide_width", width))
        height = int(getattr(presentation, "slide_height", height))
    except Exception:
        pass
    return width, height


def _slide_default_image_box(slide) -> tuple[int, int, int, int]:
    width, height = _slide_size_emu(slide)

    left = int(width * 0.08)
    top = int(height * 0.22)
    box_width = int(width * 0.84)
    box_height = int(height * 0.68)
    return left, top, box_width, box_height


def _find_largest_image_frame(slide, exclude_shape=None):
    image_placeholder = None
    image_placeholder_area = -1

    for shape in slide.shapes:
        if exclude_shape is not None and shape == exclude_shape:
            continue
        if not getattr(shape, "is_placeholder", False):
            continue
        if not _is_image_placeholder_type(_placeholder_type_id(shape)):
            continue
        area = _shape_area(shape)
        if area > image_placeholder_area:
            image_placeholder_area = area
            image_placeholder = shape

    if image_placeholder is not None:
        return image_placeholder

    candidates = []
    for shape in slide.shapes:
        if exclude_shape is not None and shape == exclude_shape:
            continue
        area = _shape_area(shape)
        if area <= 0:
            continue

        if getattr(shape, "is_placeholder", False):
            p_type = _placeholder_type_id(shape)
            if _is_title_placeholder_type(p_type) or _is_subtitle_placeholder_type(p_type):
                continue
            if _is_image_placeholder_type(p_type):
                candidates.append((area, shape))
                continue
            if _is_content_placeholder_type(p_type):
                if _shape_has_text(shape):
                    continue
                candidates.append((area, shape))
                continue
            # Ignore non-content placeholders (e.g., footer/date/number) to avoid accidental overlay.
            continue

        if getattr(shape, "has_text_frame", False):
            if _shape_has_text(shape):
                continue
            candidates.append((area, shape))

    if not candidates:
        return None
    candidates.sort(key=lambda x: x[0], reverse=True)
    return candidates[0][1]


def _shape_rect(shape) -> tuple[int, int, int, int] | None:
    try:
        left = int(getattr(shape, "left", 0))
        top = int(getattr(shape, "top", 0))
        width = int(getattr(shape, "width", 0))
        height = int(getattr(shape, "height", 0))
    except Exception:
        return None
    if width <= 0 or height <= 0:
        return None
    return left, top, width, height


def _rect_area(rect: tuple[int, int, int, int]) -> int:
    return max(0, int(rect[2])) * max(0, int(rect[3]))


def _rect_intersection_area(a: tuple[int, int, int, int], b: tuple[int, int, int, int]) -> int:
    ax1, ay1, aw, ah = a
    bx1, by1, bw, bh = b
    ax2, ay2 = ax1 + aw, ay1 + ah
    bx2, by2 = bx1 + bw, by1 + bh
    w = max(0, min(ax2, bx2) - max(ax1, bx1))
    h = max(0, min(ay2, by2) - max(ay1, by1))
    return w * h


def _collect_text_rectangles(slide, exclude_shape=None) -> List[tuple[int, int, int, int]]:
    rects: List[tuple[int, int, int, int]] = []
    for shape in slide.shapes:
        if exclude_shape is not None and shape == exclude_shape:
            continue
        if not _shape_has_text(shape):
            continue
        p_type = _placeholder_type_id(shape) if getattr(shape, "is_placeholder", False) else None
        if _is_title_placeholder_type(p_type) or _is_subtitle_placeholder_type(p_type):
            continue
        rect = _shape_rect(shape)
        if rect is not None:
            rects.append(rect)
    return rects


def _iter_candidate_image_boxes(slide, *, slide_type: str, image_intent: str):
    width, height = _slide_size_emu(slide)
    margin_x = int(width * 0.06)
    margin_top = int(height * 0.2)
    margin_bottom = int(height * 0.08)
    usable_w = max(1, width - (margin_x * 2))
    usable_h = max(1, height - margin_top - margin_bottom)

    def _clamp_box(box: tuple[int, int, int, int]) -> tuple[int, int, int, int]:
        left, top, box_w, box_h = box
        left = max(0, min(left, width - 1))
        top = max(0, min(top, height - 1))
        box_w = max(1, min(box_w, width - left))
        box_h = max(1, min(box_h, height - top))
        return left, top, box_w, box_h

    right_box = _clamp_box(
        (
            margin_x + int(usable_w * 0.52),
            margin_top,
            int(usable_w * 0.42),
            int(usable_h * 0.76),
        )
    )
    bottom_box = _clamp_box(
        (
            margin_x,
            margin_top + int(usable_h * 0.54),
            usable_w,
            int(usable_h * 0.38),
        )
    )
    center_box = _clamp_box(
        (
            margin_x + int(usable_w * 0.18),
            margin_top + int(usable_h * 0.26),
            int(usable_w * 0.64),
            int(usable_h * 0.54),
        )
    )
    hero_box = _clamp_box(
        (
            margin_x,
            margin_top + int(usable_h * 0.2),
            usable_w,
            int(usable_h * 0.68),
        )
    )

    intent = _normalize_image_intent(image_intent)
    slide_kind = _resolve_slide_type(slide_type)
    if slide_kind == "section":
        preferred = [hero_box, bottom_box, center_box]
    elif slide_kind in {"two_content", "comparison"}:
        preferred = [bottom_box, right_box, center_box]
    elif intent == "hero":
        preferred = [right_box, center_box, bottom_box]
    elif intent in {"data", "process"}:
        preferred = [bottom_box, right_box, center_box]
    else:
        preferred = [right_box, bottom_box, center_box]

    seen: set[tuple[int, int, int, int]] = set()
    for box in preferred:
        if _rect_area(box) <= 0:
            continue
        if box in seen:
            continue
        seen.add(box)
        yield box


def _pick_non_overlapping_image_box(
    slide,
    *,
    slide_type: str,
    image_intent: str,
    exclude_shape=None,
) -> tuple[int, int, int, int] | None:
    occupied = _collect_text_rectangles(slide, exclude_shape=exclude_shape)
    candidates = list(_iter_candidate_image_boxes(slide, slide_type=slide_type, image_intent=image_intent))
    if not candidates:
        return None
    if not occupied:
        return candidates[0]

    intent = _normalize_image_intent(image_intent)
    overlap_limit = 0.12 if intent == "hero" else 0.08
    best_box = None
    best_ratio = 1.0

    for box in candidates:
        area = _rect_area(box)
        if area <= 0:
            continue
        overlap = 0
        for used in occupied:
            overlap += _rect_intersection_area(box, used)
        ratio = overlap / area
        if ratio <= overlap_limit:
            return box
        if ratio < best_ratio:
            best_ratio = ratio
            best_box = box

    # Fallback for section slides: allow slightly larger overlap to keep visual richness.
    if _resolve_slide_type(slide_type) == "section" and best_box is not None and best_ratio <= 0.2:
        return best_box
    return None


def _apply_cover_crop(picture_shape, *, image_path: str, box_width: int, box_height: int):
    try:
        from PIL import Image
    except Exception:
        return

    if box_width <= 0 or box_height <= 0:
        return

    try:
        with Image.open(image_path) as img:
            src_width, src_height = img.size
    except Exception:
        return

    if not src_width or not src_height:
        return

    src_ratio = src_width / src_height
    box_ratio = box_width / box_height

    if src_ratio > box_ratio:
        crop = max(0.0, (1.0 - (box_ratio / src_ratio)) / 2.0)
        picture_shape.crop_left = crop
        picture_shape.crop_right = crop
    elif src_ratio < box_ratio:
        crop = max(0.0, (1.0 - (src_ratio / box_ratio)) / 2.0)
        picture_shape.crop_top = crop
        picture_shape.crop_bottom = crop


def _insert_generated_image(
    slide,
    image_path: str,
    exclude_shape=None,
    *,
    slide_type: str = "content",
    image_intent: str = "concept",
) -> bool:
    if not image_path or not os.path.isfile(image_path):
        return False

    target = _find_largest_image_frame(slide, exclude_shape=exclude_shape)

    if target is not None and getattr(target, "is_placeholder", False):
        if _is_image_placeholder_type(_placeholder_type_id(target)) and hasattr(target, "insert_picture"):
            try:
                target.insert_picture(image_path)
                return True
            except Exception as e:
                logger.warning("Insert picture via placeholder failed: %s", e)

    if target is None:
        fallback_box = _pick_non_overlapping_image_box(
            slide,
            slide_type=slide_type,
            image_intent=image_intent,
            exclude_shape=exclude_shape,
        )
        if fallback_box is None:
            return False
        left, top, width, height = fallback_box
    else:
        left = int(getattr(target, "left", 0))
        top = int(getattr(target, "top", 0))
        width = int(getattr(target, "width", 0))
        height = int(getattr(target, "height", 0))
        if width <= 0 or height <= 0:
            left, top, width, height = _slide_default_image_box(slide)

    try:
        picture = slide.shapes.add_picture(image_path, left, top, width=width, height=height)
        _apply_cover_crop(picture, image_path=image_path, box_width=width, box_height=height)
        return True
    except Exception as e:
        logger.warning("Insert generated image failed: %s", e)
        return False


def _generate_image_for_slide(slide_data: Dict[str, Any]) -> str | None:
    if TEXT2PPTX_IMAGE_MODE == "off":
        return None

    image_prompt = str(slide_data.get("image_prompt") or "").strip()
    image_required = _coerce_bool(slide_data.get("image_required"), default=bool(image_prompt))
    if not image_required or not image_prompt:
        return None

    aspect_ratio = _normalize_aspect_ratio(slide_data.get("aspect_ratio"))
    attempts = TEXT2PPTX_IMAGE_RETRY + 1
    for attempt in range(1, attempts + 1):
        try:
            result = generate_image(
                prompt=image_prompt,
                aspect_ratio=aspect_ratio,
                mode=TEXT2PPTX_IMAGE_MODE,
                timeout_sec=TEXT2PPTX_IMAGE_TIMEOUT_SEC,
                output_dir=GENERATED_IMAGE_DIR,
            )
            path = str(result.get("local_path") or "").strip()
            return path if path else None
        except ImageGenError as e:
            logger.warning(
                "Image generation failed code=%s retryable=%s attempt=%s/%s: %s",
                e.code,
                e.retryable,
                attempt,
                attempts,
                e,
            )
            if not e.retryable or attempt >= attempts:
                break
        except Exception as e:
            logger.warning("Image generation failed with unexpected error: %s", e)
            break
    return None


def _fill_bullets(text_frame, items: List[str]):
    text_frame.clear()
    if not items:
        p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
        p.text = "（無內容）"
        _apply_font_style(p, BODY_FONT_PT)
        return
    for idx, b in enumerate(items):
        p = text_frame.paragraphs[idx] if idx < len(text_frame.paragraphs) else text_frame.add_paragraph()
        p.text = b
        _apply_font_style(p, BODY_FONT_PT)


def _fill_column(text_frame, heading: str, items: List[str]):
    text_frame.clear()
    p0 = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    p0.text = heading
    _apply_font_style(p0, SUBTITLE_FONT_PT, bold=True)
    if not items:
        p1 = text_frame.add_paragraph()
        p1.text = "（無內容）"
        _apply_font_style(p1, BODY_FONT_PT)
        return
    for item in items:
        p = text_frame.add_paragraph()
        p.text = item
        _apply_font_style(p, BODY_FONT_PT)


def _split_two_columns(slide_data: Dict[str, Any], bullets: List[str], slide_type: str) -> tuple[str, List[str], str, List[str]]:
    left_title = str(slide_data.get("left_title") or "").strip()
    right_title = str(slide_data.get("right_title") or "").strip()
    if not left_title:
        left_title = "重點 A" if slide_type == "comparison" else "欄位 A"
    if not right_title:
        right_title = "重點 B" if slide_type == "comparison" else "欄位 B"

    left_items: List[str] = []
    right_items: List[str] = []
    for idx, b in enumerate(bullets):
        if idx % 2 == 0:
            left_items.append(b)
        else:
            right_items.append(b)
    return left_title, left_items, right_title, right_items


def _apply_font_style(paragraph, size_pt: int, bold: bool = False):
    """
    Apply font style with broad python-pptx compatibility.
    """
    from pptx.util import Pt
    from pptx.oxml.xmlchemy import OxmlElement

    if not paragraph.runs:
        paragraph.add_run()

    for run in paragraph.runs:
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.name = FONT_NAME_ASCII
        try:
            rPr = run._r.get_or_add_rPr()
            ea = rPr.first_child_found_in("a:ea")
            if ea is None:
                ea = OxmlElement("a:ea")
                ea.set("typeface", FONT_NAME_EAST_ASIAN)
                rPr.insert_element_before(
                    ea,
                    "a:cs",
                    "a:sym",
                    "a:hlinkClick",
                    "a:hlinkMouseOver",
                    "a:rtl",
                    "a:extLst",
                )
            else:
                ea.set("typeface", FONT_NAME_EAST_ASIAN)
        except Exception:
            pass


def _set_paragraph_text_preserving_style(paragraph, text: str) -> None:
    text = str(text or "")
    runs = list(getattr(paragraph, "runs", []))
    if not runs:
        try:
            run = paragraph.add_run()
            run.text = text
            return
        except Exception:
            try:
                paragraph.text = text
            except Exception:
                pass
            return

    first = True
    for run in runs:
        try:
            run.text = text if first else ""
        except Exception:
            pass
        first = False


def _set_text_frame_lines_preserving_style(text_frame, lines: List[str]) -> None:
    if text_frame is None:
        return

    items = [str(line or "").strip() for line in (lines or []) if str(line or "").strip()]
    paragraphs = list(getattr(text_frame, "paragraphs", []))
    if not paragraphs:
        try:
            paragraphs = [text_frame.add_paragraph()]
        except Exception:
            return

    for idx, paragraph in enumerate(paragraphs):
        text = items[idx] if idx < len(items) else ""
        _set_paragraph_text_preserving_style(paragraph, text)

    for idx in range(len(paragraphs), len(items)):
        try:
            paragraph = text_frame.add_paragraph()
        except Exception:
            break
        _set_paragraph_text_preserving_style(paragraph, items[idx])


def _extract_template_schema_from_template_text(template_text: str) -> Dict[str, Any] | None:
    for raw_line in str(template_text or "").splitlines():
        line = str(raw_line or "").strip()
        if not line.startswith("@template_schema "):
            continue
        payload = line[len("@template_schema ") :].strip()
        if not payload:
            return None
        try:
            schema = json.loads(payload)
        except Exception:
            return None
        return schema if isinstance(schema, dict) else None
    return None


def _restore_slide_text_from_analysis(slide, slide_data: Dict[str, Any], *, is_cover: bool = False) -> None:
    title = str(slide_data.get("title") or "").strip()
    bullets = slide_data.get("bullets") or []
    if not isinstance(bullets, list):
        bullets = [bullets]
    bullets = [str(item).strip() for item in bullets if str(item).strip()]

    title_shape = _find_title_shape(slide)
    if title_shape is not None:
        _set_text_frame_lines_preserving_style(title_shape.text_frame, [title])
    elif title:
        fallback_title = _find_largest_text_frame(slide)
        if fallback_title is not None:
            _set_text_frame_lines_preserving_style(fallback_title, [title])

    if is_cover:
        subtitle_text = str(slide_data.get("subtitle") or "").strip()
        content_lines = [line.strip() for line in subtitle_text.splitlines() if line.strip()]
        if not content_lines and subtitle_text:
            content_lines = [subtitle_text]
        content_tfs = _find_content_text_frames(slide, exclude_shape=title_shape)
        _fill_text_frames_with_lines(content_tfs, content_lines)
        return

    content_tfs = _find_content_text_frames(slide, exclude_shape=title_shape)
    if not content_tfs and bullets:
        fallback = _find_largest_text_frame(slide, exclude_shape=title_shape)
        if fallback is not None:
            _set_text_frame_lines_preserving_style(fallback, bullets)
        return

    _fill_text_frames_with_lines(content_tfs, bullets)


def _safe_filename(name: str, default: str = "slides") -> str:
    name = (name or "").strip() or default
    name = re.sub(r'[\\/:*?"<>|]+', "_", name)
    name = re.sub(r"\s+", "_", name).strip("_")
    return name or default


def _build_download_filename(raw_title: str, fallback: str = DEFAULT_MAIN_TITLE) -> str:
    base = (raw_title or "").strip() or fallback
    if base.lower().endswith(".pptx"):
        base = base[:-5]
    safe_base = _safe_filename(base, default=fallback)
    return f"{safe_base}.pptx"


def _coerce_llm_text_response(response: Any) -> str:
    if response is None:
        return ""
    if isinstance(response, str):
        return response
    content = getattr(response, "content", None)
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        parts: List[str] = []
        for item in content:
            if isinstance(item, dict):
                text = item.get("text")
                if text:
                    parts.append(str(text))
            elif isinstance(item, str):
                parts.append(item)
        if parts:
            return "\n".join(parts)
    return str(response)


def _extract_json_object(raw_text: str) -> str:
    text = (raw_text or "").strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*", "", text, flags=re.IGNORECASE)
        text = re.sub(r"\s*```$", "", text)
        text = text.strip()

    start = text.find("{")
    if start < 0:
        raise ValueError("LLM output does not contain JSON object start.")

    depth = 0
    in_string = False
    escaped = False
    for i in range(start, len(text)):
        ch = text[i]
        if in_string:
            if escaped:
                escaped = False
            elif ch == "\\":
                escaped = True
            elif ch == '"':
                in_string = False
            continue
        if ch == '"':
            in_string = True
        elif ch == "{":
            depth += 1
        elif ch == "}":
            depth -= 1
            if depth == 0:
                return text[start : i + 1]
    raise ValueError("LLM output JSON object is incomplete.")


def _fallback_from_raw_text(text: str) -> Dict[str, Any]:
    blocks = [s.strip() for s in text.split("===") if s.strip()]
    if not blocks:
        blocks = [text.strip()]
    slides = []
    for i, block in enumerate(blocks, start=1):
        lines = [ln.strip() for ln in block.splitlines() if ln.strip()]
        title = lines[0] if lines else f"投影片 {i}"
        bullets = lines[1:] if len(lines) > 1 else []
        slides.append(
            {
                "title": title,
                "slide_type": "content",
                "bullets": bullets,
                "image_required": False,
                "image_prompt": "",
                "image_intent": "",
                "aspect_ratio": "16:9",
            }
        )
    return {
        "main_title": DEFAULT_MAIN_TITLE,
        "main_subtitle": DEFAULT_MAIN_SUBTITLE,
        "slides": slides,
    }


def _normalize_analysis_result(data: Dict[str, Any], *, source_text: str) -> Dict[str, Any]:
    title = str(data.get("main_title") or "").strip() or DEFAULT_MAIN_TITLE
    subtitle = str(data.get("main_subtitle") or "").strip() or DEFAULT_MAIN_SUBTITLE

    raw_slides = data.get("slides")
    if not isinstance(raw_slides, list):
        return _fallback_from_raw_text(source_text)

    slides = []
    for idx, s in enumerate(raw_slides, start=1):
        if not isinstance(s, dict):
            continue
        slide_title = str(s.get("title") or "").strip() or f"投影片 {idx}"
        slide_type = _resolve_slide_type(s.get("slide_type"))
        left_title = str(s.get("left_title") or "").strip()
        right_title = str(s.get("right_title") or "").strip()
        image_prompt = str(s.get("image_prompt") or "").strip()
        image_intent = _normalize_image_intent(s.get("image_intent"))
        aspect_ratio = _normalize_aspect_ratio(s.get("aspect_ratio"))
        image_required = _coerce_bool(s.get("image_required"), default=bool(image_prompt))
        if not image_prompt:
            image_required = False
        bullets_raw = s.get("bullets")
        bullets: List[str] = []
        if isinstance(bullets_raw, list):
            for b in bullets_raw:
                v = str(b).strip()
                if v:
                    bullets.append(v)
        elif bullets_raw is not None:
            v = str(bullets_raw).strip()
            if v:
                bullets.append(v)
        slides.append(
            {
                "title": slide_title,
                "layout_name": str(s.get("layout_name") or "").strip(),
                "slide_type": slide_type,
                "left_title": left_title,
                "right_title": right_title,
                "bullets": bullets,
                "image_required": image_required,
                "image_prompt": image_prompt,
                "image_intent": image_intent,
                "aspect_ratio": aspect_ratio,
            }
        )

    if not slides:
        return _fallback_from_raw_text(source_text)

    return {
        "main_title": title,
        "main_subtitle": subtitle,
        "main_layout_name": str(data.get("main_layout_name") or "").strip(),
        "slides": slides,
    }


def _analyze_text_with_llm(text: str) -> Dict[str, Any]:
    """
    Uses LLM to analyze raw text and structure it for presentation.
    Returns a dict with 'title', 'subtitle', and 'slides' (list of dicts with 'title', 'bullets').
    """
    llm = get_chat_model()
    prompt = f"""
    你是一個專業的簡報生成助手。請將以下文字內容轉換為結構化的簡報大綱。輸出必須是 JSON 格式。

    規則：
    1. 輸出 JSON 必須包含三個頂層鍵："main_title", "main_subtitle", "slides"。
    2. "main_title"：整份簡報的主要標題。
    3. "main_subtitle"：整份簡報的副標題。
    4. "slides"：一個列表，每個元素代表一張投影片。
    5. 每張投影片（slides 的每個元素）必須包含三個鍵："title"、"slide_type" 和 "bullets"。
    6. "title"：該投影片的標題。
    7. "slide_type"：限定為 "content"、"section"、"two_content"、"comparison" 其中之一。
    8. "bullets"：一個列表，包含該投影片的所有條列要點。如果沒有要點，則為空列表。
    9. 若 slide_type 為 "comparison" 或 "two_content"，可額外提供 "left_title" 與 "right_title"。
    10. 盡可能語義化地拆分內容到不同的投影片，每張投影片的條列要點數量不應超過 {MAX_BULLETS_PER_SLIDE} 點。
    11. 若原文沒有明確主標題，請自動生成一個簡潔相關的標題。
    12. 若原文沒有明確副標題，請使用「文字內容自動生成簡報」作為副標題。
    13. 每張投影片都要額外輸出圖片欄位：
       - "image_required": 布林值，是否需要插圖。
       - "image_prompt": 用來生圖的簡潔描述，若不需要圖片請給空字串。
       - "image_intent": "concept"、"data"、"process"、"hero" 其一。
       - "aspect_ratio": 優先使用 "16:9"。

    請分析以下文字內容：

    ```text
    {text}
    ```

    JSON 輸出範例：
    ```json
    {{
      "main_title": "專案進度報告",
      "main_subtitle": "2026年3月更新",
      "slides": [
        {{
          "title": "第一張：專案概述",
          "slide_type": "section",
          "bullets": ["目的：提升效率", "範圍：XXX"],
          "image_required": true,
          "image_prompt": "團隊在會議室討論專案藍圖，乾淨的企業風格插圖",
          "image_intent": "hero",
          "aspect_ratio": "16:9"
        }},
        {{
          "title": "第二張：時程規劃",
          "slide_type": "comparison",
          "left_title": "已完成",
          "right_title": "待完成",
          "bullets": ["W1：需求分析", "W2：開發階段", "W3：驗收與部署"],
          "image_required": false,
          "image_prompt": "",
          "image_intent": "data",
          "aspect_ratio": "16:9"
        }}
      ]
    }}
    ```
    """

    try:
        response = llm.invoke(prompt)
        raw_text = _coerce_llm_text_response(response)
        json_text = _extract_json_object(raw_text)
        llm_output = json.loads(json_text)
        if not isinstance(llm_output, dict):
            raise ValueError("LLM JSON root must be object.")
        return _normalize_analysis_result(llm_output, source_text=text)
    except Exception as e:
        logger.warning("LLM text analysis failed for text2pptx: %s", e)
        return _fallback_from_raw_text(text)


def _clean_extracted_line(value: Any) -> str:
    text = str(value or "")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"\s+", " ", text).strip()
    text = _strip_page_number_prefix(text)
    if _is_page_number_only(text):
        return ""
    return text


def _strip_page_number_prefix(text: str) -> str:
    t = str(text or "").strip()
    if not t:
        return ""
    t = re.sub(
        r"^\s*page\s*[0-9０-９]+(?:\s*(?:/|of)\s*[0-9０-９]+)?\s*[-:：|·.)）]*\s*",
        "",
        t,
        flags=re.IGNORECASE,
    )
    t = re.sub(
        r"^\s*第\s*[0-9０-９一二三四五六七八九十百千]+\s*頁(?:\s*/\s*[0-9０-９一二三四五六七八九十百千]+\s*頁?)?\s*[-:：|·.)）]*\s*",
        "",
        t,
    )
    return t.strip()


def _is_page_number_only(text: str) -> bool:
    t = str(text or "").strip()
    if not t:
        return True
    if re.fullmatch(r"page\s*[0-9０-９]+(?:\s*(?:/|of)\s*[0-9０-９]+)?", t, flags=re.IGNORECASE):
        return True
    if re.fullmatch(r"第\s*[0-9０-９一二三四五六七八九十百千]+\s*頁(?:\s*/\s*[0-9０-９一二三四五六七八九十百千]+\s*頁?)?", t):
        return True
    return False


def _is_pure_small_number_line(text: str) -> bool:
    t = str(text or "").strip()
    if not t:
        return False
    t = t.replace(",", "").replace("，", "")
    t = t.translate(str.maketrans("０１２３４５６７８９", "0123456789"))
    return bool(re.fullmatch(r"\d{1,3}", t))


def _looks_like_footer_context_line(text: str) -> bool:
    t = str(text or "").strip()
    if not t:
        return False
    if _is_page_number_only(t):
        return True
    if re.search(r"(?:19\d{2}|20\d{2}|20XX)", t, flags=re.IGNORECASE):
        return True
    if re.search(r"(?:民國)\s*\d+\s*年", t):
        return True
    if re.search(r"(?:簡報|投影片|presentation|pitch|deck|slide)", t, flags=re.IGNORECASE):
        return True
    return False


def _remove_footer_numeric_page_lines(lines: List[str]) -> List[str]:
    if not lines:
        return []
    out: List[str] = []
    total = len(lines)
    for idx, line in enumerate(lines):
        if not _is_pure_small_number_line(line):
            out.append(line)
            continue
        if idx < max(0, total - 3):
            out.append(line)
            continue
        neighbors = lines[max(0, idx - 2) : idx] + lines[idx + 1 : min(total, idx + 2)]
        if any(_looks_like_footer_context_line(x) for x in neighbors):
            continue
        out.append(line)
    return out


def _shape_position_key(shape: Any) -> tuple[int, int]:
    try:
        return (int(getattr(shape, "top", 0)), int(getattr(shape, "left", 0)))
    except Exception:
        return (0, 0)


def _extract_text_lines_from_shape(shape: Any) -> List[str]:
    lines: List[str] = []
    if getattr(shape, "has_text_frame", False):
        text = getattr(shape.text_frame, "text", "")
        for row in str(text or "").splitlines():
            cleaned = _clean_extracted_line(row)
            if cleaned:
                lines.append(cleaned)

    # Some slide objects may raise "shape does not contain a table"
    # when reading shape.table even after has_table checks.
    try:
        has_table = bool(getattr(shape, "has_table", False))
    except Exception:
        has_table = False
    if has_table:
        try:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    for raw in str(getattr(cell, "text", "") or "").splitlines():
                        cleaned = _clean_extracted_line(raw)
                        if cleaned:
                            lines.append(cleaned)
        except Exception:
            # Keep extracting from other shapes.
            pass
    return lines


def _extract_text_lines_from_slide(slide: Any) -> List[str]:
    ordered_shapes = sorted(list(getattr(slide, "shapes", [])), key=_shape_position_key)
    lines: List[str] = []
    for shape in ordered_shapes:
        for line in _extract_text_lines_from_shape(shape):
            if lines and lines[-1] == line:
                continue
            lines.append(line)
    return _remove_footer_numeric_page_lines(lines)


def _marker_to_layout_label(marker: str) -> str:
    labels = {
        "cover": "\u6a19\u984c\u6295\u5f71\u7247",
        "section": "\u7ae0\u7bc0\u6a19\u984c",
        "content": "\u5167\u5bb9\u9801",
        "two_content": "\u96d9\u6b04\u5fc5\u8f03\u9801",
        "comparison": "\u96d9\u6b04\u5fc5\u8f03\u9801",
    }
    return labels.get(marker, labels["content"])


def _extract_slide_layout_name(slide: Any) -> str:
    try:
        name = str(getattr(getattr(slide, "slide_layout", None), "name", "") or "").strip()
    except Exception:
        name = ""
    return _clean_extracted_line(name)


def _detect_layout_marker_from_slide(slide: Any, *, extracted_index: int) -> str | None:
    if extracted_index == 1:
        return "cover"

    ph_types: List[int | None] = []
    for shape in getattr(slide, "shapes", []):
        if not getattr(shape, "is_placeholder", False):
            continue
        ph_types.append(_placeholder_type_id(shape))

    has_title = any(_is_title_placeholder_type(t) for t in ph_types)
    has_subtitle = any(_is_subtitle_placeholder_type(t) for t in ph_types)
    content_count = sum(1 for t in ph_types if _is_content_placeholder_type(t))

    if has_title and content_count >= 2:
        return "two_content"
    if has_title and has_subtitle and content_count == 0:
        return "section"
    if has_title and content_count >= 1:
        return "content"
    return None


def _pick_marker_for_extracted_slide(
    *,
    extracted_index: int,
    title: str,
    body_lines: List[str],
) -> str:
    if extracted_index == 1:
        return "cover"

    low_title = title.lower()
    combined = f"{title} {' '.join(body_lines)}"
    low_combined = combined.lower()

    comparison_keys = (
        "compare",
        "comparison",
        "vs",
        "\u6bd4\u8f03",
        "\u5c0d\u6bd4",
        "\u5dee\u7570",
        "\u65b9\u6848a",
        "\u65b9\u6848b",
    )
    section_keys = (
        "overview",
        "summary",
        "objective",
        "conclusion",
        "decision",
        "\u80cc\u666f",
        "\u6458\u8981",
        "\u7e3d\u7d50",
        "\u7d50\u8ad6",
        "\u76ee\u6a19",
        "\u6c7a\u7b56",
        "\u8acb\u6c42",
    )

    if any(k in low_title or k in low_combined for k in comparison_keys) or any(k in combined for k in comparison_keys):
        return "comparison"
    if len(body_lines) >= (MAX_BULLETS_PER_SLIDE + 3):
        return "two_content"
    if any(k in low_title or k in low_combined for k in section_keys) or any(k in combined for k in section_keys):
        return "section"
    return "content"

def _extract_sample_text_from_pptx_bytes(raw: bytes) -> Dict[str, Any]:
    from pptx import Presentation

    try:
        from webapps.text2pptx.pptx2schema.pipelines.extract_pipeline import run_extract
    except Exception:
        run_extract = None

    presentation = Presentation(io.BytesIO(raw))
    total_slides = len(presentation.slides)
    if total_slides == 0:
        raise ValueError("This PPTX file has no usable slides.")

    presentation_raw = None
    if run_extract is not None:
        try:
            presentation_raw = run_extract(raw, source_file=None)
        except Exception:
            presentation_raw = None

    blocks: List[str] = []
    extracted_count = 0
    for slide_index, slide in enumerate(presentation.slides, start=1):
        if extracted_count >= MAX_SAMPLE_SLIDES:
            break

        lines = _extract_text_lines_from_slide(slide)
        if not lines:
            continue

        title = lines[0]
        body_lines = lines[1 : 1 + (MAX_BULLETS_PER_SLIDE * 2)]
        extracted_count += 1

        heuristic_marker = _pick_marker_for_extracted_slide(
            extracted_index=extracted_count,
            title=title,
            body_lines=body_lines,
        )
        layout_marker = _detect_layout_marker_from_slide(slide, extracted_index=extracted_count)

        marker = layout_marker or heuristic_marker
        if layout_marker == "content" and heuristic_marker in {"comparison", "section"}:
            marker = heuristic_marker
        if layout_marker == "two_content" and heuristic_marker == "comparison":
            marker = "comparison"

        # Align sample labels with the extracted master numbering, not the
        # source slide layout identity.
        layout_name = f"\u7248\u578b{extracted_count}"

        block_lines = [f"[{layout_name}]{title}"]
        block_lines.extend(body_lines)

        blocks.append("\n".join(block_lines))

    if not blocks:
        raise ValueError("No extractable slide text was found.")

    sample_text = "\n===\n".join(blocks).strip()
    if len(sample_text) > MAX_SAMPLE_CHARS:
        sample_text = sample_text[:MAX_SAMPLE_CHARS].rstrip() + "\n...(content truncated)"

    sample_schema = None
    if presentation_raw is not None:
        try:
            sample_schema = presentation_raw.model_dump()
        except Exception:
            sample_schema = None

    return {
        "sample_text": sample_text,
        "sample_schema": sample_schema,
        "total_slides": total_slides,
        "used_slides": extracted_count,
    }


def _extract_master_template_text_from_pptx_bytes(raw: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(raw))
    referenced_layout_parts = _collect_referenced_layout_parts_from_presentation(presentation)
    layout_name_map = _layout_part_display_name_map(presentation, keep_layout_parts=referenced_layout_parts)
    lines: List[str] = []
    lines.append("[master structure]")

    master_count = 0
    for idx, master in enumerate(presentation.slide_masters, start=1):
        layout_lines: List[str] = []
        for layout in getattr(master, "slide_layouts", []):
            try:
                partname = str(getattr(layout.part, "partname", "") or "").strip().lstrip("/")
            except Exception:
                partname = ""
            if not partname or partname not in referenced_layout_parts:
                continue
            display_name = layout_name_map.get(partname, "")
            if not display_name:
                display_name = f"版型{len(layout_lines) + 1}"
            layout_lines.append(display_name)

        if layout_lines:
            master_count += 1
            master_name = _clean_extracted_line(getattr(master, "name", "")) or f"master {idx}"
            lines.append(f"master {idx}: {master_name}")
            for layout_name in layout_lines:
                lines.append(f"- {layout_name}")

    if master_count == 0:
        lines.append("No slide masters found.")

    try:
        template_schema = extract_template_structure(
            presentation,
            keep_layout_parts=referenced_layout_parts,
        )
    except Exception:
        template_schema = None

    if template_schema is not None:
        try:
            lines.append(
                "@template_schema "
                + json.dumps(template_schema, ensure_ascii=False, separators=(",", ":"))
            )
        except Exception:
            pass

    lines.append("===")
    lines.append("[slides]")
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        try:
            slide_layout = getattr(slide, "slide_layout", None)
            partname = str(getattr(getattr(slide_layout, "part", None), "partname", "") or "").strip().lstrip("/")
        except Exception:
            partname = ""
        if not partname or partname not in referenced_layout_parts:
            continue
        layout_name = layout_name_map.get(partname) or f"版型{slide_idx}"
        lines.append(f"slide {slide_idx}: {layout_name}")

    template_text = "\n".join(lines).strip()
    if len(template_text) > MAX_SAMPLE_CHARS:
        template_text = template_text[:MAX_SAMPLE_CHARS].rstrip() + "\n...(content truncated)"
    return template_text


def _convert_pptx_bytes_to_pptx_bytes(raw: bytes) -> bytes:
    from pptx import Presentation

    content_types_name = "[Content_Types].xml"
    ns = {"ct": "http://schemas.openxmlformats.org/package/2006/content-types"}

    # Keep total slide count, clear text, and move
    # representative-layout shapes (picture/group/freeform/textbox...) into layout parts.
    normalized_raw = raw
    try:
        src_prs = Presentation(io.BytesIO(raw))
        layout_occurrences = _collect_layout_occurrences(src_prs)

        prs = Presentation(io.BytesIO(raw))
        _retain_representative_slides_and_clear_text(prs)
        normalized_buf = io.BytesIO()
        prs.save(normalized_buf)
        normalized_raw = normalized_buf.getvalue()

        if layout_occurrences:
            normalized_prs = Presentation(io.BytesIO(normalized_raw))
            layout_part_map = _layout_identity_to_partname_map(normalized_prs)
            normalized_slide_parts = [
                str(slide.part.partname or "").lstrip("/")
                for slide in normalized_prs.slides
            ]
            for occs in layout_occurrences.values():
                for occ in occs:
                    idx = int(occ.get("slide_index", -1))
                    occ["normalized_slide_part"] = (
                        normalized_slide_parts[idx] if 0 <= idx < len(normalized_slide_parts) else ""
                    )
            normalized_raw = _inject_layout_representative_shapes(
                raw,
                normalized_raw,
                layout_occurrences=layout_occurrences,
                layout_part_map=layout_part_map,
            )
            try:
                with zipfile.ZipFile(io.BytesIO(normalized_raw), "r") as zin:
                    file_map = {name: zin.read(name) for name in zin.namelist()}
                _strip_nonplaceholder_shapes_from_slide_parts(file_map)
                out = io.BytesIO()
                with zipfile.ZipFile(out, "w", compression=zipfile.ZIP_DEFLATED) as zout:
                    for name, raw_data in file_map.items():
                        zout.writestr(name, raw_data)
                normalized_raw = out.getvalue()
            except Exception as e:
                logger.warning("strip non-placeholder slide shapes failed: %s", e)
    except Exception as e:
        logger.warning("normalize before potx conversion failed: %s", e)

    try:
        with zipfile.ZipFile(io.BytesIO(normalized_raw), "r") as zin:
            file_map: Dict[str, bytes] = {name: zin.read(name) for name in zin.namelist()}
        _prune_unused_master_layout_parts(file_map)
        out = io.BytesIO()
        with zipfile.ZipFile(out, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            for name, raw_data in file_map.items():
                zout.writestr(name, raw_data)
        return out.getvalue()
    except Exception as e:
        logger.warning("prune unused master layouts failed: %s", e)

    src = io.BytesIO(normalized_raw)
    out = io.BytesIO()
    with zipfile.ZipFile(src, "r") as zin, zipfile.ZipFile(out, "w", compression=zipfile.ZIP_DEFLATED) as zout:
        for info in zin.infolist():
            data = zin.read(info.filename)
            if info.filename == content_types_name:
                try:
                    root = ET.fromstring(data)
                    updated = False
                    for override in root.findall("ct:Override", ns):
                        if override.get("PartName") == "/ppt/presentation.xml":
                            if override.get("ContentType") != PPTX_MAIN_CONTENT_TYPE:
                                override.set("ContentType", PPTX_MAIN_CONTENT_TYPE)
                                updated = True
                    if updated:
                        data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                except Exception as e:
                    logger.warning("convert pptx content-type adjust failed: %s", e)
            zout.writestr(info, data)
    return out.getvalue()


def _convert_pptx_bytes_to_potx_bytes(raw: bytes) -> bytes:
    return _convert_pptx_bytes_to_pptx_bytes(raw)


def _save_generated_pptx_bytes(
    source_filename: str,
    raw_pptx: bytes,
    *,
    script_name: str = "",
) -> tuple[str, str]:
    source_stem = os.path.splitext(os.path.basename(source_filename or ""))[0] or "template"
    safe_stem = _safe_filename(source_stem, default="template")
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_filename = f"{safe_stem}_master_{ts}.pptx"

    rel_dir = os.path.join("text2pptx", "pptx")
    abs_dir = os.path.join(settings.MEDIA_ROOT, rel_dir)
    os.makedirs(abs_dir, exist_ok=True)
    abs_path = os.path.join(abs_dir, output_filename)
    with open(abs_path, "wb") as f:
        f.write(raw_pptx)

    media_url = str(getattr(settings, "MEDIA_URL", "/media/") or "/media/")
    if not media_url.endswith("/"):
        media_url += "/"
    if not media_url.startswith("/"):
        media_url = "/" + media_url

    script = (script_name or "").strip()
    if script.endswith("/"):
        script = script[:-1]

    rel_path = "/".join(["text2pptx", "pptx", output_filename])
    download_url = f"{script}{media_url}{quote(rel_path)}"
    return output_filename, download_url


def _save_generated_potx_bytes(
    source_filename: str,
    raw_potx: bytes,
    *,
    script_name: str = "",
) -> tuple[str, str]:
    return _save_generated_pptx_bytes(source_filename, raw_potx, script_name=script_name)


def _resolve_saved_pptx_path(filename: str) -> str:
    safe_name = _safe_filename(filename or "", default="")
    if not safe_name or not safe_name.lower().endswith(".pptx"):
        raise ValueError("invalid pptx filename")
    abs_path = os.path.join(settings.MEDIA_ROOT, "text2pptx", "pptx", safe_name)
    abs_path_norm = os.path.normpath(abs_path)
    base_dir = os.path.normpath(os.path.join(settings.MEDIA_ROOT, "text2pptx", "pptx"))
    if not abs_path_norm.startswith(base_dir):
        raise ValueError("invalid pptx path")
    return abs_path_norm


def _resolve_saved_potx_path(filename: str) -> str:
    return _resolve_saved_pptx_path(filename)


def _build_template_context() -> Dict[str, Any]:
    templates = _list_pptx_templates()
    ignored_templates = _list_ignored_template_files()
    default_template_available = DEFAULT_TEMPLATE_NAME in templates
    return {
        "templates": templates,
        "template_count": len(templates),
        "ignored_templates": ignored_templates,
        "ignored_count": len(ignored_templates),
        "default_template_name": DEFAULT_TEMPLATE_NAME,
        "default_template_available": default_template_available,
    }


def _resolve_analysis_for_image_prompts(text: str) -> Dict[str, Any]:
    marked_result = _parse_marked_text_structure(text)
    normalized_marked = (
        _normalize_analysis_result(marked_result, source_text=text) if marked_result is not None else None
    )

    llm_result = _analyze_text_with_llm(text)
    llm_slides = llm_result.get("slides") if isinstance(llm_result, dict) else []
    llm_has_prompt = False
    if isinstance(llm_slides, list):
        for slide in llm_slides:
            if not isinstance(slide, dict):
                continue
            if str(slide.get("image_prompt") or "").strip():
                llm_has_prompt = True
                break

    if llm_has_prompt:
        return _autofill_image_prompts_if_all_empty(llm_result)
    if normalized_marked is not None:
        return _autofill_image_prompts_if_all_empty(normalized_marked)
    return _autofill_image_prompts_if_all_empty(llm_result)


def _image_prompt_preview_lines(analysis_result: Dict[str, Any]) -> List[str]:
    slides = analysis_result.get("slides") if isinstance(analysis_result, dict) else []
    if not isinstance(slides, list):
        return []

    lines: List[str] = []
    for idx, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            continue
        title = str(slide.get("title") or f"Slide {idx}").strip()
        slide_type = _resolve_slide_type(slide.get("slide_type"))
        image_prompt = str(slide.get("image_prompt") or "").strip()
        image_intent = _normalize_image_intent(slide.get("image_intent"))
        aspect_ratio = _normalize_aspect_ratio(slide.get("aspect_ratio"))
        image_required = _coerce_bool(slide.get("image_required"), default=bool(image_prompt))
        if not image_prompt:
            image_required = False

        lines.append(f"[{idx}] {title} ({slide_type})")
        lines.append(f"required: {'yes' if image_required else 'no'}")
        lines.append(f"intent: {image_intent} | aspect_ratio: {aspect_ratio}")
        lines.append(f"prompt: {image_prompt or '(empty)'}")
        lines.append("")
    return lines


@require_node("pptx")
def index(request):
    return render(request, "text2pptx/index.html", _build_template_context())


@require_node("pptx")
def sample_extractor(request):
    context = _build_template_context()
    context.update(
        {
            "sample_text": "",
            "extract_meta": None,
            "uploaded_filename": "",
        }
    )

    if request.method == "POST":
        upload = request.FILES.get("pptx_file")
        if not upload:
            messages.error(request, "\u8acb\u5148\u9078\u64c7 .pptx \u6a94\u6848\u3002")
            return render(request, "text2pptx/sample_extractor.html", context)

        filename = _normalize_template_name(upload.name or "")
        context["uploaded_filename"] = filename
        if not filename.lower().endswith(".pptx"):
            messages.error(request, "\u50c5\u652f\u63f4 .pptx \u6a94\u6848\u3002")
            return render(request, "text2pptx/sample_extractor.html", context)

        if upload.size and upload.size > (MAX_SAMPLE_UPLOAD_MB * 1024 * 1024):
            messages.error(request, f"\u6a94\u6848\u5927\u5c0f\u4e0d\u53ef\u8d85\u904e {MAX_SAMPLE_UPLOAD_MB} MB\u3002")
            return render(request, "text2pptx/sample_extractor.html", context)

        try:
            raw = upload.read()
            extracted = _extract_sample_text_from_pptx_bytes(raw)
            context["sample_text"] = str(extracted.get("sample_text") or "")
            context["extract_meta"] = {
                "total_slides": int(extracted.get("total_slides") or 0),
                "used_slides": int(extracted.get("used_slides") or 0),
            }
            messages.success(request, "\u62bd\u53d6\u5b8c\u6210\uff0c\u53ef\u76f4\u63a5\u8907\u88fd\u6216\u5e36\u56de\u4e3b\u9801\u4f7f\u7528\u3002")
        except ValueError as e:
            raw_error = str(e or "").strip()
            if "shape does not contain a table" in raw_error.lower():
                messages.error(request, "\u62bd\u53d6\u5931\u6557\uff1a\u5075\u6e2c\u5230\u975e\u8868\u683c\u7269\u4ef6\uff0c\u5df2\u7565\u904e\u8a72\u7269\u4ef6\u3002\u8acb\u91cd\u8a66\u6216\u6539\u7528\u91cd\u65b0\u532f\u51fa\u7684 PPTX\u3002")
            else:
                messages.error(request, f"\u62bd\u53d6\u5931\u6557\uff1a{e}")
        except Exception as e:
            logger.warning("Sample extractor failed: %s", e)
            messages.error(request, "\u62bd\u53d6\u5931\u6557\uff1a\u8acb\u78ba\u8a8d\u6a94\u6848\u5167\u5bb9\u53ef\u8b80\u53d6\uff0c\u6216\u91cd\u65b0\u532f\u51fa\u5f8c\u518d\u8a66\u3002")

    return render(request, "text2pptx/sample_extractor.html", context)


@require_node("pptx", api=True)
def extract_sample_api(request):
    if request.method != "POST":
        return JsonResponse(
            {
                "success": False,
                "message": "method not allowed",
                "sample_text": "",
                "extract_meta": {},
            },
            status=405,
        )

    upload = request.FILES.get("pptx_file") or request.FILES.get("file")
    if not upload:
        return JsonResponse(
            {
                "success": False,
                "message": "請先選擇 .pptx 檔案。",
                "sample_text": "",
                "extract_meta": {},
            },
            status=400,
        )

    filename = _normalize_template_name(upload.name or "")
    if not filename.lower().endswith(".pptx"):
        return JsonResponse(
            {
                "success": False,
                "message": "僅支援 .pptx 檔案。",
                "sample_text": "",
                "extract_meta": {},
            },
            status=400,
        )

    if upload.size and upload.size > (MAX_SAMPLE_UPLOAD_MB * 1024 * 1024):
        return JsonResponse(
            {
                "success": False,
                "message": f"檔案大小不可超過 {MAX_SAMPLE_UPLOAD_MB} MB。",
                "sample_text": "",
                "extract_meta": {},
            },
            status=400,
        )

    try:
        raw = upload.read()
        extracted = _extract_sample_text_from_pptx_bytes(raw)
        return JsonResponse(
            {
                "success": True,
                "message": "抽取完成",
                "sample_text": str(extracted.get("sample_text") or ""),
                "extract_meta": {
                    "total_slides": int(extracted.get("total_slides") or 0),
                    "used_slides": int(extracted.get("used_slides") or 0),
                },
            },
            status=200,
        )
    except Exception as e:
        logger.warning("extract_sample_api failed: %s", e)
        return JsonResponse(
            {
                "success": False,
                "message": f"抽取失敗：{e}",
                "sample_text": "",
                "extract_meta": {},
            },
            status=500,
        )


@require_node("pptx", api=True)
def extract_template(request):
    if request.method != "POST":
        return JsonResponse(
            {
                "success": False,
                "message": "method not allowed",
                "template_text": "",
                "output_filename": "",
                "download_url": "",
            },
            status=405,
        )

    upload = request.FILES.get("pptx_file") or request.FILES.get("file")
    if not upload:
        return JsonResponse(
            {
                "success": False,
                "message": "請先選擇 PPTX 檔案",
                "template_text": "",
                "output_filename": "",
                "download_url": "",
            },
            status=400,
        )

    filename = _normalize_template_name(upload.name or "")
    if not filename.lower().endswith(".pptx"):
        return JsonResponse(
            {
                "success": False,
                "message": "只支援 .pptx 檔案",
                "template_text": "",
                "output_filename": "",
                "download_url": "",
            },
            status=400,
        )

    if upload.size and upload.size > (MAX_SAMPLE_UPLOAD_MB * 1024 * 1024):
        return JsonResponse(
            {
                "success": False,
                "message": f"檔案大小不可超過 {MAX_SAMPLE_UPLOAD_MB} MB",
                "template_text": "",
                "output_filename": "",
                "download_url": "",
            },
            status=400,
        )

    try:
        raw_pptx = upload.read()
        template_text = _extract_master_template_text_from_pptx_bytes(raw_pptx)
        pptx_bytes = _convert_pptx_bytes_to_pptx_bytes(raw_pptx)
        script_name = getattr(request, "script_name", "") or request.META.get("SCRIPT_NAME", "")
        output_filename, media_download_url = _save_generated_pptx_bytes(
            filename,
            pptx_bytes,
            script_name=script_name,
        )
        script = (script_name or "").rstrip("/")
        api_download_url = f"{script}/text2pptx/download-pptx/{quote(output_filename)}"
        return JsonResponse(
            {
                "success": True,
                "message": "母片模板抽取成功",
                "template_text": template_text,
                "output_filename": output_filename,
                "template_pptx_filename": output_filename,
                "download_url": api_download_url,
                "media_download_url": media_download_url,
            },
            status=200,
        )
    except Exception as e:
        logger.warning("extract_template failed: %s", e)
        return JsonResponse(
            {
                "success": False,
                "message": f"母片模板抽取失敗：{e}",
                "template_text": "",
                "output_filename": "",
                "download_url": "",
            },
            status=500,
        )


def _iter_pptx_files_in_dir(source_dir: str, recursive: bool = True) -> List[str]:
    source_dir = os.path.abspath(str(source_dir or "").strip())
    if not os.path.isdir(source_dir):
        return []

    files: List[str] = []
    if recursive:
        for root, _dirs, names in os.walk(source_dir):
            for name in names:
                if str(name).lower().endswith(".pptx"):
                    files.append(os.path.join(root, name))
    else:
        for name in os.listdir(source_dir):
            full = os.path.join(source_dir, name)
            if os.path.isfile(full) and str(name).lower().endswith(".pptx"):
                files.append(full)
    files.sort(key=lambda p: str(p).lower())
    return files


@require_node("pptx", api=True)
def extract_template_batch(request):
    if request.method != "POST":
        return JsonResponse(
            {
                "success": False,
                "message": "method not allowed",
                "processed": 0,
                "success_count": 0,
                "failed_count": 0,
                "results": [],
            },
            status=405,
        )

    try:
        payload = json.loads((request.body or b"").decode("utf-8") or "{}")
        if not isinstance(payload, dict):
            payload = {}
    except Exception:
        payload = {}

    source_dir = os.path.abspath(str(payload.get("source_dir") or "").strip())
    output_dir = os.path.abspath(str(payload.get("output_dir") or "").strip())
    recursive = bool(payload.get("recursive", True))

    if not source_dir:
        return JsonResponse(
            {
                "success": False,
                "message": "source_dir required",
                "processed": 0,
                "success_count": 0,
                "failed_count": 0,
                "results": [],
            },
            status=400,
        )
    if not os.path.isdir(source_dir):
        return JsonResponse(
            {
                "success": False,
                "message": f"source_dir not found: {source_dir}",
                "processed": 0,
                "success_count": 0,
                "failed_count": 0,
                "results": [],
            },
            status=400,
        )

    if not output_dir:
        output_dir = os.path.join(source_dir, "_extracted")
    try:
        os.makedirs(output_dir, exist_ok=True)
    except Exception as e:
        return JsonResponse(
            {
                "success": False,
                "message": f"failed to create output_dir: {e}",
                "processed": 0,
                "success_count": 0,
                "failed_count": 0,
                "results": [],
            },
            status=400,
        )

    pptx_files = _iter_pptx_files_in_dir(source_dir, recursive=recursive)
    if not pptx_files:
        return JsonResponse(
            {
                "success": True,
                "message": "no pptx files found",
                "processed": 0,
                "success_count": 0,
                "failed_count": 0,
                "output_dir": output_dir,
                "results": [],
            },
            status=200,
        )

    results: List[Dict[str, Any]] = []
    success_count = 0
    failed_count = 0
    max_bytes = MAX_SAMPLE_UPLOAD_MB * 1024 * 1024

    for idx, file_path in enumerate(pptx_files, start=1):
        filename = os.path.basename(file_path)
        rel_path = os.path.relpath(file_path, source_dir)
        safe_stem = _safe_filename(os.path.splitext(filename)[0], default=f"pptx_{idx}")
        item_result: Dict[str, Any] = {
            "index": idx,
            "source_file": file_path,
            "source_relpath": rel_path,
            "filename": filename,
            "ok": False,
            "sample_txt_path": "",
            "master_txt_path": "",
            "master_pptx_path": "",
            "error": "",
        }

        try:
            file_size = os.path.getsize(file_path)
            if file_size > max_bytes:
                raise ValueError(f"file too large (> {MAX_SAMPLE_UPLOAD_MB} MB)")

            with open(file_path, "rb") as f:
                raw = f.read()

            extracted = _extract_sample_text_from_pptx_bytes(raw)
            sample_text = str(extracted.get("sample_text") or "").strip()
            master_text = _extract_master_template_text_from_pptx_bytes(raw)
            master_pptx_bytes = _convert_pptx_bytes_to_pptx_bytes(raw)

            sample_txt_path = os.path.join(output_dir, f"{safe_stem}_sample.txt")
            master_txt_path = os.path.join(output_dir, f"{safe_stem}_master_template.txt")
            master_pptx_path = os.path.join(output_dir, f"{safe_stem}_master_template.pptx")

            with open(sample_txt_path, "w", encoding="utf-8") as f:
                f.write(sample_text)
            with open(master_txt_path, "w", encoding="utf-8") as f:
                f.write(master_text)
            with open(master_pptx_path, "wb") as f:
                f.write(master_pptx_bytes)

            item_result["ok"] = True
            item_result["sample_txt_path"] = sample_txt_path
            item_result["master_txt_path"] = master_txt_path
            item_result["master_pptx_path"] = master_pptx_path
            success_count += 1
        except Exception as e:
            item_result["error"] = str(e)
            failed_count += 1

        results.append(item_result)

    return JsonResponse(
        {
            "success": True,
            "message": "batch extraction finished",
            "processed": len(pptx_files),
            "success_count": success_count,
            "failed_count": failed_count,
            "output_dir": output_dir,
            "results": results,
        },
        status=200,
    )


@require_node("pptx", api=True)
def download_pptx(request, filename: str):
    try:
        abs_path = _resolve_saved_pptx_path(filename)
    except ValueError:
        return HttpResponse("invalid filename", status=400)

    if not os.path.isfile(abs_path):
        return HttpResponse("file not found", status=404)

    ctype, _ = mimetypes.guess_type(abs_path)
    ctype = ctype or "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    with open(abs_path, "rb") as f:
        raw = f.read()
    resp = HttpResponse(raw, content_type=ctype)
    safe_name = os.path.basename(abs_path)
    resp["Content-Disposition"] = f"attachment; filename=\"{safe_name}\"; filename*=UTF-8''{quote(safe_name)}"
    return resp


@require_node("pptx", api=True)
def download_potx(request, filename: str):
    return download_pptx(request, filename)


@require_node("pptx")
def template_admin(request):
    return render(request, "text2pptx/template_admin.html", _build_template_context())


@require_node("pptx")
def import_template(request):
    if request.method != "POST":
        return redirect("pptx_template_admin")

    upload = request.FILES.get("template_file")
    if not upload:
        messages.error(request, "請選擇要匯入的 .pptx 範本檔。")
        return redirect("pptx_template_admin")

    filename = _normalize_template_name(upload.name or "")
    if not filename.lower().endswith(".pptx"):
        messages.error(request, "只支援匯入 .pptx 檔案。")
        return redirect("pptx_template_admin")

    if upload.size and upload.size > (MAX_TEMPLATE_UPLOAD_MB * 1024 * 1024):
        messages.error(request, f"範本檔案過大，請小於 {MAX_TEMPLATE_UPLOAD_MB} MB。")
        return redirect("pptx_template_admin")

    raw = upload.read()
    try:
        audit = _audit_template_bytes(raw)
    except Exception as e:
        messages.error(request, f"範本解析失敗：{e}")
        return redirect("pptx_template_admin")

    if not audit.get("ok"):
        missing = audit.get("missing") or []
        messages.error(request, "匯入失敗，缺少必要版型：" + "、".join(missing))
        return redirect("pptx_template_admin")

    saved_name = _save_template_bytes(filename, raw)
    found = audit.get("found") or {}
    messages.success(
        request,
        "匯入成功："
        f"{saved_name}。"
        f"標題投影片={found.get('cover')} / "
        f"章節標題={found.get('section')} / "
        f"內容頁={found.get('content')} / "
        f"雙欄必較頁={found.get('two_content')}",
    )
    return redirect("pptx_template_admin")


@require_node("pptx", api=True)
def analyze_image_prompts(request):
    if request.method != "POST":
        return JsonResponse({"ok": False, "error": "method not allowed"}, status=405)

    text = (request.POST.get("text") or "").strip()
    if not text:
        return JsonResponse({"ok": False, "error": "請先輸入內容。"}, status=400)
    if len(text) > MAX_INPUT_CHARS:
        return JsonResponse(
            {"ok": False, "error": f"輸入內容長度不可超過 {MAX_INPUT_CHARS} 字。"},
            status=400,
        )

    try:
        analysis_result = _resolve_analysis_for_image_prompts(text)
        slides = analysis_result.get("slides") if isinstance(analysis_result, dict) else []
        normalized_slides: List[Dict[str, Any]] = []
        if isinstance(slides, list):
            for idx, slide in enumerate(slides, start=1):
                if not isinstance(slide, dict):
                    continue
                prompt = str(slide.get("image_prompt") or "").strip()
                normalized_slides.append(
                    {
                        "index": idx,
                        "title": str(slide.get("title") or f"Slide {idx}").strip(),
                        "slide_type": _resolve_slide_type(slide.get("slide_type")),
                        "image_required": _coerce_bool(slide.get("image_required"), default=bool(prompt)),
                        "image_intent": _normalize_image_intent(slide.get("image_intent")),
                        "aspect_ratio": _normalize_aspect_ratio(slide.get("aspect_ratio")),
                        "image_prompt": prompt,
                    }
                )

        return JsonResponse(
            {
                "ok": True,
                "main_title": str(analysis_result.get("main_title") or DEFAULT_MAIN_TITLE),
                "main_subtitle": str(analysis_result.get("main_subtitle") or DEFAULT_MAIN_SUBTITLE),
                "slides": normalized_slides,
                "preview_text": "\n".join(_image_prompt_preview_lines(analysis_result)).strip(),
            }
        )
    except Exception as e:
        logger.warning("Analyze image prompts failed: %s", e)
        return JsonResponse({"ok": False, "error": str(e)}, status=500)


@require_node("pptx", api=True)
def generate_restored_pptx(request):
    if request.method != "POST":
        return JsonResponse(
            {
                "success": False,
                "message": "method not allowed",
                "output_filename": "",
                "download_url": "",
            },
            status=405,
        )

    try:
        payload = json.loads((request.body or b"").decode("utf-8") or "{}")
    except Exception:
        payload = {}

    extracted_text = str(payload.get("extracted_text") or "").strip()
    template_text = str(payload.get("template_text") or "").strip()
    template_pptx_filename = _normalize_template_name(str(payload.get("template_pptx_filename") or "").strip())
    source_filename = _normalize_template_name(str(payload.get("source_filename") or "").strip())

    if not extracted_text:
        return JsonResponse(
            {
                "success": False,
                "message": "請先取得抽取文字結果",
                "output_filename": "",
                "download_url": "",
            },
            status=400,
        )
    if not template_text:
        return JsonResponse(
            {
                "success": False,
                "message": "請先匯入母片模板",
                "output_filename": "",
                "download_url": "",
            },
            status=400,
        )

    source_stem = os.path.splitext(os.path.basename(source_filename))[0] if source_filename else "source"
    safe_stem = _safe_filename(source_stem, default="source")
    output_filename = f"restored_{safe_stem}.pptx"

    template_pptx_path = ""
    if template_pptx_filename:
        try:
            template_pptx_path = _resolve_saved_pptx_path(template_pptx_filename)
        except Exception as e:
            logger.warning("restore template filename invalid: %s", e)
            template_pptx_path = ""

    if not template_pptx_path or not os.path.isfile(template_pptx_path):
        return JsonResponse(
            {
                "success": False,
                "message": "請先使用抽取母片後產生的 PPTX 模板再執行還原",
                "output_filename": "",
                "download_url": "",
            },
            status=400,
        )

    try:
        from pptx import Presentation

        analysis_result = _parse_marked_text_structure(extracted_text)
        if analysis_result is None:
            analysis_result = _fallback_from_raw_text(extracted_text)

        prs = Presentation(template_pptx_path)
        if len(prs.slides) > 0:
            cover_data = {
                "title": str(analysis_result.get("main_title") or DEFAULT_MAIN_TITLE).strip(),
                "subtitle": str(analysis_result.get("main_subtitle") or DEFAULT_MAIN_SUBTITLE).strip(),
            }
            _restore_slide_text_from_analysis(prs.slides[0], cover_data, is_cover=True)

        slides = analysis_result.get("slides") if isinstance(analysis_result, dict) else []
        if not isinstance(slides, list):
            slides = []
        for idx, slide_data in enumerate(slides, start=1):
            if idx >= len(prs.slides):
                break
            if not isinstance(slide_data, dict):
                continue
            _restore_slide_text_from_analysis(prs.slides[idx], slide_data)

        out = io.BytesIO()
        prs.save(out)
        restored_bytes = out.getvalue()
    except Exception as e:
        logger.warning("restore generation failed: %s", e)
        return JsonResponse(
            {
                "success": False,
                "message": f"還原失敗：{e}",
                "output_filename": "",
                "download_url": "",
            },
            status=500,
        )

    script_name = (getattr(request, "script_name", "") or request.META.get("SCRIPT_NAME", "")).rstrip("/")
    saved_name, download_url = _save_generated_pptx_bytes(
        source_filename or "source",
        restored_bytes,
        script_name=script_name,
    )
    return JsonResponse(
        {
            "success": True,
            "message": "還原簡報生成成功",
            "output_filename": saved_name or output_filename,
            "download_url": download_url,
        },
        status=200,
    )


@require_node("pptx", api=True)
def schema_extract(request):
    if request.method != "POST":
        return JsonResponse({"ok": False, "error": "method not allowed"}, status=405)

    upload = request.FILES.get("pptx_file") or request.FILES.get("file")
    if not upload:
        return JsonResponse({"ok": False, "error": "請上傳 .pptx 檔案"}, status=400)
    filename = _normalize_template_name(upload.name or "")
    if not filename.lower().endswith(".pptx"):
        return JsonResponse({"ok": False, "error": "只支援 .pptx 檔案"}, status=400)

    try:
        from webapps.text2pptx.pptx2schema.pipelines.extract_pipeline import run_extract

        result = run_extract(upload.read(), source_file=filename)
        return JsonResponse({"ok": True, "raw": result.model_dump()}, status=200)
    except Exception as e:
        logger.warning("schema_extract failed: %s", e)
        return JsonResponse({"ok": False, "error": str(e)}, status=500)


@require_node("pptx", api=True)
def schema_analyze(request):
    if request.method != "POST":
        return JsonResponse({"ok": False, "error": "method not allowed"}, status=405)

    upload = request.FILES.get("pptx_file") or request.FILES.get("file")
    if not upload:
        return JsonResponse({"ok": False, "error": "請上傳 .pptx 檔案"}, status=400)
    filename = _normalize_template_name(upload.name or "")
    if not filename.lower().endswith(".pptx"):
        return JsonResponse({"ok": False, "error": "只支援 .pptx 檔案"}, status=400)

    try:
        import tempfile

        from webapps.text2pptx.pptx2schema.pipelines.analyze_pipeline import run_analyze

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tf:
            tf.write(upload.read())
            temp_path = tf.name
        bundle = run_analyze(temp_path)
        try:
            os.remove(temp_path)
        except Exception:
            pass
        return JsonResponse({"ok": True, "bundle": bundle.model_dump()}, status=200)
    except Exception as e:
        logger.warning("schema_analyze failed: %s", e)
        return JsonResponse({"ok": False, "error": str(e)}, status=500)


@require_node("pptx", api=True)
def generate_pptx(request):
    if request.method != "POST":
        return JsonResponse({"ok": False, "error": "method not allowed"}, status=405)

    text = (request.POST.get("text") or "").strip()
    user_title = (request.POST.get("title") or "").strip()
    tpl_name = (request.POST.get("template") or "").strip() or DEFAULT_TEMPLATE_NAME

    if not text:
        return JsonResponse({"ok": False, "error": "請輸入文字內容"}, status=400)
    if len(text) > MAX_INPUT_CHARS:
        return JsonResponse({"ok": False, "error": f"文字內容超過上限 {MAX_INPUT_CHARS} 字"}, status=400)

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as e:
        return JsonResponse({"ok": False, "error": f"缺少 python-pptx 套件：{e}"}, status=500)

    try:
        # --- AI 結構化分析 ---
        marked_result = _parse_marked_text_structure(text)
        analysis_result = marked_result if marked_result is not None else _analyze_text_with_llm(text)
        main_title = user_title or analysis_result.get("main_title", DEFAULT_MAIN_TITLE)
        main_subtitle = analysis_result.get("main_subtitle", DEFAULT_MAIN_SUBTITLE)
        llm_slides = analysis_result.get("slides", [])
        
        tpl_path = _safe_select_template(tpl_name)
        prs = Presentation(tpl_path) if tpl_path else Presentation()
        if tpl_path:
            _clear_all_slides(prs)

        cover_layout = _pick_layout_by_display_name(prs, str(analysis_result.get("main_layout_name") or "").strip())
        if cover_layout is None:
            cover_layout = _pick_layout_adaptive(prs, "cover")

        # --- 封面製作 ---
        cover = prs.slides.add_slide(cover_layout)
        cover_title = _find_title_shape(cover)
        if not cover_title:
            tb = cover.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1.5))
            cover_title = tb
        
        cover_title.text_frame.text = main_title # Use LLM-generated main title
        _apply_font_style(cover_title.text_frame.paragraphs[0], TITLE_FONT_PT, bold=True)

        subtitle_lines = [line.strip() for line in str(main_subtitle or "").splitlines() if line.strip()]
        if not subtitle_lines and str(main_subtitle or "").strip():
            subtitle_lines = [str(main_subtitle).strip()]
        cover_text_frames = _find_content_text_frames(cover, exclude_shape=cover_title)
        _fill_text_frames_with_lines(cover_text_frames, subtitle_lines)

        # Use LLM-structured slides
        for i, slide_data in enumerate(llm_slides, start=1):
            slide_title = str(slide_data.get("title") or f"投影片 {i}")
            slide_type = _resolve_slide_type(slide_data.get("slide_type"))
            bullets = slide_data.get("bullets", [])
            if not isinstance(bullets, list):
                bullets = [str(bullets)]
            bullets = [str(b).strip() for b in bullets if str(b).strip()]
            generated_image_path = _generate_image_for_slide(slide_data)

            if slide_type == "section":
                slide = prs.slides.add_slide(_pick_layout_for_slide_data(prs, slide_data, default_type="section"))
                title_shape = _find_title_shape(slide)
                if not title_shape:
                    title_shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.5), Inches(1.6))
                title_shape.text_frame.text = slide_title
                _apply_font_style(title_shape.text_frame.paragraphs[0], TITLE_FONT_PT, bold=True)

                section_tfs = _find_content_text_frames(slide, exclude_shape=title_shape)
                if section_tfs:
                    _fill_text_frames_with_lines(section_tfs, bullets)
                if generated_image_path:
                    _insert_generated_image(
                        slide,
                        generated_image_path,
                        exclude_shape=title_shape,
                        slide_type=slide_type,
                        image_intent=str(slide_data.get("image_intent") or ""),
                    )
                continue

            if slide_type in ("two_content", "comparison"):
                column_chunks = list(_chunk_list(bullets, MAX_BULLETS_PER_SLIDE * 2)) if bullets else [[]]
                for part_idx, chunk in enumerate(column_chunks):
                    slide = prs.slides.add_slide(_pick_layout_for_slide_data(prs, slide_data, default_type=slide_type))
                    shown_title = slide_title if part_idx == 0 else f"{slide_title}（續）"
                    title_shape = _find_title_shape(slide)
                    if not title_shape:
                        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
                    title_shape.text_frame.text = shown_title
                    _apply_font_style(title_shape.text_frame.paragraphs[0], TITLE_FONT_PT, bold=True)

                    content_tfs = _find_content_text_frames(slide, exclude_shape=title_shape)
                    if len(content_tfs) < 2:
                        content_tfs = [
                            slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(4.2), Inches(4.8)).text_frame,
                            slide.shapes.add_textbox(Inches(5.0), Inches(1.6), Inches(4.2), Inches(4.8)).text_frame,
                        ]

                    left_title, left_items, right_title, right_items = _split_two_columns(slide_data, chunk, slide_type)
                    _fill_column(content_tfs[0], left_title, left_items)
                    _fill_column(content_tfs[1], right_title, right_items)
                    if generated_image_path and part_idx == 0:
                        _insert_generated_image(
                            slide,
                            generated_image_path,
                            exclude_shape=title_shape,
                            slide_type=slide_type,
                            image_intent=str(slide_data.get("image_intent") or ""),
                        )
                continue

            bullet_chunks = list(_chunk_list(bullets, MAX_BULLETS_PER_SLIDE)) if bullets else [[]]
            for part_idx, part in enumerate(bullet_chunks):
                slide = prs.slides.add_slide(_pick_layout_for_slide_data(prs, slide_data, default_type="content"))
                shown_title = slide_title if part_idx == 0 else f"{slide_title}（續）"
                title_shape = _find_title_shape(slide)
                if not title_shape:
                    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
                title_shape.text_frame.text = shown_title
                _apply_font_style(title_shape.text_frame.paragraphs[0], TITLE_FONT_PT, bold=True)

                content_tfs = _find_content_text_frames(slide, exclude_shape=title_shape)
                if content_tfs:
                    _fill_text_frames_with_lines(content_tfs, part)
                if generated_image_path and part_idx == 0:
                    _insert_generated_image(
                        slide,
                        generated_image_path,
                        exclude_shape=title_shape,
                        slide_type=slide_type,
                        image_intent=str(slide_data.get("image_intent") or ""),
                    )

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        filename = _build_download_filename(user_title or main_title, fallback=DEFAULT_MAIN_TITLE)
        ascii_fallback = _safe_filename(
            filename.encode("ascii", errors="ignore").decode("ascii"),
            default="slides",
        )
        if not ascii_fallback.lower().endswith(".pptx"):
            ascii_fallback = f"{ascii_fallback}.pptx"
        quoted_filename = quote(filename)
        resp = HttpResponse(buf.getvalue(), content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        resp["Content-Disposition"] = (
            f"attachment; filename=\"{ascii_fallback}\"; "
            f"filename*=UTF-8''{quoted_filename}"
        )
        return resp

    except Exception as e:
        return JsonResponse({"ok": False, "error": str(e)}, status=500)

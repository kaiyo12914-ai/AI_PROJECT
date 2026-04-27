from __future__ import annotations

import io
import base64
import os
import json
import re
import posixpath
import tempfile
import zipfile
import unicodedata
from xml.etree import ElementTree as ET
from urllib.parse import unquote
from unittest.mock import patch

from django.core.files.uploadedfile import SimpleUploadedFile
from django.test import SimpleTestCase, override_settings

from webapps.text2pptx import views
from webapps.text2pptx import image_service


def _build_minimal_pptx_bytes_for_upload() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title:
        slide.shapes.title.text = "Template Source"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_multi_slide_pptx_bytes_for_upload(slide_count: int = 3) -> bytes:
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[0]
    for idx in range(max(1, int(slide_count))):
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = f"Template Source {idx + 1}"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_textbox_pptx_bytes_for_schema() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    if slide.shapes.title:
        slide.shapes.title.text = "Schema Title"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Schema Body"
    textbox = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(4.0), Inches(1.0))
    textbox.text_frame.text = "Standalone textbox"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_chart_pptx_bytes_for_schema() -> bytes:
    from pptx import Presentation
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(4.0), Inches(0.6))
    title.text_frame.text = "Chart Slide"
    chart_data = ChartData()
    chart_data.categories = ["A", "B", "C"]
    chart_data.add_series("Series 1", (1, 2, 3))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.0),
        Inches(1.0),
        Inches(6.0),
        Inches(4.0),
        chart_data,
    )
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _tiny_png_bytes_variant_a() -> bytes:
    return base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+yXioAAAAASUVORK5CYII="
    )


def _tiny_png_bytes_variant_b() -> bytes:
    return base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAAAAAA6fptVAAAACklEQVR42mNkAAIAAAoAAdYk4XkAAAAASUVORK5CYII="
    )


class Text2PptxHelperTests(SimpleTestCase):
    def test_clean_extracted_line_removes_page_number_only(self):
        self.assertEqual(views._clean_extracted_line("Page 01"), "")
        self.assertEqual(views._clean_extracted_line("第1頁"), "")

    def test_clean_extracted_line_removes_page_prefix_but_keeps_content(self):
        self.assertEqual(views._clean_extracted_line("Page 02 Design Philosophy"), "Design Philosophy")
        self.assertEqual(views._clean_extracted_line("第3頁 專案目標"), "專案目標")

    def test_clean_extracted_line_keeps_normal_numeric_content(self):
        self.assertEqual(views._clean_extracted_line("M3-M4：功能開發與整合"), "M3-M4：功能開發與整合")

    def test_remove_footer_numeric_page_lines_filters_trailing_page_number(self):
        lines = ["市場缺口", "20XX 年", "募資簡報標題", "3"]
        out = views._remove_footer_numeric_page_lines(lines)
        self.assertEqual(out, ["市場缺口", "20XX 年", "募資簡報標題"])

    def test_remove_footer_numeric_page_lines_keeps_numeric_data(self):
        lines = ["收入", "10", "20", "30"]
        out = views._remove_footer_numeric_page_lines(lines)
        self.assertEqual(out, lines)

    def test_extract_json_object_from_fenced_text(self):
        raw = """```json
        {"main_title":"A","main_subtitle":"B","slides":[]}
        ```"""
        extracted = views._extract_json_object(raw)
        self.assertEqual(extracted, '{"main_title":"A","main_subtitle":"B","slides":[]}')

    def test_normalize_analysis_result_invalid_slide_type_falls_back_to_content(self):
        data = {
            "main_title": "主標",
            "main_subtitle": "副標",
            "slides": [
                {
                    "title": "頁1",
                    "slide_type": "unknown_type",
                    "bullets": ["a", "b"],
                }
            ],
        }
        out = views._normalize_analysis_result(data, source_text="頁1\na\nb")
        self.assertEqual(out["slides"][0]["slide_type"], "content")

    def test_normalize_analysis_result_sets_image_fields_defaults(self):
        data = {
            "main_title": "主標",
            "main_subtitle": "副標",
            "slides": [
                {
                    "title": "頁1",
                    "slide_type": "content",
                    "bullets": ["a"],
                }
            ],
        }
        out = views._normalize_analysis_result(data, source_text="頁1\na")
        slide = out["slides"][0]
        self.assertFalse(slide["image_required"])
        self.assertEqual(slide["image_prompt"], "")
        self.assertEqual(slide["image_intent"], "concept")
        self.assertEqual(slide["aspect_ratio"], "16:9")

    def test_normalize_analysis_result_invalid_image_intent_defaults_to_concept(self):
        data = {
            "main_title": "T",
            "main_subtitle": "S",
            "slides": [
                {
                    "title": "A",
                    "slide_type": "content",
                    "bullets": ["x"],
                    "image_required": True,
                    "image_prompt": "abstract workflow diagram",
                    "image_intent": "poster",
                }
            ],
        }
        out = views._normalize_analysis_result(data, source_text="A\nx")
        slide = out["slides"][0]
        self.assertEqual(slide["image_intent"], "concept")

    def test_normalize_analysis_result_reads_image_fields(self):
        data = {
            "main_title": "主標",
            "main_subtitle": "副標",
            "slides": [
                {
                    "title": "頁1",
                    "slide_type": "content",
                    "bullets": ["a"],
                    "image_required": "true",
                    "image_prompt": "modern office team discussing timeline",
                    "image_intent": "Hero",
                    "aspect_ratio": "4:3",
                }
            ],
        }
        out = views._normalize_analysis_result(data, source_text="頁1\na")
        slide = out["slides"][0]
        self.assertTrue(slide["image_required"])
        self.assertEqual(slide["image_prompt"], "modern office team discussing timeline")
        self.assertEqual(slide["image_intent"], "hero")
        self.assertEqual(slide["aspect_ratio"], "4:3")

    def test_generate_image_for_slide_mock_mode(self):
        slide_data = {
            "image_required": True,
            "image_prompt": "simple blue abstract background",
            "aspect_ratio": "16:9",
        }
        with tempfile.TemporaryDirectory(dir=os.getcwd(), ignore_cleanup_errors=True) as td:
            with patch.object(views, "GENERATED_IMAGE_DIR", td), patch.object(views, "TEXT2PPTX_IMAGE_MODE", "mock"):
                path = views._generate_image_for_slide(slide_data)
                self.assertIsNotNone(path)
                assert path is not None
                self.assertTrue(os.path.exists(path))

    def test_generate_image_google_accepts_ok_image_path(self):
        with tempfile.TemporaryDirectory(dir=os.getcwd(), ignore_cleanup_errors=True) as td:
            image_path = os.path.join(td, "provider_out.png")
            with open(image_path, "wb") as f:
                f.write(base64.b64decode("iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/w8AAusB9Y0A3L8AAAAASUVORK5CYII="))
            with patch("webapps.text2pptx.image_service._call_google_provider", return_value={"ok": True, "image_path": image_path}):
                result = image_service.generate_image(prompt="x", mode="google", output_dir=td)
            self.assertTrue(result["ok"])
            self.assertEqual(result["local_path"], image_path)

    def test_generate_image_google_uses_output_path_when_ok_true(self):
        with tempfile.TemporaryDirectory(dir=os.getcwd(), ignore_cleanup_errors=True) as td:
            def _provider(**kwargs):
                out = str(kwargs["output_path"])
                with open(out, "wb") as f:
                    f.write(base64.b64decode("iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/w8AAusB9Y0A3L8AAAAASUVORK5CYII="))
                return {"ok": True}

            with patch("webapps.text2pptx.image_service._call_google_provider", side_effect=_provider):
                result = image_service.generate_image(prompt="x", mode="google", output_dir=td, aspect_ratio="3:2")
            self.assertTrue(result["ok"])
            self.assertTrue(os.path.isfile(result["local_path"]))

    def test_generate_image_google_raises_when_ok_false(self):
        with tempfile.TemporaryDirectory(dir=os.getcwd(), ignore_cleanup_errors=True) as td:
            with patch(
                "webapps.text2pptx.image_service._call_google_provider",
                return_value={"ok": False, "error": "invalid API key"},
            ):
                with self.assertRaises(image_service.ImageGenError) as cm:
                    image_service.generate_image(prompt="x", mode="google", output_dir=td)
        self.assertEqual(cm.exception.code, "IMG_E_PROVIDER")

    def test_sample_extraction_preserves_schema_metadata(self):
        raw = _build_textbox_pptx_bytes_for_schema()
        result = views._extract_sample_text_from_pptx_bytes(raw)
        self.assertNotIn("@slide_schema ", result["sample_text"])
        self.assertIsInstance(result["sample_schema"], dict)
        self.assertEqual(result["sample_schema"]["meta"]["slide_count"], 1)
        self.assertTrue(
            any(
                shape.get("text_frame")
                for shape in result["sample_schema"]["slides"][0]["shapes"]
                if isinstance(shape, dict)
            )
        )

    def test_apply_font_style_writes_valid_east_asian_font_xml(self):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(4.0), Inches(1.0))
        paragraph = textbox.text_frame.paragraphs[0]
        paragraph.text = "測試"
        views._apply_font_style(paragraph, 18, bold=True)

        buf = io.BytesIO()
        prs.save(buf)

        with zipfile.ZipFile(io.BytesIO(buf.getvalue()), "r") as zf:
            xml = zf.read("ppt/slides/slide1.xml").decode("utf-8", errors="strict")
        self.assertNotIn('ea="', xml)
        self.assertIn(f'<a:ea typeface="{views.FONT_NAME_EAST_ASIAN}"', xml)

    def test_clear_all_text_on_shape_preserves_textbox_geometry_and_styles(self):
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(1.25), Inches(2.0), Inches(3.5), Inches(1.1))
        frame = textbox.text_frame
        paragraph = frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = "Preserve me"
        run.font.name = "Arial"
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x12, 0x34, 0x56)

        views._clear_all_text_on_shape(textbox)

        buf = io.BytesIO()
        prs.save(buf)
        xml = zipfile.ZipFile(io.BytesIO(buf.getvalue()), "r").read("ppt/slides/slide1.xml").decode("utf-8", errors="strict")
        self.assertIn("<a:xfrm>", xml)
        self.assertIn("<a:off", xml)
        self.assertIn("<a:ext", xml)
        self.assertIn('typeface="Arial"', xml)
        self.assertIn('val="123456"', xml)

    def test_master_template_extraction_embeds_template_schema_metadata(self):
        raw = _build_minimal_pptx_bytes_for_upload()
        template_text = views._extract_master_template_text_from_pptx_bytes(raw)
        self.assertIn("@template_schema ", template_text)
        self.assertIn("版型1", template_text)
        self.assertNotIn("版型2", template_text)

    def test_convert_pptx_to_potx_prunes_unused_layouts(self):
        raw = _build_minimal_pptx_bytes_for_upload()
        potx = views._convert_pptx_bytes_to_potx_bytes(raw)

        with zipfile.ZipFile(io.BytesIO(potx), "r") as zf:
            layout_files = sorted(
                name
                for name in zf.namelist()
                if re.match(r"^ppt/slideLayouts/slideLayout\d+\.xml$", name)
            )
            master_files = sorted(
                name
                for name in zf.namelist()
                if re.match(r"^ppt/slideMasters/slideMaster\d+\.xml$", name)
            )
            self.assertEqual(layout_files, ["ppt/slideLayouts/slideLayout1.xml"])
            self.assertEqual(len(master_files), 1)

            master_root = ET.fromstring(zf.read(master_files[0]))
            ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
            layout_ids = master_root.findall(".//p:sldLayoutId", ns)
            self.assertEqual(len(layout_ids), 1)

            content_types_root = ET.fromstring(zf.read("[Content_Types].xml"))
            ct_ns = {"ct": "http://schemas.openxmlformats.org/package/2006/content-types"}
            presentation_override = content_types_root.find(
                'ct:Override[@PartName="/ppt/presentation.xml"]',
                ct_ns,
            )
            self.assertIsNotNone(presentation_override)
            self.assertEqual(
                presentation_override.get("ContentType"),
                views.PPTX_MAIN_CONTENT_TYPE,
            )

    def test_convert_pptx_to_potx_renames_layouts_sequentially(self):
        from pptx import Presentation

        prs = Presentation()
        slide_a = prs.slides.add_slide(prs.slide_layouts[0])
        slide_b = prs.slides.add_slide(prs.slide_layouts[1])
        if slide_a.shapes.title:
            slide_a.shapes.title.text = "A"
        if slide_b.shapes.title:
            slide_b.shapes.title.text = "B"

        buf = io.BytesIO()
        prs.save(buf)
        potx = views._convert_pptx_bytes_to_potx_bytes(buf.getvalue())

        p_ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        with zipfile.ZipFile(io.BytesIO(potx), "r") as zf:
            layout_files = sorted(
                name
                for name in zf.namelist()
                if re.match(r"^ppt/slideLayouts/slideLayout\d+\.xml$", name)
            )
            layout_names = []
            for name in layout_files:
                root = ET.fromstring(zf.read(name))
                c_sld = root.find("p:cSld", p_ns)
                layout_names.append(str(c_sld.get("name") if c_sld is not None else ""))

        self.assertEqual(layout_names[:2], ["版型1", "版型2"])

    def test_convert_pptx_to_potx_preserves_picture_placeholder_picture(self):
        from pptx import Presentation

        prs = Presentation()
        layout = prs.slide_layouts[8]
        slide = prs.slides.add_slide(layout)

        picture_placeholder = next(
            (
                ph
                for ph in slide.placeholders
                if int(getattr(getattr(ph, "placeholder_format", None), "type", -1)) == 18
            ),
            None,
        )
        self.assertIsNotNone(picture_placeholder)
        assert picture_placeholder is not None

        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp.write(_tiny_png_bytes_variant_a())
            tmp_path = tmp.name
        try:
            picture_placeholder.insert_picture(tmp_path)
            buf = io.BytesIO()
            prs.save(buf)
            potx = views._convert_pptx_bytes_to_potx_bytes(buf.getvalue())
        finally:
            try:
                os.remove(tmp_path)
            except Exception:
                pass

        with zipfile.ZipFile(io.BytesIO(potx), "r") as zf:
            layout_files = sorted(
                name
                for name in zf.namelist()
                if re.match(r"^ppt/slideLayouts/slideLayout\d+\.xml$", name)
            )
            self.assertEqual(len(layout_files), 1)
            layout_root = ET.fromstring(zf.read(layout_files[0]))
            p_ns = {
                "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
                "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            }
            pic_nodes = layout_root.findall(".//p:pic", p_ns)
            self.assertEqual(len(pic_nodes), 1)
            blip = layout_root.find(".//a:blip", p_ns)
            self.assertIsNotNone(blip)
            assert blip is not None
            self.assertTrue(blip.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"))

    def test_convert_pptx_to_potx_copies_chart_related_parts_recursively(self):
        raw = _build_chart_pptx_bytes_for_schema()
        potx = views._convert_pptx_bytes_to_potx_bytes(raw)

        rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
        with zipfile.ZipFile(io.BytesIO(potx), "r") as zf:
            names = set(zf.namelist())
            chart_parts = sorted(name for name in names if re.match(r"^ppt/charts/chart\d+\.xml$", name))
            self.assertTrue(chart_parts)
            chart_rels_parts = [name for name in names if re.match(r"^ppt/charts/_rels/chart\d+\.xml\.rels$", name)]
            self.assertTrue(chart_rels_parts)
            embedded_parts = [name for name in names if re.match(r"^ppt/embeddings/.*\.(xlsx|bin)$", name)]
            self.assertTrue(embedded_parts)

            for rels_name in sorted(name for name in names if name.endswith(".rels")):
                base_part = "" if rels_name == "_rels/.rels" else rels_name.replace("/_rels/", "/").removesuffix(".rels")
                rels_root = ET.fromstring(zf.read(rels_name))
                for rel in rels_root.findall(f"{{{rel_ns}}}Relationship"):
                    if str(rel.get("TargetMode") or "").lower() == "external":
                        continue
                    target = str(rel.get("Target") or "").strip()
                    if not target:
                        continue
                    if target.startswith("/"):
                        resolved = target.lstrip("/")
                    else:
                        resolved = posixpath.normpath(posixpath.join(posixpath.dirname(base_part), target)).lstrip("./")
                    self.assertIn(resolved, names, msg=f"missing target {resolved} from {rels_name}")

    def test_insert_generated_image_places_picture_in_available_area(self):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        blank_layout = next((ly for ly in prs.slide_layouts if len(getattr(ly, "placeholders", [])) == 0), prs.slide_layouts[-1])
        slide = prs.slides.add_slide(blank_layout)
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        title_shape.text_frame.text = "Title"

        with tempfile.TemporaryDirectory(dir=os.getcwd(), ignore_cleanup_errors=True) as td:
            image_path = os.path.join(td, "probe.png")
            png = base64.b64decode(
                "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8"
                "/w8AAusB9Y0A3L8AAAAASUVORK5CYII="
            )
            with open(image_path, "wb") as f:
                f.write(png)

            inserted = views._insert_generated_image(
                slide,
                image_path,
                exclude_shape=title_shape,
                slide_type="content",
                image_intent="hero",
            )

        self.assertTrue(inserted)
        self.assertTrue(any(hasattr(sh, "image") for sh in slide.shapes))

    def test_insert_generated_image_skips_when_text_almost_fullslide(self):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        blank_layout = next((ly for ly in prs.slide_layouts if len(getattr(ly, "placeholders", [])) == 0), prs.slide_layouts[-1])
        slide = prs.slides.add_slide(blank_layout)
        title_shape = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(9.3), Inches(0.8))
        title_shape.text_frame.text = "Title"
        body = slide.shapes.add_textbox(Inches(0.0), Inches(0.0), Inches(10.0), Inches(7.5))
        body.text_frame.text = "This area is intentionally occupied by text."

        with tempfile.TemporaryDirectory(dir=os.getcwd(), ignore_cleanup_errors=True) as td:
            image_path = os.path.join(td, "probe.png")
            png = base64.b64decode(
                "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8"
                "/w8AAusB9Y0A3L8AAAAASUVORK5CYII="
            )
            with open(image_path, "wb") as f:
                f.write(png)

            with patch.object(views, "_find_largest_image_frame", return_value=None):
                inserted = views._insert_generated_image(
                    slide,
                    image_path,
                    exclude_shape=title_shape,
                    slide_type="content",
                    image_intent="data",
                )

        self.assertFalse(inserted)

    @patch("webapps.text2pptx.views.get_chat_model")
    def test_analyze_text_with_llm_parses_json_with_extra_explanation(self, mock_get_chat_model):
        class DummyLLM:
            def invoke(self, _prompt):
                return (
                    "以下是分析結果：\n"
                    "```json\n"
                    '{"main_title":"T","main_subtitle":"S","slides":[{"title":"A","slide_type":"content","bullets":["x"]}]}\n'
                    "```\n"
                    "結束。"
                )

        mock_get_chat_model.return_value = DummyLLM()
        out = views._analyze_text_with_llm("原文")
        self.assertEqual(out["main_title"], "T")
        self.assertEqual(out["slides"][0]["title"], "A")
        self.assertEqual(out["slides"][0]["slide_type"], "content")

    @patch("webapps.text2pptx.views.get_chat_model")
    def test_analyze_text_with_llm_falls_back_when_non_json(self, mock_get_chat_model):
        class DummyLLM:
            def invoke(self, _prompt):
                return "這是一段沒有 JSON 的回覆"

        mock_get_chat_model.return_value = DummyLLM()
        out = views._analyze_text_with_llm("頁1\n點1\n===\n頁2\n點2")
        self.assertEqual(out["main_title"], views.DEFAULT_MAIN_TITLE)
        self.assertEqual(len(out["slides"]), 2)
        self.assertEqual(out["slides"][0]["title"], "頁1")
        self.assertEqual(out["slides"][0]["slide_type"], "content")

    @patch("webapps.text2pptx.views.get_chat_model")
    def test_analyze_text_with_llm_falls_back_when_incomplete_json(self, mock_get_chat_model):
        class DummyLLM:
            def invoke(self, _prompt):
                return '{"main_title":"X","main_subtitle":"Y","slides":[{"title":"A"'

        mock_get_chat_model.return_value = DummyLLM()
        out = views._analyze_text_with_llm("單頁標題\n重點")
        self.assertEqual(out["main_title"], views.DEFAULT_MAIN_TITLE)
        self.assertEqual(len(out["slides"]), 1)
        self.assertEqual(out["slides"][0]["title"], "單頁標題")

    def test_parse_marked_text_structure_chinese_markers(self):
        text = """[標題投影片] 智慧文件流程優化專案提案
提案單位：資訊處
===
[章節標題] 章節一：現況與目標
本章重點摘要
===
[內容頁] 提案背景
痛點一
痛點二
===
[雙欄必較頁] 成本與時程
左欄：自建成本高、時程長
右欄：採購成本中、時程短
===
[雙欄必較頁] 風險比較
左欄：自建需技術團隊
右欄：採購受供應商排程影響"""
        out = views._parse_marked_text_structure(text)
        self.assertIsNotNone(out)
        assert out is not None
        self.assertEqual(out["main_title"], "智慧文件流程優化專案提案")
        self.assertEqual(out["main_subtitle"], "提案單位：資訊處")
        self.assertEqual(len(out["slides"]), 4)
        self.assertEqual(out["slides"][0]["slide_type"], "section")
        self.assertEqual(out["slides"][1]["slide_type"], "content")
        self.assertEqual(out["slides"][2]["slide_type"], "two_content")
        self.assertEqual(out["slides"][3]["slide_type"], "comparison")

    def test_parse_marked_text_structure_ignores_layout_meta_line(self):
        text = """[內容頁] 提案背景
版面名稱：標題及內容
痛點一
痛點二"""
        out = views._parse_marked_text_structure(text)
        self.assertIsNotNone(out)
        assert out is not None
        self.assertEqual(len(out["slides"]), 1)
        self.assertEqual(out["slides"][0]["title"], "提案背景")
        self.assertEqual(out["slides"][0]["bullets"], ["痛點一", "痛點二"])

    def test_extract_sample_text_includes_layout_name(self):
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = "封面標題"
        buf = io.BytesIO()
        prs.save(buf)
        out = views._extract_sample_text_from_pptx_bytes(buf.getvalue())
        sample_text = str(out.get("sample_text") or "")
        self.assertIn("[版型1]封面標題", sample_text)
        self.assertNotIn("@slide_schema ", sample_text)

    def test_extract_sample_text_numbers_slides_sequentially(self):
        raw = _build_multi_slide_pptx_bytes_for_upload(slide_count=2)
        out = views._extract_sample_text_from_pptx_bytes(raw)
        sample_text = str(out.get("sample_text") or "")
        self.assertIn("[版型1]", sample_text)
        self.assertIn("[版型2]", sample_text)
        self.assertEqual(sample_text.count("[版型1]"), 1)
        self.assertEqual(sample_text.count("[版型2]"), 1)

    def test_parse_marked_text_structure_accepts_custom_layout_marker(self):
        text = """[感謝您]Mirjam Nilsson
206-555-0146
mirjam@contoso.com"""
        out = views._parse_marked_text_structure(text)
        self.assertIsNotNone(out)
        assert out is not None
        self.assertEqual(len(out["slides"]), 1)
        self.assertEqual(out["slides"][0]["title"], "Mirjam Nilsson")
        self.assertEqual(out["slides"][0]["layout_name"], "")
        self.assertEqual(out["slides"][0]["bullets"], ["206-555-0146", "mirjam@contoso.com"])

    def test_clear_all_slides_removes_existing_template_slides(self):
        from pptx import Presentation

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
        self.assertEqual(len(prs.slides), 2)
        views._clear_all_slides(prs)
        self.assertEqual(len(prs.slides), 0)

    def test_convert_pptx_to_potx_keeps_total_slide_count(self):
        raw = _build_multi_slide_pptx_bytes_for_upload(slide_count=3)
        potx_raw = views._convert_pptx_bytes_to_potx_bytes(raw)

        ns = {"ct": "http://schemas.openxmlformats.org/package/2006/content-types"}
        with zipfile.ZipFile(io.BytesIO(potx_raw), "r") as zf:
            slide_parts = [name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
            self.assertEqual(len(slide_parts), 3)

            slide_root = ET.fromstring(zf.read(slide_parts[0]))
            a_ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            text_nodes = [str(n.text or "").strip() for n in slide_root.findall(".//a:t", a_ns)]
            self.assertFalse(any(text_nodes))

            root = ET.fromstring(zf.read("[Content_Types].xml"))
            override = None
            for node in root.findall("ct:Override", ns):
                if node.get("PartName") == "/ppt/presentation.xml":
                    override = node
                    break

            self.assertIsNotNone(override)
            assert override is not None
            self.assertEqual(
                override.get("ContentType"),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
            )

    def test_convert_pptx_to_potx_keeps_representative_slides_for_each_layout(self):
        from pptx import Presentation

        prs = Presentation()
        layout_a = prs.slide_layouts[0]
        layout_b = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else layout_a

        s1 = prs.slides.add_slide(layout_a)
        if s1.shapes.title:
            s1.shapes.title.text = "A1"
        s2 = prs.slides.add_slide(layout_a)
        if s2.shapes.title:
            s2.shapes.title.text = "A2"

        s3 = prs.slides.add_slide(layout_b)
        if s3.shapes.title:
            s3.shapes.title.text = "B1"
        s4 = prs.slides.add_slide(layout_b)
        if s4.shapes.title:
            s4.shapes.title.text = "B2"

        src = io.BytesIO()
        prs.save(src)
        potx_raw = views._convert_pptx_bytes_to_potx_bytes(src.getvalue())

        expected_count = 4
        with zipfile.ZipFile(io.BytesIO(potx_raw), "r") as zf:
            slide_parts = [name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
            self.assertEqual(len(slide_parts), expected_count)

            a_ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            for slide_part in slide_parts:
                slide_root = ET.fromstring(zf.read(slide_part))
                text_nodes = [str(n.text or "").strip() for n in slide_root.findall(".//a:t", a_ns)]
                self.assertFalse(any(text_nodes))

    def test_convert_pptx_to_potx_same_layout_keeps_last_picture(self):
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

        prs = Presentation()
        layout = next((l for l in prs.slide_layouts if "blank" in (getattr(l, "name", "") or "").lower()), prs.slide_layouts[0])

        slide1 = prs.slides.add_slide(layout)
        slide1.shapes.add_textbox(Inches(1.6), Inches(0.5), Inches(2.6), Inches(0.6)).text_frame.text = "OLD_TEXT_1"
        shp1 = slide1.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.6), Inches(2.1), Inches(2.0), Inches(0.8))
        shp1.text_frame.text = "OLD_SHAPE_TEXT"
        img_a = _tiny_png_bytes_variant_a()
        slide1.shapes.add_picture(io.BytesIO(img_a), Inches(0.8), Inches(1.2), Inches(0.8), Inches(0.8))

        slide2 = prs.slides.add_slide(layout)
        slide2.shapes.add_textbox(Inches(1.6), Inches(0.5), Inches(2.6), Inches(0.6)).text_frame.text = "LATEST_TEXT_2"
        shp2 = slide2.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.6), Inches(2.1), Inches(2.0), Inches(0.8))
        shp2.text_frame.text = "LATEST_SHAPE_TEXT"
        img_b = _tiny_png_bytes_variant_b()
        slide2.shapes.add_picture(io.BytesIO(img_b), Inches(0.8), Inches(1.2), Inches(0.8), Inches(0.8))

        src = io.BytesIO()
        prs.save(src)
        potx_raw = views._convert_pptx_bytes_to_potx_bytes(src.getvalue())

        rel_ns = {"r": "http://schemas.openxmlformats.org/package/2006/relationships"}
        p_ns = {
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        }
        with zipfile.ZipFile(io.BytesIO(potx_raw), "r") as zf:
            slide_parts = [name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
            self.assertEqual(len(slide_parts), 2)

            layout_parts_for_slides: List[str] = []
            for slide_xml_name in sorted(slide_parts):
                slide_idx = os.path.splitext(os.path.basename(slide_xml_name))[0].replace("slide", "")
                slide_rels_path = f"ppt/slides/_rels/slide{slide_idx}.xml.rels"
                slide_rels_root = ET.fromstring(zf.read(slide_rels_path))
                layout_target = None
                for rel in slide_rels_root.findall("r:Relationship", rel_ns):
                    if rel.get("Type") == "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout":
                        layout_target = rel.get("Target")
                        break
                self.assertIsNotNone(layout_target)
                assert layout_target is not None
                layout_part_name = ("ppt/slides/" + layout_target).replace("\\", "/")
                if "/../" in layout_part_name:
                    parts = []
                    for part in layout_part_name.split("/"):
                        if part == "..":
                            if parts:
                                parts.pop()
                        elif part and part != ".":
                            parts.append(part)
                    layout_part_name = "/".join(parts)
                layout_parts_for_slides.append(layout_part_name)

            # Same layout should be split into two variants.
            self.assertEqual(len(set(layout_parts_for_slides)), 2)

            extracted_images: List[bytes] = []
            layout_names: List[str] = []
            for layout_part_name in layout_parts_for_slides:
                layout_root = ET.fromstring(zf.read(layout_part_name))
                c_sld = layout_root.find("p:cSld", p_ns)
                layout_names.append(str(c_sld.get("name") if c_sld is not None else ""))

                pic_nodes = layout_root.findall(".//p:pic", p_ns)
                self.assertEqual(len(pic_nodes), 1)

                sp_nodes = layout_root.findall(".//p:sp", p_ns)
                non_placeholder_sp_nodes = [n for n in sp_nodes if n.find("p:nvSpPr/p:nvPr/p:ph", p_ns) is None]
                self.assertTrue(len(non_placeholder_sp_nodes) >= 2)

                for node in non_placeholder_sp_nodes:
                    node_text_values = [str(t.text or "").strip() for t in node.findall(".//a:t", p_ns)]
                    self.assertFalse(any(node_text_values))

                layout_xml_text = zf.read(layout_part_name).decode("utf-8", errors="ignore")
                self.assertNotIn("OLD_TEXT_1", layout_xml_text)
                self.assertNotIn("LATEST_TEXT_2", layout_xml_text)
                self.assertNotIn("OLD_SHAPE_TEXT", layout_xml_text)
                self.assertNotIn("LATEST_SHAPE_TEXT", layout_xml_text)

                blip = layout_root.find(".//a:blip", p_ns)
                self.assertIsNotNone(blip)
                assert blip is not None
                image_rel_id = blip.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
                self.assertTrue(image_rel_id)

                layout_filename = os.path.basename(layout_part_name)
                rels_path = f"ppt/slideLayouts/_rels/{layout_filename}.rels"
                rels_root = ET.fromstring(zf.read(rels_path))
                image_target = None
                for rel in rels_root.findall("r:Relationship", rel_ns):
                    if rel.get("Id") == image_rel_id:
                        image_target = rel.get("Target")
                        break
                self.assertIsNotNone(image_target)
                assert image_target is not None

                if image_target.startswith("../"):
                    image_path = "ppt/" + image_target[3:]
                else:
                    image_path = "ppt/slideLayouts/" + image_target
                extracted_images.append(zf.read(image_path))

            # Verify both source images are preserved across layout variants.
            self.assertEqual(set(extracted_images), {img_a, img_b})
            self.assertTrue(any(name == "版型1" for name in layout_names))
            self.assertTrue(any(name == "版型2" for name in layout_names))

    def test_safe_select_template_supports_chinese_filename(self):
        with tempfile.TemporaryDirectory(dir=os.getcwd(), ignore_cleanup_errors=True) as td:
            fname_nfc = "中心測試.pptx"
            with open(os.path.join(td, fname_nfc), "wb") as f:
                f.write(b"x")
            with patch.object(views, "PPTX_TEMPLATE_DIR", td):
                selected = views._safe_select_template(unicodedata.normalize("NFD", fname_nfc))
            self.assertEqual(selected, os.path.join(td, fname_nfc))

    def test_save_template_bytes_preserves_chinese_filename(self):
        with tempfile.TemporaryDirectory(dir=os.getcwd(), ignore_cleanup_errors=True) as td:
            with patch.object(views, "PPTX_TEMPLATE_DIR", td):
                saved1 = views._save_template_bytes("中心測試.pptx", b"a")
                saved2 = views._save_template_bytes("中心測試.pptx", b"b")
            self.assertEqual(saved1, "中心測試.pptx")
            self.assertEqual(saved2, "中心測試_1.pptx")
            self.assertTrue(os.path.exists(os.path.join(td, saved1)))
            self.assertTrue(os.path.exists(os.path.join(td, saved2)))

    def test_build_download_filename_avoids_double_extension(self):
        self.assertEqual(views._build_download_filename("測試主題"), "測試主題.pptx")
        self.assertEqual(views._build_download_filename("測試主題.pptx"), "測試主題.pptx")


@override_settings(PORTAL_ACL_ENABLED=False)
class Text2PptxViewTests(SimpleTestCase):
    databases = {"default"}

    def test_analyze_image_prompts_method_not_allowed(self):
        resp = self.client.get("/text2pptx/analyze-image-prompts/")
        self.assertEqual(resp.status_code, 405)

    @patch("webapps.text2pptx.views._save_generated_pptx_bytes")
    def test_extract_template_returns_pptx_download_info(self, mock_save):
        mock_save.return_value = (
            "default_master_20260325_000000.pptx",
            "/media/text2pptx/pptx/default_master_20260325_000000.pptx",
        )
        upload = SimpleUploadedFile(
            "default.pptx",
            _build_minimal_pptx_bytes_for_upload(),
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        resp = self.client.post("/text2pptx/extract_template/", {"pptx_file": upload})
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertTrue(data.get("success"))
        self.assertIn("template_text", data)
        self.assertEqual(data.get("output_filename"), "default_master_20260325_000000.pptx")
        self.assertIn("/text2pptx/download-pptx/default_master_20260325_000000.pptx", data.get("download_url", ""))
        self.assertIn("/media/text2pptx/pptx/default_master_20260325_000000.pptx", data.get("media_download_url", ""))

    def test_save_generated_pptx_bytes_applies_script_name_prefix(self):
        with tempfile.TemporaryDirectory(dir=os.getcwd(), ignore_cleanup_errors=True) as td:
            with patch.object(views.settings, "MEDIA_ROOT", td), patch.object(views.settings, "MEDIA_URL", "/media/"):
                filename, download_url = views._save_generated_pptx_bytes(
                    "default.pptx",
                    b"fake-pptx",
                    script_name="/djangoai",
                )
        self.assertTrue(filename.endswith(".pptx"))
        self.assertTrue(download_url.startswith("/djangoai/media/text2pptx/pptx/"))

    def test_download_pptx_returns_file(self):
        with tempfile.TemporaryDirectory(dir=os.getcwd(), ignore_cleanup_errors=True) as td:
            base = os.path.join(td, "text2pptx", "pptx")
            os.makedirs(base, exist_ok=True)
            fn = "sample_master.pptx"
            with open(os.path.join(base, fn), "wb") as f:
                f.write(b"pptx-bytes")
            with patch.object(views.settings, "MEDIA_ROOT", td):
                resp = self.client.get(f"/text2pptx/download-pptx/{fn}")
        self.assertEqual(resp.status_code, 200)
        self.assertEqual(resp.content, b"pptx-bytes")
        self.assertIn("attachment", resp.get("Content-Disposition", ""))

    def test_sample_extractor_page_shows_restore_buttons(self):
        resp = self.client.get("/text2pptx/sample-extractor/")
        self.assertEqual(resp.status_code, 200)
        self.assertContains(resp, "抽取母片模板")
        self.assertContains(resp, "抽取範例及母片")
        self.assertContains(resp, "匯入母片模板")
        self.assertContains(resp, "生成還原簡報")
        self.assertContains(resp, "母片模板內容")

    def test_generate_restored_pptx_requires_post(self):
        resp = self.client.get("/text2pptx/generate_restored_pptx/")
        self.assertEqual(resp.status_code, 405)
        self.assertFalse(resp.json().get("success", True))

    def test_generate_restored_pptx_requires_extracted_text(self):
        resp = self.client.post(
            "/text2pptx/generate_restored_pptx/",
            data=json.dumps({"extracted_text": "", "template_text": "x"}),
            content_type="application/json",
        )
        self.assertEqual(resp.status_code, 400)
        self.assertEqual(resp.json().get("message"), "請先取得抽取文字結果")

    def test_generate_restored_pptx_success(self):
        self.skipTest("replaced by restored PPTX integration test")
        resp = self.client.post(
            "/text2pptx/generate_restored_pptx/",
            data=json.dumps(
                {
                    "extracted_text": "[內容頁] 測試",
                    "template_text": "{\"layouts\":[]}",
                    "source_filename": "default.pptx",
                }
            ),
            content_type="application/json",
        )
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertTrue(data.get("success"))
        self.assertEqual(data.get("output_filename"), "restored_default.pptx")
        self.assertIn("/download/restored_default.pptx", data.get("download_url", ""))

    def test_generate_restored_pptx_with_saved_template_pptx(self):
        from pptx import Presentation

        with tempfile.TemporaryDirectory(dir=os.getcwd(), ignore_cleanup_errors=True) as td:
            media_pptx_dir = os.path.join(td, "text2pptx", "pptx")
            os.makedirs(media_pptx_dir, exist_ok=True)
            template_filename = "template_source.pptx"
            template_path = os.path.join(media_pptx_dir, template_filename)

            prs = Presentation()
            slide0 = prs.slides.add_slide(prs.slide_layouts[0])
            if slide0.shapes.title:
                slide0.shapes.title.text = "placeholder"
            if len(slide0.placeholders) > 1:
                slide0.placeholders[1].text = "placeholder"
            slide1 = prs.slides.add_slide(prs.slide_layouts[1])
            if slide1.shapes.title:
                slide1.shapes.title.text = "placeholder"
            if len(slide1.placeholders) > 1:
                slide1.placeholders[1].text = "placeholder"
            prs.save(template_path)

            payload = {
                "extracted_text": "[版型1]封面標題\n副標\n===\n[版型2]內頁標題\nA\nB",
                "template_text": "{\"slide_masters\":[]}",
                "template_pptx_filename": template_filename,
                "source_filename": "default.pptx",
            }
            with patch.object(views.settings, "MEDIA_ROOT", td):
                resp = self.client.post(
                    "/text2pptx/generate_restored_pptx/",
                    data=json.dumps(payload),
                    content_type="application/json",
                )
                self.assertEqual(resp.status_code, 200)
                data = resp.json()
                self.assertTrue(data.get("success"))
                self.assertTrue(str(data.get("output_filename") or "").endswith(".pptx"))
                self.assertIn("/media/text2pptx/pptx/", data.get("download_url", ""))

                saved_path = os.path.join(td, "text2pptx", "pptx", data["output_filename"])
                prs = Presentation(saved_path)
                self.assertGreaterEqual(len(prs.slides), 2)
                slide0_text = "\n".join(
                    (getattr(sh, "text", "") or "").strip()
                    for sh in prs.slides[0].shapes
                    if getattr(sh, "has_text_frame", False)
                )
                slide1_text = "\n".join(
                    (getattr(sh, "text", "") or "").strip()
                    for sh in prs.slides[1].shapes
                    if getattr(sh, "has_text_frame", False)
                )
                self.assertIn("封面標題", slide0_text)
                self.assertIn("內頁標題", slide1_text)
                self.assertIn("A", slide1_text)
                return

        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertTrue(data.get("success"))
        self.assertTrue(str(data.get("output_filename") or "").endswith(".pptx"))
        self.assertIn("/media/text2pptx/pptx/", data.get("download_url", ""))

        saved_path = os.path.join(td, "text2pptx", "pptx", data["output_filename"])
        prs = Presentation(saved_path)
        self.assertGreaterEqual(len(prs.slides), 2)
        slide0_text = "\n".join(
            (getattr(sh, "text", "") or "").strip()
            for sh in prs.slides[0].shapes
            if getattr(sh, "has_text_frame", False)
        )
        slide1_text = "\n".join(
            (getattr(sh, "text", "") or "").strip()
            for sh in prs.slides[1].shapes
            if getattr(sh, "has_text_frame", False)
        )
        self.assertIn("封面標題", slide0_text)
        self.assertIn("內頁標題", slide1_text)
        self.assertIn("A", slide1_text)

    def test_analyze_image_prompts_rejects_empty_text(self):
        resp = self.client.post("/text2pptx/analyze-image-prompts/", {"text": ""})
        self.assertEqual(resp.status_code, 400)
        self.assertFalse(resp.json().get("ok", True))

    @patch("webapps.text2pptx.views._resolve_analysis_for_image_prompts")
    def test_analyze_image_prompts_returns_preview(self, mock_resolve):
        mock_resolve.return_value = {
            "main_title": "T",
            "main_subtitle": "S",
            "slides": [
                {
                    "title": "Slide A",
                    "slide_type": "content",
                    "image_required": True,
                    "image_prompt": "a clean enterprise illustration",
                    "image_intent": "hero",
                    "aspect_ratio": "16:9",
                }
            ],
        }
        resp = self.client.post("/text2pptx/analyze-image-prompts/", {"text": "demo"})
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertTrue(data.get("ok"))
        self.assertEqual(len(data.get("slides", [])), 1)
        self.assertIn("a clean enterprise illustration", data.get("preview_text", ""))

    def test_generate_rejects_empty_text(self):
        resp = self.client.post("/text2pptx/generate/", {"text": ""})
        self.assertEqual(resp.status_code, 400)
        self.assertIn("請輸入文字內容", resp.json().get("error", ""))

    def test_generate_rejects_too_long_text(self):
        with patch.object(views, "MAX_INPUT_CHARS", 3):
            resp = self.client.post("/text2pptx/generate/", {"text": "1234"})
        self.assertEqual(resp.status_code, 400)
        self.assertIn("文字內容超過上限", resp.json().get("error", ""))

    def test_index_lists_chinese_template_filename(self):
        with tempfile.TemporaryDirectory(dir=os.getcwd(), ignore_cleanup_errors=True) as td:
            with open(os.path.join(td, "中心測試.pptx"), "wb") as f:
                f.write(b"x")
            with patch.object(views, "PPTX_TEMPLATE_DIR", td):
                resp = self.client.get("/text2pptx/")
        self.assertEqual(resp.status_code, 200)
        self.assertContains(resp, "投影片可引用版面範例")
        self.assertContains(resp, "中心測試.pptx")
        self.assertContains(resp, "預設範本缺失，改用內建")
        self.assertContains(resp, "引用範例並產生 PPTX")
        self.assertContains(resp, "解析內容轉生圖提示詞")
        self.assertContains(resp, "前往範本管理（管理者）")
        self.assertNotContains(resp, "匯入範本（含四種版型審查）")

    def test_index_uses_default_template_when_available(self):
        with tempfile.TemporaryDirectory(dir=os.getcwd(), ignore_cleanup_errors=True) as td:
            with open(os.path.join(td, "預設範本.pptx"), "wb") as f:
                f.write(b"x")
            with open(os.path.join(td, "中心測試.pptx"), "wb") as f:
                f.write(b"x")
            with patch.object(views, "PPTX_TEMPLATE_DIR", td):
                resp = self.client.get("/text2pptx/")
        self.assertEqual(resp.status_code, 200)
        self.assertContains(resp, "（預設範本：預設範本.pptx）")

    def test_template_admin_page_renders_import_panel(self):
        resp = self.client.get("/text2pptx/template-admin/")
        self.assertEqual(resp.status_code, 200)
        self.assertContains(resp, "TEXT → PPTX 範本管理")
        self.assertContains(resp, "匯入範本（含四種版型審查）")

    @patch("webapps.text2pptx.views._audit_template_bytes")
    def test_import_template_rejects_when_required_layout_missing(self, mock_audit):
        mock_audit.return_value = {
            "ok": False,
            "missing": ["章節標題（Section Header）"],
            "found": {},
        }
        f = SimpleUploadedFile(
            "bad_template.pptx",
            b"fake-pptx-content",
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        resp = self.client.post("/text2pptx/import-template/", {"template_file": f}, follow=True)
        self.assertEqual(resp.status_code, 200)
        self.assertContains(resp, "TEXT → PPTX 範本管理")
        self.assertContains(resp, "匯入失敗，缺少必要版型")

    @patch("webapps.text2pptx.views._save_template_bytes")
    @patch("webapps.text2pptx.views._audit_template_bytes")
    def test_import_template_success_message(self, mock_audit, mock_save):
        mock_audit.return_value = {
            "ok": True,
            "missing": [],
            "found": {
                "cover": "Title Slide",
                "section": "Section Header",
                "content": "Title and Content",
                "two_content": "Two Content",
            },
        }
        mock_save.return_value = "uploaded_ok.pptx"
        f = SimpleUploadedFile(
            "ok_template.pptx",
            b"fake-pptx-content",
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        resp = self.client.post("/text2pptx/import-template/", {"template_file": f}, follow=True)
        self.assertEqual(resp.status_code, 200)
        self.assertContains(resp, "TEXT → PPTX 範本管理")
        self.assertContains(resp, "匯入成功：uploaded_ok.pptx")

    @patch("webapps.text2pptx.views._analyze_text_with_llm")
    def test_generate_uses_markers_without_calling_llm(self, mock_analyze):
        resp = self.client.post(
            "/text2pptx/generate/",
            {
                "text": "[標題投影片] 封面標題\n副標\n===\n[內容頁] 內容頁\n重點一\n重點二",
                "title": "",
                "template": "",
            },
        )
        self.assertEqual(resp.status_code, 200)
        mock_analyze.assert_not_called()

    @patch("webapps.text2pptx.views._analyze_text_with_llm")
    def test_section_slide_keeps_all_bullets(self, mock_analyze):
        mock_analyze.return_value = {
            "main_title": "主標",
            "main_subtitle": "副標",
            "slides": [
                {
                    "title": "章節標題",
                    "slide_type": "section",
                    "bullets": ["重點一", "重點二", "重點三"],
                }
            ],
        }
        resp = self.client.post("/text2pptx/generate/", {"text": "x", "title": "", "template": ""})
        self.assertEqual(resp.status_code, 200)

        from pptx import Presentation

        prs = Presentation(io.BytesIO(resp.content))
        self.assertEqual(len(prs.slides), 2)  # cover + section
        all_text = []
        for sh in prs.slides[1].shapes:
            if getattr(sh, "has_text_frame", False):
                all_text.append((sh.text or "").strip())
        joined = "\n".join(all_text)
        self.assertIn("重點一", joined)
        self.assertIn("重點二", joined)
        self.assertIn("重點三", joined)

    @patch("webapps.text2pptx.views._analyze_text_with_llm")
    def test_generate_supports_multiple_slide_types_and_title_override(self, mock_analyze):
        mock_analyze.return_value = {
            "main_title": "AI 產生主標",
            "main_subtitle": "AI 副標",
            "slides": [
                {
                    "title": "章節標題",
                    "slide_type": "section",
                    "bullets": ["重點摘要"],
                },
                {
                    "title": "雙欄必較頁",
                    "slide_type": "comparison",
                    "left_title": "已完成",
                    "right_title": "待完成",
                    "bullets": [f"項目{i}" for i in range(1, 18)],
                },
                {
                    "title": "一般內容",
                    "slide_type": "content",
                    "bullets": ["A", "B", "C"],
                },
            ],
        }

        resp = self.client.post(
            "/text2pptx/generate/",
            {
                "text": "任意原文",
                "title": "使用者自訂主標",
                "template": "",
            },
        )
        self.assertEqual(resp.status_code, 200)
        self.assertEqual(
            resp["Content-Type"],
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        cd_raw = resp["Content-Disposition"]
        self.assertIn("filename*=", cd_raw)
        encoded_name = cd_raw.split("filename*=UTF-8''", 1)[1]
        self.assertIn("使用者自訂主標.pptx", unquote(encoded_name))

        from pptx import Presentation

        prs = Presentation(io.BytesIO(resp.content))
        # cover(1) + section(1) + comparison(17 bullets => 2 pages) + content(1) = 5
        self.assertEqual(len(prs.slides), 5)

    @patch("webapps.text2pptx.views._generate_image_for_slide", return_value=None)
    def test_generate_pptx_uses_marked_layout_names(self, _mock_image):
        from pptx import Presentation

        with tempfile.TemporaryDirectory(dir=os.getcwd(), ignore_cleanup_errors=True) as td:
            template_prs = Presentation()
            if len(template_prs.slide_layouts) > 0:
                template_prs.slide_layouts[0].name = "版型1"
            if len(template_prs.slide_layouts) > 1:
                template_prs.slide_layouts[1].name = "版型2"
            template_path = os.path.join(td, "marked_layouts_template.pptx")
            template_prs.save(template_path)

            marked_text = "\n===\n".join(
                [
                    "[版型1]Cover Title\nSubtitle A\nSubtitle B",
                    "[版型2]Body Title\nPoint A\nPoint B",
                ]
            )

            with patch("webapps.text2pptx.views._safe_select_template", return_value=template_path):
                resp = self.client.post(
                    "/text2pptx/generate/",
                    {
                        "text": marked_text,
                        "title": "",
                        "template": "ignored",
                    },
                )

        self.assertEqual(resp.status_code, 200)
        prs = Presentation(io.BytesIO(resp.content))
        self.assertGreaterEqual(len(prs.slides), 2)
        self.assertEqual(prs.slides[0].slide_layout.name, "版型1")
        self.assertEqual(prs.slides[1].slide_layout.name, "版型2")

        with zipfile.ZipFile(io.BytesIO(resp.content), "r") as zf:
            slide1_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8", errors="strict")
            slide2_xml = zf.read("ppt/slides/slide2.xml").decode("utf-8", errors="strict")
        self.assertIn("Subtitle A", slide1_xml)
        self.assertIn("Point A", slide2_xml)

    def test_restore_slide_text_distributes_across_multiple_textboxes(self):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = "old title"

        slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(4.0), Inches(0.8)).text_frame.text = "box1"
        slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(4.0), Inches(0.8)).text_frame.text = "box2"

        views._restore_slide_text_from_analysis(
            slide,
            {
                "title": "New Title",
                "subtitle": "Line 1\nLine 2\nLine 3\nLine 4\nLine 5",
            },
            is_cover=True,
        )

        self.assertEqual(getattr(slide.shapes.title, "text", ""), "New Title")
        frames = views._find_content_text_frames(slide, exclude_shape=slide.shapes.title)
        self.assertGreaterEqual(len(frames), 3)
        frame_texts = [getattr(frame, "text", "") for frame in frames[:3]]
        self.assertIn("Line 1", frame_texts[0])
        self.assertIn("Line 2", frame_texts[0])
        self.assertIn("Line 3", frame_texts[1])
        self.assertIn("Line 4", frame_texts[1])
        self.assertIn("Line 5", frame_texts[2])

    def test_convert_pptx_keeps_textbox_shape_body_on_slide(self):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(3.0), Inches(1.0))
        textbox.text = "Keep me"

        buf = io.BytesIO()
        prs.save(buf)

        converted = views._convert_pptx_bytes_to_pptx_bytes(buf.getvalue())
        with zipfile.ZipFile(io.BytesIO(converted), "r") as zf:
            slide_root = ET.fromstring(zf.read("ppt/slides/slide1.xml"))
            ns = {
                "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
                "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            }
            sp_nodes = slide_root.findall(".//p:sp", ns)
            self.assertTrue(sp_nodes)
            self.assertTrue(any(node.find("p:txBody", ns) is not None for node in sp_nodes))
            text_nodes = [str(n.text or "").strip() for n in slide_root.findall(".//a:t", ns)]
            self.assertFalse(any(text_nodes))

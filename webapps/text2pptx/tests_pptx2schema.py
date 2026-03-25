from __future__ import annotations

import io
import tempfile

from django.core.files.uploadedfile import SimpleUploadedFile
from django.test import Client, SimpleTestCase, override_settings

from webapps.text2pptx.pptx2schema.pipelines.extract_pipeline import run_extract


def _build_minimal_pptx_bytes() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title:
        slide.shapes.title.text = "Schema Demo"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Subtitle"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class Pptx2SchemaPipelineTests(SimpleTestCase):
    def test_run_extract_from_bytes(self):
        raw = run_extract(_build_minimal_pptx_bytes(), source_file="demo.pptx")
        self.assertEqual(raw.meta.slide_count, 1)
        self.assertEqual(raw.meta.source_file, "demo.pptx")
        self.assertGreaterEqual(len(raw.slides[0].shapes), 1)


@override_settings(PORTAL_ACL_ENABLED=False)
class Pptx2SchemaApiTests(SimpleTestCase):
    databases = {"default"}

    def setUp(self):
        self.client = Client()

    def test_schema_extract_requires_post(self):
        resp = self.client.get("/djangoai/text2pptx/schema/extract/")
        self.assertEqual(resp.status_code, 405)

    def test_schema_extract_success(self):
        upload = SimpleUploadedFile(
            "demo.pptx",
            _build_minimal_pptx_bytes(),
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        resp = self.client.post("/djangoai/text2pptx/schema/extract/", data={"pptx_file": upload})
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertTrue(data.get("ok"))
        self.assertEqual(data["raw"]["meta"]["slide_count"], 1)

    def test_schema_analyze_success(self):
        upload = SimpleUploadedFile(
            "demo.pptx",
            _build_minimal_pptx_bytes(),
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        resp = self.client.post("/djangoai/text2pptx/schema/analyze/", data={"pptx_file": upload})
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertTrue(data.get("ok"))
        self.assertIn("template_dna", data["bundle"])

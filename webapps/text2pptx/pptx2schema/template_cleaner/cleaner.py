from __future__ import annotations

import io
import json
import statistics
import zipfile
from collections import Counter
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Tuple
from xml.etree import ElementTree as ET

from pptx import Presentation

from webapps.text2pptx.pptx2schema.analyzers.color_tokenizer import extract_color_tokens
from webapps.text2pptx.pptx2schema.analyzers.typography_analyzer import extract_font_tokens
from webapps.text2pptx.pptx2schema.pipelines.extract_pipeline import run_extract
from webapps.text2pptx.pptx2schema.template_cleaner.classifier import ObservedShape, classify_shape
from webapps.text2pptx.pptx2schema.template_cleaner.models import (
    CleanerConfig,
    KeepRemovePolicySummary,
    LayoutAsset,
    MasterAsset,
    PlaceholderAsset,
    ShapeDecision,
    TemplateCleanerSchema,
    ThemeAssets,
)
from webapps.text2pptx.pptx2schema.models.raw import BBox


TITLE_PLACEHOLDER_TYPES = {"TITLE", "CENTER_TITLE"}
BODY_PLACEHOLDER_TYPES = {"BODY", "OBJECT", "PICTURE"}
CAPTION_PLACEHOLDER_TYPES = {"SUBTITLE", "DATE", "FOOTER", "SLIDE_NUMBER"}


def _shape_type_name(shape: object) -> str:
    return str(getattr(getattr(shape, "shape_type", None), "name", "") or "OTHER")


def _placeholder_type_name(shape: object) -> Optional[str]:
    try:
        if not getattr(shape, "is_placeholder", False):
            return None
        return str(getattr(getattr(shape.placeholder_format, "type", None), "name", "") or "") or None
    except Exception:
        return None


def _shape_bbox(shape: object) -> BBox:
    return BBox(
        x=float(getattr(shape, "left", 0) or 0),
        y=float(getattr(shape, "top", 0) or 0),
        w=float(getattr(shape, "width", 0) or 0),
        h=float(getattr(shape, "height", 0) or 0),
    )


def _shape_has_text(shape: object) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    try:
        paragraphs = getattr(shape.text_frame, "paragraphs", [])
        return any(str(getattr(p, "text", "") or "").strip() for p in paragraphs)
    except Exception:
        return False


def _shape_position_key(shape: object, *, slide_w: float, slide_h: float, digits: int) -> Tuple[object, ...]:
    bbox = _shape_bbox(shape)
    ptype = _placeholder_type_name(shape) or ""
    sw = slide_w or 1.0
    sh = slide_h or 1.0
    return (
        _shape_type_name(shape),
        ptype,
        round(bbox.x / sw, digits),
        round(bbox.y / sh, digits),
        round(bbox.w / sw, digits),
        round(bbox.h / sh, digits),
    )


def _build_repetition_counter(prs: Presentation, *, cfg: CleanerConfig) -> Counter[Tuple[object, ...]]:
    slide_w = float(prs.slide_width or 1.0)
    slide_h = float(prs.slide_height or 1.0)
    key_counter: Counter[Tuple[object, ...]] = Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            key = _shape_position_key(
                shape,
                slide_w=slide_w,
                slide_h=slide_h,
                digits=cfg.repeat_position_round_digits,
            )
            key_counter[key] += 1
    return key_counter


def _clear_shape_text(shape: object) -> None:
    if getattr(shape, "has_text_frame", False):
        try:
            shape.text_frame.clear()
        except Exception:
            pass
    if getattr(shape, "has_table", False):
        try:
            for row in shape.table.rows:
                for cell in row.cells:
                    cell.text = ""
        except Exception:
            pass


def _remove_shape(shape: object) -> None:
    try:
        node = shape._element  # pylint: disable=protected-access
        node.getparent().remove(node)
    except Exception:
        pass


def _infer_layout_type(layout_name: str, *, title_count: int, body_count: int) -> str:
    low = (layout_name or "").lower()
    if "section" in low or "章節" in low:
        return "section"
    if "comparison" in low or "比較" in low:
        return "comparison"
    if "title slide" in low or "標題投影片" in low:
        return "cover"
    if body_count >= 2:
        return "two_content"
    if title_count >= 1 and body_count >= 1:
        return "content"
    if title_count >= 1 and body_count == 0:
        return "section"
    return "content"


def _collect_master_assets(prs: Presentation) -> List[MasterAsset]:
    out: List[MasterAsset] = []
    for master in prs.slide_masters:
        master_name = str(getattr(master, "name", "") or "") or "(unnamed-master)"
        layout_names = [str(getattr(layout, "name", "") or "") or "(unnamed-layout)" for layout in master.slide_layouts]
        out.append(MasterAsset(master_name=master_name, layout_names=layout_names))
    return out


def _collect_layout_assets(prs: Presentation) -> List[LayoutAsset]:
    out: List[LayoutAsset] = []
    seen: set[Tuple[str, str]] = set()
    for master in prs.slide_masters:
        master_name = str(getattr(master, "name", "") or "") or "(unnamed-master)"
        for layout in master.slide_layouts:
            layout_name = str(getattr(layout, "name", "") or "") or "(unnamed-layout)"
            key = (master_name, layout_name)
            if key in seen:
                continue
            seen.add(key)

            placeholders: List[PlaceholderAsset] = []
            title_regions: List[BBox] = []
            body_regions: List[BBox] = []
            caption_regions: List[BBox] = []

            for ph in getattr(layout, "placeholders", []):
                ptype = _placeholder_type_name(ph)
                bbox = _shape_bbox(ph)
                pidx = None
                try:
                    pidx = int(getattr(ph.placeholder_format, "idx", 0))
                except Exception:
                    pidx = None

                placeholders.append(
                    PlaceholderAsset(
                        placeholder_type=ptype,
                        placeholder_idx=pidx,
                        bbox=bbox,
                    )
                )
                if ptype in TITLE_PLACEHOLDER_TYPES:
                    title_regions.append(bbox)
                elif ptype in BODY_PLACEHOLDER_TYPES:
                    body_regions.append(bbox)
                elif ptype in CAPTION_PLACEHOLDER_TYPES:
                    caption_regions.append(bbox)
                else:
                    body_regions.append(bbox)

            layout_type = _infer_layout_type(
                layout_name,
                title_count=len(title_regions),
                body_count=len(body_regions),
            )
            out.append(
                LayoutAsset(
                    master_name=master_name,
                    layout_name=layout_name,
                    layout_type=layout_type,
                    placeholders=placeholders,
                    title_regions=title_regions,
                    body_regions=body_regions,
                    caption_regions=caption_regions,
                )
            )
    return out


def _extract_theme_assets_from_bytes(raw_bytes: bytes) -> ThemeAssets:
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    with zipfile.ZipFile(io.BytesIO(raw_bytes), "r") as zf:
        theme_parts = sorted([name for name in zf.namelist() if name.startswith("ppt/theme/theme") and name.endswith(".xml")])
        if not theme_parts:
            return ThemeAssets()
        theme_part = theme_parts[0]
        xml_raw = zf.read(theme_part)

    root = ET.fromstring(xml_raw)
    font_scheme: Dict[str, str] = {}
    color_scheme: Dict[str, str] = {}

    major = root.find(".//a:fontScheme/a:majorFont/a:latin", ns)
    minor = root.find(".//a:fontScheme/a:minorFont/a:latin", ns)
    if major is not None and major.get("typeface"):
        font_scheme["major_latin"] = str(major.get("typeface"))
    if minor is not None and minor.get("typeface"):
        font_scheme["minor_latin"] = str(minor.get("typeface"))

    for key in ["dk1", "lt1", "dk2", "lt2", "accent1", "accent2", "accent3", "accent4", "accent5", "accent6", "hlink", "folHlink"]:
        node = root.find(f".//a:clrScheme/a:{key}", ns)
        if node is None:
            continue
        rgb_node = node.find("a:srgbClr", ns)
        sys_node = node.find("a:sysClr", ns)
        value = None
        if rgb_node is not None and rgb_node.get("val"):
            value = rgb_node.get("val")
        elif sys_node is not None:
            value = sys_node.get("lastClr") or sys_node.get("val")
        if value:
            color_scheme[key] = str(value).upper()

    return ThemeAssets(
        theme_part=theme_part,
        font_scheme=font_scheme,
        color_scheme=color_scheme,
    )


def _safe_median(values: Iterable[float]) -> float:
    seq = [float(x) for x in values if x is not None]
    if not seq:
        return 0.0
    return float(statistics.median(seq))


def _compute_spacing_rules(layouts: List[LayoutAsset]) -> Dict[str, float]:
    title_body_gaps: List[float] = []
    body_gaps: List[float] = []

    for layout in layouts:
        if layout.title_regions and layout.body_regions:
            title_bottom = max([b.y + b.h for b in layout.title_regions])
            first_body_top = min([b.y for b in layout.body_regions])
            title_body_gaps.append(max(0.0, first_body_top - title_bottom))

        sorted_body = sorted(layout.body_regions, key=lambda b: (b.y, b.x))
        for idx in range(len(sorted_body) - 1):
            gap = sorted_body[idx + 1].y - (sorted_body[idx].y + sorted_body[idx].h)
            if gap >= 0:
                body_gaps.append(gap)

    return {
        "title_to_body_gap": _safe_median(title_body_gaps),
        "body_vertical_gap": _safe_median(body_gaps),
    }


def _build_policy_summary(decisions: List[ShapeDecision]) -> KeepRemovePolicySummary:
    by_category: Dict[str, int] = {}
    by_strategy: Dict[str, int] = {}
    for d in decisions:
        by_category[d.category] = by_category.get(d.category, 0) + 1
        by_strategy[d.strategy] = by_strategy.get(d.strategy, 0) + 1

    return KeepRemovePolicySummary(
        total_shapes=len(decisions),
        kept_shapes=by_strategy.get("keep_all", 0),
        cleared_text_shapes=by_strategy.get("clear_text_keep_shape", 0),
        removed_shapes=by_strategy.get("remove_shape", 0),
        by_category=by_category,
        by_strategy=by_strategy,
    )


def _sort_shape_id(value: str) -> Tuple[int, str]:
    text = str(value or "")
    try:
        return (0, str(int(text)))
    except Exception:
        return (1, text)


def clean_template_assets(
    input_pptx: str | Path,
    output_cleaned_pptx: str | Path,
    output_schema_json: str | Path,
    *,
    source_file: str | None = None,
    config: CleanerConfig | None = None,
) -> TemplateCleanerSchema:
    cfg = config or CleanerConfig()
    input_path = Path(input_pptx)
    output_pptx_path = Path(output_cleaned_pptx)
    output_schema_path = Path(output_schema_json)

    raw_bytes = input_path.read_bytes()
    prs = Presentation(str(input_path))
    raw = run_extract(str(input_path), source_file=source_file or str(input_path))

    slide_w = float(prs.slide_width or 1.0)
    slide_h = float(prs.slide_height or 1.0)
    repetition_counter = _build_repetition_counter(prs, cfg=cfg)

    decisions: List[ShapeDecision] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            bbox = _shape_bbox(shape)
            key = _shape_position_key(
                shape,
                slide_w=slide_w,
                slide_h=slide_h,
                digits=cfg.repeat_position_round_digits,
            )
            repeated_count = int(repetition_counter.get(key, 1))
            area_ratio = (bbox.w * bbox.h) / (slide_w * slide_h) if slide_w > 0 and slide_h > 0 else 0.0
            y_ratio = (bbox.y / slide_h) if slide_h > 0 else 0.0
            observed = ObservedShape(
                is_placeholder=bool(getattr(shape, "is_placeholder", False)),
                has_text=_shape_has_text(shape),
                is_picture=(_shape_type_name(shape) == "PICTURE"),
                repeated_count=repeated_count,
                area_ratio=area_ratio,
                y_ratio=y_ratio,
                shape_type=_shape_type_name(shape),
            )
            category, strategy, reasons = classify_shape(observed, config=cfg)
            decisions.append(
                ShapeDecision(
                    slide_index=slide_index,
                    shape_id=str(getattr(shape, "shape_id", "")),
                    name=str(getattr(shape, "name", "") or "") or None,
                    shape_type=_shape_type_name(shape),
                    placeholder_type=_placeholder_type_name(shape),
                    bbox=bbox,
                    has_text=observed.has_text,
                    is_picture=observed.is_picture,
                    repeated_count=repeated_count,
                    category=category,
                    strategy=strategy,
                    reasons=reasons,
                )
            )

    decision_map: Dict[Tuple[int, str], ShapeDecision] = {
        (d.slide_index, d.shape_id): d for d in decisions
    }
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in list(slide.shapes)[::-1]:
            key = (slide_index, str(getattr(shape, "shape_id", "")))
            d = decision_map.get(key)
            if d is None:
                continue
            if d.strategy == "remove_shape":
                _remove_shape(shape)
            elif d.strategy == "clear_text_keep_shape":
                _clear_shape_text(shape)

    output_pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx_path))

    theme_assets = _extract_theme_assets_from_bytes(raw_bytes)
    layout_assets = _collect_layout_assets(prs)
    summary = _build_policy_summary(decisions)

    schema = TemplateCleanerSchema(
        source_file=source_file or str(input_path),
        cleaned_template_file=str(output_pptx_path),
        theme=theme_assets,
        slide_masters=_collect_master_assets(prs),
        layouts=layout_assets,
        font_tokens=extract_font_tokens(raw),
        color_tokens=extract_color_tokens(raw),
        spacing_rules=_compute_spacing_rules(layout_assets),
        keep_remove_policy_summary=summary,
        shape_decisions=sorted(decisions, key=lambda x: (x.slide_index, _sort_shape_id(x.shape_id))),
    )

    output_schema_path.parent.mkdir(parents=True, exist_ok=True)
    output_schema_path.write_text(
        json.dumps(schema.model_dump(), ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return schema

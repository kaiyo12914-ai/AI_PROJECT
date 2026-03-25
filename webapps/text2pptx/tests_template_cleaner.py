from __future__ import annotations

import base64
import json
import os
import tempfile
from pathlib import Path

from django.test import SimpleTestCase, override_settings

from webapps.text2pptx.pptx2schema.template_cleaner.classifier import ObservedShape, classify_shape
from webapps.text2pptx.pptx2schema.template_cleaner.cleaner import clean_template_assets
from webapps.text2pptx.pptx2schema.template_cleaner.models import CleanerConfig


def _tiny_png_bytes() -> bytes:
    return base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8"
        "/w8AAusB9Y0A3L8AAAAASUVORK5CYII="
    )


class TemplateCleanerRuleTests(SimpleTestCase):
    def test_classify_placeholder_is_structural_clear_text(self):
        category, strategy, _reasons = classify_shape(
            ObservedShape(
                is_placeholder=True,
                has_text=True,
                is_picture=False,
                repeated_count=1,
                area_ratio=0.1,
                y_ratio=0.2,
                shape_type="TEXT_BOX",
            ),
            config=CleanerConfig(),
        )
        self.assertEqual(category, "structural")
        self.assertEqual(strategy, "clear_text_keep_shape")

    def test_classify_single_picture_is_removed(self):
        category, strategy, _reasons = classify_shape(
            ObservedShape(
                is_placeholder=False,
                has_text=False,
                is_picture=True,
                repeated_count=1,
                area_ratio=0.2,
                y_ratio=0.1,
                shape_type="PICTURE",
            ),
            config=CleanerConfig(),
        )
        self.assertEqual(category, "content")
        self.assertEqual(strategy, "remove_shape")


@override_settings(PORTAL_ACL_ENABLED=False)
class TemplateCleanerE2ETests(SimpleTestCase):
    databases = {"default"}

    def test_clean_template_assets_outputs_cleaned_pptx_and_schema(self):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        with tempfile.TemporaryDirectory(dir=os.getcwd(), ignore_cleanup_errors=True) as td:
            work = Path(td)
            input_pptx = work / "source.pptx"
            output_pptx = work / "cleaned_template.pptx"
            output_schema = work / "template_schema.json"
            image_path = work / "logo.png"
            image_path.write_bytes(_tiny_png_bytes())

            prs = Presentation()
            slide1 = prs.slides.add_slide(prs.slide_layouts[1])
            if slide1.shapes.title:
                slide1.shapes.title.text = "首頁標題"
            if len(slide1.placeholders) > 1:
                slide1.placeholders[1].text = "首頁內容"
            slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(6.8), Inches(1.5), Inches(0.3))

            slide2 = prs.slides.add_slide(prs.slide_layouts[1])
            if slide2.shapes.title:
                slide2.shapes.title.text = "第二頁標題"
            if len(slide2.placeholders) > 1:
                slide2.placeholders[1].text = "第二頁內容"
            slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(6.8), Inches(1.5), Inches(0.3))

            unique = slide2.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(4.0), Inches(1.0))
            unique.text_frame.text = "UNIQUE-CONTENT-REMOVE-ME"
            slide2.shapes.add_picture(str(image_path), Inches(6.0), Inches(1.0), Inches(2.0), Inches(2.0))

            prs.save(str(input_pptx))

            schema = clean_template_assets(
                input_pptx=input_pptx,
                output_cleaned_pptx=output_pptx,
                output_schema_json=output_schema,
                source_file=str(input_pptx),
            )

            self.assertTrue(output_pptx.exists())
            self.assertTrue(output_schema.exists())
            self.assertGreater(schema.keep_remove_policy_summary.total_shapes, 0)
            self.assertGreaterEqual(schema.keep_remove_policy_summary.removed_shapes, 1)
            self.assertGreaterEqual(schema.keep_remove_policy_summary.cleared_text_shapes, 1)
            self.assertGreaterEqual(len(schema.layouts), 1)
            self.assertIn("title_to_body_gap", schema.spacing_rules)

            from pptx import Presentation as CleanedPresentation

            cleaned = CleanedPresentation(str(output_pptx))
            all_text = []
            all_shape_types = []
            for slide in cleaned.slides:
                for shape in slide.shapes:
                    st = str(getattr(getattr(shape, "shape_type", None), "name", "") or "")
                    all_shape_types.append(st)
                    if getattr(shape, "has_text_frame", False):
                        all_text.append(str(getattr(shape, "text", "") or ""))

            merged_text = "\n".join(all_text)
            self.assertNotIn("UNIQUE-CONTENT-REMOVE-ME", merged_text)
            self.assertNotIn("首頁標題", merged_text)
            self.assertNotIn("第二頁標題", merged_text)
            self.assertNotIn("PICTURE", all_shape_types)

            schema_json = json.loads(output_schema.read_text(encoding="utf-8"))
            self.assertIn("layouts", schema_json)
            self.assertIn("font_tokens", schema_json)
            self.assertIn("color_tokens", schema_json)
            self.assertIn("keep_remove_policy_summary", schema_json)


import re
import time
from urllib.parse import urlparse
from typing import Any, Dict, List
from collections import Counter
from django.http import HttpRequest
from django.http import HttpResponseForbidden
from django.shortcuts import render
from django.views.decorators.csrf import csrf_exempt
from django.db import ProgrammingError
from django.db.models import Q
from pgvector.django import L2Distance
from webapps.portal.decorators import require_node
from webapps.portal import decorators as portal_decorators
from webapps.llm.llm_factory import get_chat_model

from .models import (
    Project,
    Source,
    Document,
    DocumentVersion,
    DocumentChunk,
    Conversation,
    Message,
    MessageCitation,
    ActivityLog,
)
from .lang_guard import prefer_traditional_chinese, is_zh_dominant
from .citation_guard import ensure_sentence_citations, detect_citation_conflicts
from .retrieval_policy import (
    is_definition_query,
    definition_chunk_boost,
    generic_source_penalty,
    build_sparse_terms,
    rerank_candidates,
)
from .query_rewrite import rewrite_query_for_retrieval
from .api_helpers import (
    api_error as _api_error,
    is_bad_utf8_request as _is_bad_utf8_request,
    read_json_body as _read_json_body,
    safe_json_response as _safe_json_response,
    safe_text as _safe_text,
    to_int as _to_int,
)
from .embedding_service import mock_embedding as _mock_embedding
from .text_processing import (
    build_chunks as _build_chunks,
    clean_label as _clean_label,
    decode_text_bytes_best_effort as _decode_text_bytes_best_effort,
    preprocess_rag_text as _preprocess_rag_text,
)

# --- UTILS ---

def _current_user_id(request: HttpRequest) -> str:
    # Use login_user if available (from custom middleware), else request.user
    uid = getattr(request, "login_user", None)
    if uid:
        return str(uid)
    user = getattr(request, "user", None)
    if user and not user.is_anonymous:
        return user.username
    return ""

def _can_manage_projects(request: HttpRequest) -> bool:
    return portal_decorators.can_access(getattr(request, "user", None), "portal")


def _log_activity(
    *,
    project_id: int,
    action: str,
    user_id: str = "",
    target_type: str = "",
    target_id: int | None = None,
    detail: Dict[str, Any] | None = None,
) -> None:
    if project_id <= 0 or not action:
        return
    try:
        ActivityLog.objects.create(
            project_id=project_id,
            user_id=_safe_text(user_id),
            action=_safe_text(action),
            target_type=_safe_text(target_type),
            target_id=target_id if target_id and target_id > 0 else None,
            detail_json=detail or {},
        )
    except Exception:
        # Activity logging must not break user flows.
        return


def _db_write_error_response(exc: Exception, fallback: str = "database write failed"):
    msg = _safe_text(exc)
    if "permission denied" in msg.lower():
        return _api_error(
            "projectnotes database permission denied; please grant INSERT/UPDATE/DELETE on required tables",
            error_code="db_permission_denied",
            status=500,
        )
    return _api_error(fallback, error_code="db_write_failed", status=500)


def _suggest_conversation_title(question: str, conversation_id: int = 0) -> str:
    text = _safe_text(question).replace("\r", " ").replace("\n", " ").strip()
    text = re.sub(r"\s+", " ", text)
    text = text.lstrip("Qq：:;,.?？!！ ")
    title = _clean_label(text, "", 60)
    if not title:
        return f"對話 {conversation_id}" if conversation_id > 0 else "新對話"
    return title


def _maybe_update_conversation_title(conv: Conversation, question: str) -> str:
    if not conv:
        return ""
    current = _safe_text(getattr(conv, "title", "")).strip()
    normalized = current.lower()
    defaults = {"", "new chat", "新對話"}
    if normalized not in defaults:
        return current
    title = _suggest_conversation_title(question, getattr(conv, "id", 0))
    if title and title != current:
        conv.title = title
        conv.save(update_fields=["title", "updated_at"])
    return _safe_text(conv.title)


def _persist_message_citations(message: Message, citations: List[Dict[str, Any]]) -> None:
    if not message or not citations:
        return
    chunk_keys = []
    for item in citations:
        source_id = _to_int(item.get("source_id"), 0)
        chunk_index = _to_int(item.get("chunk_index"), -1)
        if source_id <= 0 or chunk_index < 0:
            continue
        chunk_keys.append((source_id, chunk_index, _safe_text(item.get("ref"))))
    if not chunk_keys:
        return

    source_ids = list({sid for sid, _, _ in chunk_keys})
    chunk_indexes = list({idx for _, idx, _ in chunk_keys})
    rows = (
        DocumentChunk.objects
        .filter(
            chunk_index__in=chunk_indexes,
            document_version__document__source_id__in=source_ids,
        )
        .select_related("document_version__document")
    )
    key_map: Dict[tuple[int, int], DocumentChunk] = {}
    for row in rows:
        key = (row.document_version.document.source_id, row.chunk_index)
        if key not in key_map:
            key_map[key] = row

    inserts: List[MessageCitation] = []
    for source_id, chunk_index, ref in chunk_keys:
        chunk = key_map.get((source_id, chunk_index))
        if not chunk:
            continue
        inserts.append(
            MessageCitation(
                message=message,
                document_chunk=chunk,
                citation_text=ref,
            )
        )
    if inserts:
        MessageCitation.objects.bulk_create(inserts)


def _build_project_overview(project_id: int) -> Dict[str, Any]:
    source_rows = list(Source.objects.filter(project_id=project_id).order_by("name", "id"))
    source_count = len(source_rows)

    latest_versions = list(
        DocumentVersion.objects
        .filter(document__source__project_id=project_id)
        .select_related("document__source")
        .order_by("document_id", "-version_number", "-id")
    )
    latest_by_document: Dict[int, DocumentVersion] = {}
    for row in latest_versions:
        if row.document_id not in latest_by_document:
            latest_by_document[row.document_id] = row

    version_ids = [row.id for row in latest_by_document.values()]
    chunk_rows = list(
        DocumentChunk.objects
        .filter(document_version_id__in=version_ids)
        .values("chunk_index", "content", "document_version__document__source__name")
    )
    chunk_count = len(chunk_rows)

    activity_rows = list(
        ActivityLog.objects
        .filter(project_id=project_id)
        .order_by("-created_at", "-id")[:300]
    )
    chat_rows = [r for r in activity_rows if _safe_text(r.action) == "chat_query"]
    upload_rows = [r for r in activity_rows if _safe_text(r.action) == "source_upload"]

    summary_bits = [f"目前共有 {source_count} 份來源"]
    if chunk_count > 0:
        summary_bits.append(f"{chunk_count} 個可檢索 chunks")
    if chat_rows:
        summary_bits.append(f"近期待辦互動 {len(chat_rows)} 筆查詢")
    if upload_rows:
        summary_bits.append(f"{len(upload_rows)} 筆來源上傳紀錄")

    top_sources = [(_safe_text(s.name) or f"source#{s.id}") for s in source_rows[:3]]
    if top_sources:
        summary_bits.append("主要來源：" + "、".join(top_sources))
    summary = "；".join(summary_bits) + "。"

    faq: List[str] = []
    seen_queries = set()
    for row in chat_rows:
        detail = row.detail_json if isinstance(row.detail_json, dict) else {}
        query = _safe_text(detail.get("query"))
        if not query or query in seen_queries:
            continue
        seen_queries.add(query)
        faq.append(query)
        if len(faq) >= 3:
            break

    keyword_counter: Counter[str] = Counter()
    for row in source_rows:
        for token in _query_tokens(row.name):
            if len(token) >= 2:
                keyword_counter[token] += 3
    for row in chunk_rows[:200]:
        for token in _query_tokens(_safe_text(row.get("content"))):
            if len(token) >= 2:
                keyword_counter[token] += 1
    stopwords = {
        "project", "notes", "source", "chunk", "this", "that", "with", "from",
        "以及", "相關", "內容", "資料", "來源", "文件", "專案", "工作", "辦理", "目前",
    }
    keywords = [tok for tok, _ in keyword_counter.most_common(8) if tok not in stopwords][:8]

    decisions: List[str] = []
    decision_patterns = ("決議", "結論", "辦理", "應", "須", "得", "採購", "履約", "決定")
    for row in chunk_rows[:200]:
        content = _safe_text(row.get("content"))
        if not content:
            continue
        if not any(p in content for p in decision_patterns):
            continue
        clean = re.sub(r"\s+", " ", content).strip()
        if len(clean) > 80:
            clean = clean[:80].rstrip() + "..."
        if clean and clean not in decisions:
            decisions.append(clean)
        if len(decisions) >= 3:
            break

    return {
        "summary": summary,
        "faq": faq,
        "keywords": keywords,
        "decisions": decisions,
    }

_TOKEN_RE = re.compile(r"[A-Za-z0-9\u4e00-\u9fff]{2,}")


def _query_tokens(text: str) -> List[str]:
    base = [t.lower() for t in _TOKEN_RE.findall(_safe_text(text))]
    out: List[str] = []
    for tok in base:
        out.append(tok)
        has_cjk = bool(re.search(r"[\u4e00-\u9fff]", tok))
        has_ascii = bool(re.search(r"[A-Za-z0-9]", tok))
        if has_cjk and has_ascii:
            parts = re.findall(r"[\u4e00-\u9fff]{2,}|[A-Za-z0-9]{2,}", tok)
            for p in parts:
                p2 = p.lower()
                if p2 and p2 not in out:
                    out.append(p2)
    return out


def _is_source_lookup_query(query: str) -> bool:
    q = _safe_text(query).lower()
    if not q:
        return False
    flags = ["source", "where", "from", "which file", "citation"]
    return any(f in q for f in flags)

def _extract_lookup_target(query: str) -> str:
    q = _safe_text(query)
    if not q:
        return ""
    cleaned = q
    for token in ["source", "where", "from", "which file", "citation", "?"]:
        cleaned = cleaned.replace(token, " ")
    cleaned = re.sub(r"\s+", " ", cleaned).strip(" :;,.!?")
    return cleaned



def _keyword_score(tokens: List[str], content: str) -> float:
    if not tokens:
        return 0.0
    c = (content or "").lower()
    return float(sum(c.count(t) for t in tokens))


def _query_match_score(tokens: List[str], content: str) -> float:
    if not tokens:
        return 0.5
    text = (content or "").lower()
    uniq = list(dict.fromkeys([t for t in tokens if t]))
    if not uniq:
        return 0.5
    hit = [t for t in uniq if t in text]
    coverage = len(hit) / max(1, len(uniq))
    freq_raw = sum(text.count(t) for t in uniq)
    freq_norm = min(1.0, freq_raw / max(1.0, len(uniq) * 2.0))
    return min(1.0, max(0.0, (coverage * 0.7) + (freq_norm * 0.3)))


def _content_quality_score(content: str) -> float:
    t = _safe_text(content)
    if not t:
        return 0.0
    low = t.lower()
    readable = len(re.findall(r"[\u4e00-\u9fffA-Za-z0-9]", t))
    symbols = len(re.findall(r"[^ \t\u4e00-\u9fffA-Za-z0-9]", t))
    read_ratio = readable / max(1, readable + symbols)

    score = 0.45 + (0.45 * read_ratio)
    penalties = 0.0
    if "lorem ipsum" in low:
        penalties += 0.35
    if "reallygreatsite.com" in low:
        penalties += 0.25
    if re.search(r"\bwww\.[a-z0-9.-]+\.[a-z]{2,}\b", low):
        penalties += 0.12
    if re.search(r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}", t):
        penalties += 0.12
    if re.search(r"\+?\d[\d\-\s]{7,}\d", t):
        penalties += 0.10
    for w in ("our team", "thank you", "resource page", "placeholder", "text_frame"):
        if w in low:
            penalties += 0.08
    if len(re.findall(r"Page \d+", t)) >= 3 and len(t) < 1200:
        penalties += 0.10
    if symbols > readable * 1.2:
        penalties += 0.20
    return min(1.0, max(0.0, score - penalties))


def _looks_noisy_content(text: str) -> bool:
    t = _safe_text(text)
    if not t:
        return True
    # Exclude JSON/template-like chunks that often dominate retrieval noise.
    punct = sum(1 for ch in t if ch in "{}[]<>:\"\\")
    ratio = punct / max(1, len(t))
    if ratio > 0.08:
        return True
    if "placeholder" in t.lower() and "text_frame" in t.lower():
        return True
    low = t.lower()
    if "lorem ipsum" in low:
        return True
    if "reallygreatsite.com" in low:
        return True
    if "table of content" in low and "introduction to artificial intelligence" in low:
        return True
    return False


def _build_answer_from_evidence(query: str, evidence: List[Dict[str, Any]]) -> str:
    if not evidence:
        return "\u76ee\u524d\u5728\u5df2\u9078\u4f86\u6e90\u4e2d\u627e\u4e0d\u5230\u53ef\u7528\u8b49\u64da\u3002"
    lines = [f"\u4f9d\u64da\u5df2\u9078\u4f86\u6e90\uff0c\u91dd\u5c0d\u300c{query}\u300d\u6574\u7406\u5982\u4e0b\uff1a"]
    for i, ev in enumerate(evidence[:3], start=1):
        src = _safe_text(ev.get("source_title"))
        if not src:
            src = f"\u4f86\u6e90{i}"
        lines.append(f"{i}. \u5df2\u53c3\u8003\uff1a{src}")
    lines.append("\u8a73\u7d30 CHUNK \u53c3\u8003\u8acb\u53c3\u8003\u4e0b\u65b9\u300cCHUNK \u67e5\u8a62\u7d00\u9304\u300d\u5340\u584a\u3002")
    return "\n".join(lines)


def _build_citation_tail(citations: List[Dict[str, Any]]) -> str:
    if not citations:
        return ""
    parts: List[str] = []
    for c in citations:
        ref = _safe_text(c.get("ref")) or "C"
        conf_raw = c.get("confidence")
        try:
            conf = f"{float(conf_raw):.2f}"
        except Exception:
            conf = "--"
        chunk = _to_int(c.get("chunk_index"), 0)
        title = _safe_text(c.get("source_title")) or "\u672a\u77e5\u4f86\u6e90"
        parts.append(f"{ref}({conf})#{chunk} 『{title}』#{chunk}")
    return "來源依據：" + "、".join(parts)



def _llm_to_text(resp: Any) -> str:
    if resp is None:
        return ""
    
    content = getattr(resp, "content", resp)
    
    if isinstance(content, tuple):
        content = content[0]
        
    if isinstance(content, str):
        c_stripped = content.strip()
        if (c_stripped.startswith("('") or c_stripped.startswith('("')) and c_stripped.endswith(")"):
            import ast
            try:
                parsed = ast.literal_eval(c_stripped)
                if isinstance(parsed, tuple):
                    content = parsed[0]
            except Exception:
                pass
                
    if isinstance(content, list):
        parts = []
        for p in content:
            if isinstance(p, str):
                parts.append(p)
            elif isinstance(p, dict) and "text" in p:
                parts.append(str(p["text"]))
        content = "".join(parts)
        
    return _safe_text(content)


def _clean_evidence_for_llm(text: str) -> str:
    t = (text or "").replace("\r\n", "\n").replace("\r", "\n")
    bad_patterns = [
        r"lorem ipsum",
        r"reallygreatsite",
        r"\bwww\.[a-z0-9.-]+\.[a-z]{2,}\b",
        r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}",
        r"\+?\d[\d\-\s]{7,}\d",
        r"\bshape_id\b|\bbbox\b|\bplaceholder\b|\bparagraphs\b|\bruns\b|\bstyle\b",
    ]
    out: List[str] = []
    for raw in t.split("\n"):
        line = _safe_text(raw)
        if not line:
            continue
        low = line.lower()
        if re.match(r"^\s*Page\s*\d+\s*$", line):
            continue
        if any(re.search(p, low, flags=re.IGNORECASE) for p in bad_patterns):
            continue
        if re.search(r"[{}\[\]<>]", line) and len(re.findall(r"[\u4e00-\u9fffA-Za-z0-9]", line)) < 10:
            continue
        out.append(line)
    return "\n".join(out).strip()


def _post_clean_llm_answer(text: str) -> str:
    t = _safe_text(text)
    if not t:
        return ""
    lines: List[str] = []
    for raw in t.split("\n"):
        line = _safe_text(raw)
        if not line:
            continue
        low = line.lower()
        if "lorem ipsum" in low or "reallygreatsite" in low:
            continue
        if re.search(r"\bwww\.[a-z0-9.-]+\.[a-z]{2,}\b", low):
            continue
        if re.search(r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}", line):
            continue
        lines.append(line)
    return "\n".join(lines).strip()


def _select_citations_for_llm(question: str, citations: List[Dict[str, Any]], max_items: int = 4) -> List[Dict[str, Any]]:
    if not citations:
        return []
    tokens = _query_tokens(question)
    scored: List[Dict[str, Any]] = []
    for c in citations:
        excerpt = _clean_evidence_for_llm(_safe_text(c.get("excerpt")))
        if not excerpt:
            continue
        title = _safe_text(c.get("source_title"))
        merged = f"{title}\n{excerpt}"
        match = _query_match_score(tokens, merged)
        quality = _content_quality_score(merged)
        base_conf = float(c.get("confidence") or 0.0)
        final = (match * 0.60) + (quality * 0.20) + (base_conf * 0.20)
        if tokens and match < 0.10:
            continue
        if quality < 0.12:
            continue
        item = dict(c)
        item["_excerpt_clean"] = excerpt
        item["_llm_score"] = final
        scored.append(item)

    if not scored:
        # fallback: keep the top confident citations but still clean excerpt
        fallback = sorted(citations, key=lambda x: float(x.get("confidence") or 0.0), reverse=True)[:max_items]
        out = []
        for c in fallback:
            excerpt = _clean_evidence_for_llm(_safe_text(c.get("excerpt")))
            if not excerpt:
                continue
            item = dict(c)
            item["_excerpt_clean"] = excerpt
            item["_llm_score"] = float(c.get("confidence") or 0.0)
            out.append(item)
        return out

    scored.sort(key=lambda x: float(x.get("_llm_score") or 0.0), reverse=True)
    return scored[:max_items]


def _llm_synthesize_answer(question: str, rule_answer: str, citations: List[Dict[str, Any]]) -> Dict[str, str]:
    if not citations:
        return {"answer": _safe_text(rule_answer), "prompt": "", "warnings": []}
    selected = _select_citations_for_llm(question, citations, max_items=4)
    if not selected:
        return {"answer": "\u76ee\u524d\u5728\u5df2\u9078\u4f86\u6e90\u4e2d\u627e\u4e0d\u5230\u53ef\u7528\u8b49\u64da\u3002", "prompt": "", "warnings": []}
    ev_lines: List[str] = []
    for c in selected:
        ref = _safe_text(c.get("ref")) or "C"
        title = _safe_text(c.get("source_title"))
        idx = _to_int(c.get("chunk_index"), 0)
        excerpt = _safe_text(c.get("_excerpt_clean"))
        if len(excerpt) > 260:
            excerpt = excerpt[:260].rstrip() + "..."
        if not excerpt:
            continue
        ev_lines.append(f"[{ref}] {title}#{idx}\n{excerpt}")
    if not ev_lines:
        return {"answer": _safe_text(rule_answer), "prompt": "", "warnings": []}
    evidence_text = "\n\n".join(ev_lines)
    prompt = f"""
\u4f60\u662f\u516c\u6587\u8207\u77e5\u8b58\u6574\u5408\u52a9\u7406\u3002\u8acb\u4f9d\u64da\u4e0b\u65b9 evidence \u56de\u7b54\uff0c\u4e14\u5fc5\u9808\u9075\u5b88\uff1a
1. \u50c5\u80fd\u4f7f\u7528 evidence \u5167\u5bb9\u4f5c\u7b54\uff0c\u4e0d\u53ef\u81ea\u884c\u88dc\u5145\u672a\u63d0\u4f9b\u7684\u4e8b\u5be6\u3002
2. \u56de\u7b54\u8a9e\u8a00\u4e00\u5f8b\u4f7f\u7528\u7e41\u9ad4\u4e2d\u6587\u3002
3. \u56de\u7b54\u683c\u5f0f\u63a1\u4e00\u554f\u4e00\u7b54\uff0c\u5167\u5bb9\u76f4\u63a5\u3001\u6e05\u695a\u3002
4. \u82e5 evidence \u4e0d\u8db3\u4ee5\u5b8c\u6574\u56de\u7b54\uff0c\u8acb\u660e\u78ba\u8aaa\u660e\u4e0d\u8db3\u8655\uff0c\u907f\u514d\u81c6\u6e2c\u3002
5. \u82e5\u5f15\u7528\u8b49\u64da\uff0c\u8acb\u5728\u53e5\u672b\u4ee5 [C1] \u9019\u7a2e\u683c\u5f0f\u6a19\u793a\u3002

\u554f\u984c\uff1a
{question}

evidence\uff1a
{evidence_text}
""".strip()

    def _rewrite_to_traditional_chinese(llm_obj: Any, text: str) -> str:
        raw = _safe_text(text)
        if not raw:
            return ""
        rewrite_prompt = f"""
隢?銝??批捆?孵神?箇?擃葉??靽???嚗????冽?閮? [C1]?C2]??
撠????臭????雿擗?餈啗?隞亦?擃葉???整?
隢頛詨?孵神敺摰對?銝?憿?隤芣???

?批捆嚗?
{raw}
""".strip()
        try:
            out2 = llm_obj.invoke(rewrite_prompt)
            return _post_clean_llm_answer(_llm_to_text(out2))
        except Exception:
            return ""

    try:
        llm = get_chat_model(temperature=0.1, timeout=90)
        out = llm.invoke(prompt)
        txt = _post_clean_llm_answer(_llm_to_text(out))
        if not is_zh_dominant(txt):
            rewritten = _rewrite_to_traditional_chinese(llm, txt)
            if rewritten:
                txt = rewritten
        fallback_zh = _safe_text(rule_answer)
        if fallback_zh and not is_zh_dominant(fallback_zh):
            rewritten_fb = _rewrite_to_traditional_chinese(llm, fallback_zh)
            if rewritten_fb:
                fallback_zh = rewritten_fb
        final_txt = prefer_traditional_chinese(txt, fallback_zh)
        final_txt = ensure_sentence_citations(final_txt, citations)
        warnings = detect_citation_conflicts(citations)
        return {"answer": final_txt, "prompt": prompt, "warnings": warnings}
    except Exception:
        fallback_txt = ensure_sentence_citations(_safe_text(rule_answer), citations)
        warnings = detect_citation_conflicts(citations)
        return {"answer": fallback_txt, "prompt": prompt, "warnings": warnings}


def _is_http_url(url: str) -> bool:
    try:
        p = urlparse(_safe_text(url))
        return p.scheme in ("http", "https") and bool(p.netloc)
    except Exception:
        return False

# --- VIEWS ---

@require_node("projectnotes")
def index(request: HttpRequest):
    return render(
        request,
        "projectnotes/index.html",
        {"can_manage_projects": _can_manage_projects(request)},
    )

@require_node("projectnotes")
def manage_page(request: HttpRequest):
    if not _can_manage_projects(request):
        return HttpResponseForbidden("system administrator permission required")
    return render(
        request,
        "projectnotes/manage.html",
        {"can_manage_projects": True},
    )

@csrf_exempt
@require_node("projectnotes", api=True)
def api_projects(request: HttpRequest):
    if request.method == "GET":
        rows = Project.objects.all().order_by("-updated_at")[:50]
        data = []
        for p in rows:
            data.append(
                {
                    "id": p.id,
                    "name": p.name,
                    "description": p.description,
                    "strict_source_only": True,
                    "permission_mode": "auth",
                    "source_count": Source.objects.filter(project_id=p.id).count(),
                    "updated_at": p.updated_at.isoformat() if p.updated_at else "",
                }
            )
        return _safe_json_response(
            {"ok": True, "projects": data, "can_manage_projects": _can_manage_projects(request)}
        )

    if request.method == "POST":
        if not _can_manage_projects(request):
            return _api_error("system administrator permission required", error_code="forbidden", status=403)
        body = _read_json_body(request)
        if _is_bad_utf8_request(request):
            return _api_error("request body must be UTF-8", error_code="invalid_encoding", status=400)
        name = _clean_label(body.get("name"), "New Project")
        if not name:
            return _api_error("name is required")
        
        try:
            p = Project.objects.create(
                name=name,
                description=_safe_text(body.get("description")),
                created_by=_current_user_id(request),
            )
        except ProgrammingError as exc:
            return _db_write_error_response(exc, fallback="project create failed")
        _log_activity(
            project_id=p.id,
            action="project_create",
            user_id=_current_user_id(request),
            target_type="project",
            target_id=p.id,
            detail={"name": p.name},
        )
        return _safe_json_response({"ok": True, "project": {"id": p.id, "name": p.name}})
    
    return _api_error("method not allowed", status=405)

@csrf_exempt
@require_node("projectnotes", api=True)
def api_sources(request: HttpRequest):
    if request.method == "GET":
        project_id = _to_int(request.GET.get("project_id"))
        if project_id <= 0:
            return _api_error("project_id required")
        
        rows = Source.objects.filter(project_id=project_id).order_by("-created_at")
        data = []
        for s in rows:
            latest_doc = Document.objects.filter(source_id=s.id).order_by("-id").first()
            latest_ver = (
                DocumentVersion.objects.filter(document_id=latest_doc.id).order_by("-version_number").first()
                if latest_doc
                else None
            )
            ref_url = latest_doc.path if (latest_doc and _is_http_url(latest_doc.path)) else ""
            file_name = ""
            if latest_doc and latest_doc.path and not _is_http_url(latest_doc.path):
                file_name = latest_doc.path
            chunk_count = DocumentChunk.objects.filter(document_version_id=latest_ver.id).count() if latest_ver else 0
            data.append(
                {
                    "id": s.id,
                    "project_id": s.project_id,
                    "title": s.name,
                    "source_type": s.source_type,
                    "reference_url": ref_url,
                    "file_name": file_name,
                    "snapshot_no": latest_ver.version_number if latest_ver else 1,
                    "source_version": f"v{latest_ver.version_number}" if latest_ver else "v1",
                    "is_enabled": True,
                    "chunk_count": chunk_count,
                    "updated_at": s.updated_at.isoformat() if s.updated_at else "",
                }
            )
        return _safe_json_response({"ok": True, "sources": data})

    if request.method == "POST":
        if not _can_manage_projects(request):
            return _api_error("system administrator permission required", error_code="forbidden", status=403)
        project_id = _to_int(request.POST.get("project_id"))
        title = _clean_label(request.POST.get("title"), "Uploaded Source", 240)
        upload = request.FILES.get("file")
        
        if project_id <= 0 or not title or not upload:
            return _api_error("project_id, title, and file are required")
        
        # Read file and preprocess into readable plain text for RAG.
        detected_encoding = "utf-8"
        is_non_utf8 = False
        file_name_lower = upload.name.lower()
        if file_name_lower.endswith(".pdf"):
            try:
                from webapps.pdf.views import _extract_pdf_text_auto
                raw_text, _ = _extract_pdf_text_auto(upload)
                detected_encoding = "extracted"
            except Exception as e:
                return _api_error(f"PDF extraction failed: {e}", status=400)
        elif file_name_lower.endswith(".docx"):
            try:
                import docx
                d = docx.Document(upload)
                raw_text = "\n".join([p.text for p in d.paragraphs])
                detected_encoding = "extracted"
            except Exception as e:
                return _api_error(f"DOCX extraction failed: {e}", status=400)
        elif file_name_lower.endswith((".xlsx", ".xls")):
            try:
                import pandas as pd
                dfs = pd.read_excel(upload, sheet_name=None)
                parts = []
                for sheet, df in dfs.items():
                    parts.append(f"--- Sheet: {sheet} ---")
                    parts.append(df.to_csv(index=False, sep='\t'))
                raw_text = "\n".join(parts)
                detected_encoding = "extracted"
            except Exception as e:
                return _api_error(f"Excel extraction failed: {e}", status=400)
        elif file_name_lower.endswith(".doc"):
            try:
                import tempfile
                import os
                import win32com.client
                import pythoncom
                
                upload.seek(0)
                temp = tempfile.NamedTemporaryFile(delete=False, suffix=".doc")
                temp.write(upload.read())
                temp.close()
                
                pythoncom.CoInitialize()
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                w_doc = word.Documents.Open(os.path.abspath(temp.name))
                raw_text = w_doc.Content.Text
                w_doc.Close()
                word.Quit()
                os.unlink(temp.name)
            except Exception as e:
                # ????????????????????Word COM ???????????????????湔?????????????
                upload.seek(0)
                try:
                    raw_text, detected_encoding = _decode_text_bytes_best_effort(upload.read())
                    is_non_utf8 = detected_encoding not in ("utf-8", "utf-8-sig")
                except ValueError as de:
                    return _api_error(f"DOC fallback decode failed: {de}", status=400)
        elif file_name_lower.endswith(".pptx"):
            try:
                import pptx
                prs = pptx.Presentation(upload)
                parts = []
                for i, slide in enumerate(prs.slides, start=1):
                    parts.append(f"--- Slide {i} ---")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            parts.append(shape.text)
                raw_text = "\n".join(parts)
                detected_encoding = "extracted"
            except Exception as e:
                return _api_error(f"PPTX extraction failed: {e}", status=400)
        elif file_name_lower.endswith(".ppt"):
            try:
                import tempfile
                import os
                import win32com.client
                import pythoncom
                
                upload.seek(0)
                temp = tempfile.NamedTemporaryFile(delete=False, suffix=".ppt")
                temp.write(upload.read())
                temp.close()
                
                pythoncom.CoInitialize()
                ppt_app = win32com.client.Dispatch("PowerPoint.Application")
                
                # PowerPoint COM params: FileName, ReadOnly, Untitled, WithWindow
                p_doc = ppt_app.Presentations.Open(os.path.abspath(temp.name), True, False, False)
                
                text_parts = []
                for i, slide in enumerate(p_doc.Slides, start=1):
                    text_parts.append(f"--- Slide {i} ---")
                    for shape in slide.Shapes:
                        if shape.HasTextFrame and shape.TextFrame.HasText:
                            text_parts.append(shape.TextFrame.TextRange.Text)
                raw_text = "\n".join(text_parts)
                
                p_doc.Close()
                # If there are no other open presentations, quit PowerPoint
                if ppt_app.Presentations.Count == 0:
                    ppt_app.Quit()
                os.unlink(temp.name)
            except Exception as e:
                upload.seek(0)
                try:
                    raw_text, detected_encoding = _decode_text_bytes_best_effort(upload.read())
                    is_non_utf8 = detected_encoding not in ("utf-8", "utf-8-sig")
                except ValueError as de:
                    return _api_error(f"PPT fallback decode failed: {de}", status=400)
        else:
            try:
                raw_text, detected_encoding = _decode_text_bytes_best_effort(upload.read())
                is_non_utf8 = detected_encoding not in ("utf-8", "utf-8-sig")
            except ValueError as de:
                return _api_error(f"Text decode failed: {de}", status=400)
            
        # Ensure no NUL bytes before database insertion
        raw_text = raw_text.replace("\x00", "").replace("\u0000", "")
        
        cleaned_text = _preprocess_rag_text(raw_text)
        if not cleaned_text:
            return _api_error("empty or unreadable content after preprocessing", status=400)
        
        # Create records
        try:
            source = Source.objects.create(project_id=project_id, name=title, source_type="text")
            doc = Document.objects.create(source=source, title=title, path=upload.name)
            doc_version = DocumentVersion.objects.create(
                document=doc,
                version_number=1,
                raw_text=cleaned_text,
                uploaded_by=_current_user_id(request),
            )
        except ProgrammingError as exc:
            return _db_write_error_response(exc, fallback="source create failed")
        
        # Chunking & embedding
        chunks_text = _build_chunks(cleaned_text)
        chunk_objs = []
        for i, text_seg in enumerate(chunks_text):
            emb = _mock_embedding(text_seg)
            chunk_objs.append(DocumentChunk(
                document_version=doc_version,
                chunk_index=i,
                token_count=len(text_seg),
                content=text_seg,
                embedding=emb
            ))
        try:
            DocumentChunk.objects.bulk_create(chunk_objs)
        except ProgrammingError as exc:
            return _db_write_error_response(exc, fallback="chunk create failed")
        _log_activity(
            project_id=project_id,
            action="source_upload",
            user_id=_current_user_id(request),
                target_type="source",
                target_id=source.id,
                detail={
                "source_title": _safe_text(getattr(source, "name", "")) or title,
                    "file_name": upload.name,
                    "chunk_count": len(chunks_text),
                    "detected_encoding": detected_encoding,
                    "is_non_utf8": is_non_utf8,
                },
        )
        
        return _safe_json_response(
            {
                "ok": True,
                "source": {"id": source.id, "chunk_count": len(chunks_text)},
                "detected_encoding": detected_encoding,
                "is_non_utf8": is_non_utf8,
            }
        )
    
    return _api_error("method not allowed", status=405)

@csrf_exempt
@require_node("projectnotes", api=True)
def api_source_delete(request: HttpRequest, source_id: int):
    if request.method not in ("DELETE", "POST"):
        return _api_error("method not allowed", status=405)
    if not _can_manage_projects(request):
        return _api_error("system administrator permission required", error_code="forbidden", status=403)

    sid = _to_int(source_id, 0)
    if sid <= 0:
        return _api_error("source_id required", error_code="missing_source_id", status=400)

    source = Source.objects.filter(id=sid).first()
    if not source:
        return _api_error("source not found", error_code="source_not_found", status=404)

    doc_ids = list(Document.objects.filter(source_id=sid).values_list("id", flat=True))
    version_ids = list(DocumentVersion.objects.filter(document_id__in=doc_ids).values_list("id", flat=True)) if doc_ids else []
    chunk_count = DocumentChunk.objects.filter(document_version_id__in=version_ids).count() if version_ids else 0

    out = {
        "source_id": sid,
        "source_title": _safe_text(source.name),
        "project_id": source.project_id,
        "deleted_documents": len(doc_ids),
        "deleted_versions": len(version_ids),
        "deleted_chunks": chunk_count,
    }

    try:
        source.delete()
    except ProgrammingError as exc:
        return _db_write_error_response(exc, fallback="source delete failed")
    _log_activity(
        project_id=out["project_id"],
        action="source_delete",
        user_id=_current_user_id(request),
        target_type="source",
        target_id=sid,
        detail=out,
    )
    return _safe_json_response({"ok": True, "deleted": out})

@csrf_exempt
@require_node("projectnotes", api=True)
def api_chat(request: HttpRequest):
    if request.method != "POST":
        return _api_error("method not allowed", status=405)

    t0 = time.perf_counter()
    body = _read_json_body(request)
    if _is_bad_utf8_request(request):
        return _api_error("request body must be UTF-8", error_code="invalid_encoding", status=400)
    project_id = _to_int(body.get("project_id"))
    query = _safe_text(body.get("question")) or _safe_text(body.get("query"))
    retrieval_query = rewrite_query_for_retrieval(query)
    conversation_id = _to_int(body.get("conversation_id"))
    llm_synthesis = bool(body.get("llm_synthesis"))
    selected_source_ids = body.get("selected_source_ids") if isinstance(body.get("selected_source_ids"), list) else []
    selected_source_ids = [_to_int(x, 0) for x in selected_source_ids]
    selected_source_ids = [x for x in selected_source_ids if x > 0]
    
    if project_id <= 0 or not query:
        return _api_error("project_id and query required")
        
    query_vector = _mock_embedding(retrieval_query or query)
    conv = None
    if conversation_id > 0:
        conv = Conversation.objects.filter(id=conversation_id, project_id=project_id).first()
    if not conv:
        try:
            conv = Conversation.objects.create(
                project_id=project_id,
                title="New Chat",
                created_by=_current_user_id(request),
            )
        except ProgrammingError as exc:
            return _db_write_error_response(exc, fallback="conversation create failed")
    
    qs = DocumentChunk.objects.filter(
        document_version__document__source__project_id=project_id
    ).select_related("document_version__document__source")
    if selected_source_ids:
        qs = qs.filter(document_version__document__source_id__in=selected_source_ids)

    # Source lookup mode: when user asks "which file / source", use exact phrase hits first.
    lookup_target = _extract_lookup_target(query)
    if _is_source_lookup_query(query) and lookup_target:
        exact_hits = []
        for c in qs.order_by("document_version__document__source_id", "chunk_index")[:500]:
            content = _safe_text(c.content)
            if not content:
                continue
            if lookup_target.lower() in content.lower():
                exact_hits.append(c)
            if len(exact_hits) >= 8:
                break

        if exact_hits:
            citations = []
            evidences = []
            lines = [f"Found exact evidence for: {lookup_target}"]
            for idx, c in enumerate(exact_hits, start=1):
                title = _safe_text(c.document_version.document.title) or f"source_{c.document_version.document.source_id}"
                excerpt = _safe_text(c.content)
                if len(excerpt) > 140:
                    excerpt = excerpt[:140].rstrip() + "..."
                lines.append(f"{idx}. {title} / chunk {c.chunk_index}")
                evidences.append(
                    {
                        "chunk_id": c.id,
                        "content": excerpt,
                        "document": title,
                    }
                )
                src_path = _safe_text(c.document_version.document.path)
                citations.append(
                    {
                        "ref": f"C{idx}",
                        "source_id": c.document_version.document.source_id,
                        "source_title": title,
                        "chunk_index": c.chunk_index,
                        "confidence": 0.9,
                        "source_url": src_path if _is_http_url(src_path) else "",
                        "excerpt": _safe_text(c.content),
                    }
                )

            rule_answer = "\n".join(lines)
            llm_payload = _llm_synthesize_answer(query, rule_answer, citations) if llm_synthesis else {"answer": "", "prompt": "", "warnings": []}
            llm_answer = _safe_text(llm_payload.get("answer"))
            llm_prompt_preview = _safe_text(llm_payload.get("prompt"))
            citation_warnings = llm_payload.get("warnings") if isinstance(llm_payload.get("warnings"), list) else []
            ai_answer = rule_answer
            citation_tail = _build_citation_tail(citations)
            try:
                conversation_title = _maybe_update_conversation_title(conv, query)
                Message.objects.create(
                    conversation=conv,
                    sender_type="user",
                    sender_id=_current_user_id(request),
                    content=query,
                )
                assistant_msg = Message.objects.create(
                    conversation=conv,
                    sender_type="assistant",
                    sender_id="system",
                    content=ai_answer,
                )
                _persist_message_citations(assistant_msg, citations)
            except ProgrammingError as exc:
                return _db_write_error_response(exc, fallback="message create failed")
            latency_ms = round((time.perf_counter() - t0) * 1000.0, 2)
            _log_activity(
                project_id=project_id,
                action="chat_query",
                user_id=_current_user_id(request),
                target_type="conversation",
                target_id=conv.id,
                detail={
                    "conversation_id": conv.id,
                    "query": query,
                    "retrieval_query": retrieval_query,
                    "lookup_mode": True,
                    "selected_source_ids": selected_source_ids,
                    "citation_count": len(citations),
                    "evidence_count": len(evidences),
                    "llm_synthesis": llm_synthesis,
                    "llm_synthesis_used": bool(llm_answer),
                    "status": "ok" if citations else "insufficient",
                    "latency_ms": latency_ms,
                },
            )
            return _safe_json_response({
                "ok": True,
                "conversation_id": conv.id,
                "conversation_title": conversation_title or _safe_text(conv.title),
                "retrieval_query": retrieval_query,
                "answer": ai_answer,
                "rule_answer": rule_answer,
                "llm_answer": llm_answer,
                "llm_prompt_preview": llm_prompt_preview,
                "citation_warnings": citation_warnings,
                "citation_tail": citation_tail,
                "confidence": 0.9,
                "llm_synthesis_used": bool(llm_answer),
                "citations": citations,
                "turn": {
                    "query": query,
                    "answer": ai_answer,
                    "evidence": evidences
                }
            })

    # Hybrid recall v1: dense(vector) + sparse(lexical contains), then rerank.
    tokens = _query_tokens(retrieval_query or query)
    top_chunks_dense = list(qs.order_by(L2Distance("embedding", query_vector))[:40])
    sparse_terms = build_sparse_terms(tokens, max_terms=6)
    top_chunks_sparse: List[DocumentChunk] = []
    if sparse_terms:
        sparse_q = Q()
        for term in sparse_terms:
            sparse_q |= Q(content__icontains=term)
        top_chunks_sparse = list(qs.filter(sparse_q).order_by("id")[:40])

    top_chunks: List[DocumentChunk] = []
    seen_chunk_ids = set()
    for c in top_chunks_dense + top_chunks_sparse:
        cid = int(getattr(c, "id", 0) or 0)
        if cid <= 0 or cid in seen_chunk_ids:
            continue
        seen_chunk_ids.add(cid)
        top_chunks.append(c)
        if len(top_chunks) >= 60:
            break

    def_query = is_definition_query(retrieval_query or query)
    ranked: List[Dict[str, Any]] = []
    max_kscore = 0.0
    for c in top_chunks:
        content = _safe_text(c.content)
        if not content:
            continue
        noisy = _looks_noisy_content(content)
        kscore = _keyword_score(tokens, content)
        source_title = _safe_text(c.document_version.document.title)
        match_score = _query_match_score(tokens, content)
        quality_score = _content_quality_score(content)
        generic_penalty = generic_source_penalty(query, source_title, content)
        def_boost = definition_chunk_boost(content) if def_query else 0.0
        if kscore > max_kscore:
            max_kscore = kscore
        # Hybrid ranking: lexical + semantic-match + quality + intent boost - penalties.
        score = (kscore * 1.6) + (match_score * 1.2) + (quality_score * 0.6)
        score += (0.20 if not noisy else -0.35)
        score += def_boost
        score -= generic_penalty
        ranked.append(
            {
                "score": score,
                "chunk": c,
                "excerpt": content,
                "source_title": source_title,
                "noisy": noisy,
                "kscore": kscore,
                "match_score": match_score,
                "quality_score": quality_score,
                "generic_penalty": generic_penalty,
            }
        )
    if tokens and max_kscore <= 0:
        # Fallback: keep a smaller rerank pool when lexical signal is weak.
        picked = rerank_candidates(retrieval_query or query, ranked, top_k=3)
    else:
        picked = rerank_candidates(retrieval_query or query, ranked, top_k=5)

    citations = []
    evidences = []
    kept_items: List[Dict[str, Any]] = []
    kept_confs: List[float] = []
    min_citation_conf = 0.15
    # Confidence v2: query-match + content-quality + rank decay - noise penalty.
    # Goal: lower confidence on template/ad text and repeated boilerplate.

    for rank, item in enumerate(picked, start=1):
        c = item["chunk"]
        excerpt = _safe_text(item["excerpt"])
        match_score = float(item.get("match_score", 0.0))
        quality_score = float(item.get("quality_score", 0.0))
        rank_factor = max(0.40, 1.0 - (rank - 1) * 0.14)
        noisy_penalty = 0.18 if bool(item.get("noisy")) else 0.0
        generic_penalty = float(item.get("generic_penalty") or 0.0) * 0.50
        zero_match_penalty = 0.18 if (tokens and float(item.get("kscore", 0.0)) <= 0.0) else 0.0
        cit_conf = (match_score * 0.55) + (quality_score * 0.30) + (rank_factor * 0.15)
        cit_conf = cit_conf - noisy_penalty - zero_match_penalty - generic_penalty
        cit_conf = min(0.99, max(0.0, cit_conf))
        if cit_conf < min_citation_conf:
            continue
        kept_items.append(item)
        kept_confs.append(cit_conf)
        evidences.append(
            {
                "chunk_id": c.id,
                "content": excerpt,
                "document": item["source_title"],
            }
        )
    for idx, item in enumerate(kept_items, start=1):
        c = item["chunk"]
        excerpt = _safe_text(item["excerpt"])
        cit_conf = kept_confs[idx - 1]
        src_path = _safe_text(c.document_version.document.path)
        citations.append(
            {
                "ref": f"C{idx}",
                "source_id": c.document_version.document.source_id,
                "source_title": item["source_title"],
                "chunk_index": c.chunk_index,
                "confidence": round(cit_conf, 2),
                "source_url": src_path if _is_http_url(src_path) else "",
                "excerpt": excerpt,
            }
        )

    rule_answer = _build_answer_from_evidence(query, kept_items)
    llm_payload = _llm_synthesize_answer(query, rule_answer, citations) if llm_synthesis else {"answer": "", "prompt": "", "warnings": []}
    llm_answer = _safe_text(llm_payload.get("answer"))
    llm_prompt_preview = _safe_text(llm_payload.get("prompt"))
    citation_warnings = llm_payload.get("warnings") if isinstance(llm_payload.get("warnings"), list) else []
    ai_answer = rule_answer
    citation_tail = _build_citation_tail(citations)
    
    try:
        conversation_title = _maybe_update_conversation_title(conv, query)
        Message.objects.create(
            conversation=conv,
            sender_type="user",
            sender_id=_current_user_id(request),
            content=query,
        )
        assistant_msg = Message.objects.create(
            conversation=conv,
            sender_type="assistant",
            sender_id="system",
            content=ai_answer,
        )
        _persist_message_citations(assistant_msg, citations)
    except ProgrammingError as exc:
        return _db_write_error_response(exc, fallback="message create failed")
    latency_ms = round((time.perf_counter() - t0) * 1000.0, 2)
    _log_activity(
        project_id=project_id,
        action="chat_query",
        user_id=_current_user_id(request),
        target_type="conversation",
        target_id=conv.id,
        detail={
            "conversation_id": conv.id,
            "query": query,
            "retrieval_query": retrieval_query,
            "lookup_mode": False,
            "selected_source_ids": selected_source_ids,
            "citation_count": len(citations),
            "evidence_count": len(evidences),
            "llm_synthesis": llm_synthesis,
            "llm_synthesis_used": bool(llm_answer),
            "status": "ok" if kept_items else "insufficient",
            "latency_ms": latency_ms,
        },
    )

    return _safe_json_response({
        "ok": True,
        "conversation_id": conv.id,
        "conversation_title": conversation_title or _safe_text(conv.title),
        "retrieval_query": retrieval_query,
        "answer": ai_answer,
        "rule_answer": rule_answer,
        "llm_answer": llm_answer,
        "llm_prompt_preview": llm_prompt_preview,
        "citation_warnings": citation_warnings,
        "citation_tail": citation_tail,
        "confidence": 0.6 if kept_items else 0.0,
        "llm_synthesis_used": bool(llm_answer),
        "citations": citations,
        "turn": {
            "query": query,
            "answer": ai_answer,
            "evidence": evidences
        }
    })

@csrf_exempt
@require_node("projectnotes", api=True)
def api_conversations(request):
    if request.method == "GET":
        project_id = _to_int(request.GET.get("project_id"))
        if project_id <= 0:
            return _api_error("project_id required", status=400)
        rows = Conversation.objects.filter(project_id=project_id).order_by("-updated_at", "-id")[:100]
        data = []
        for c in rows:
            turn_count = Message.objects.filter(conversation_id=c.id, sender_type="assistant").count()
            safe_title = _clean_label(c.title, f"\u5c0d\u8a71 {c.id}", 200)
            if safe_title != (c.title or ""):
                c.title = safe_title
                c.save(update_fields=["title", "updated_at"])
            data.append(
                {
                    "id": c.id,
                    "title": safe_title,
                    "turn_count": turn_count,
                    "updated_at": c.updated_at.isoformat() if c.updated_at else "",
                }
            )
        return _safe_json_response({"ok": True, "conversations": data})

    if request.method == "POST":
        body = _read_json_body(request)
        if _is_bad_utf8_request(request):
            return _api_error("request body must be UTF-8", error_code="invalid_encoding", status=400)
        project_id = _to_int(body.get("project_id"))
        if project_id <= 0:
            return _api_error("project_id required", status=400)
        title = _clean_label(body.get("title"), "\u65b0\u5c0d\u8a71", 200)
        try:
            conv = Conversation.objects.create(
                project_id=project_id,
                title=title,
                created_by=_current_user_id(request),
            )
        except ProgrammingError as exc:
            return _db_write_error_response(exc, fallback="conversation create failed")
        _log_activity(
            project_id=project_id,
            action="conversation_create",
            user_id=_current_user_id(request),
            target_type="conversation",
            target_id=conv.id,
            detail={"title": conv.title},
        )
        return _safe_json_response({"ok": True, "conversation": {"id": conv.id, "title": conv.title}})

    if request.method in ("PATCH", "PUT"):
        body = _read_json_body(request)
        if _is_bad_utf8_request(request):
            return _api_error("request body must be UTF-8", error_code="invalid_encoding", status=400)
        conversation_id = _to_int(body.get("conversation_id"))
        if conversation_id <= 0:
            return _api_error("conversation_id required", status=400)
        conv = Conversation.objects.filter(id=conversation_id).first()
        if not conv:
            return _api_error("conversation not found", status=404)
        title = _clean_label(body.get("title"), "", 200)
        if not title:
            return _api_error("title required", status=400)
        try:
            conv.title = title
            conv.save(update_fields=["title", "updated_at"])
        except ProgrammingError as exc:
            return _db_write_error_response(exc, fallback="conversation rename failed")
        _log_activity(
            project_id=conv.project_id,
            action="conversation_rename",
            user_id=_current_user_id(request),
            target_type="conversation",
            target_id=conv.id,
            detail={"title": conv.title},
        )
        return _safe_json_response({"ok": True, "conversation": {"id": conv.id, "title": conv.title}})

    if request.method == "DELETE":
        body = _read_json_body(request)
        if _is_bad_utf8_request(request):
            return _api_error("request body must be UTF-8", error_code="invalid_encoding", status=400)
        conversation_id = _to_int(body.get("conversation_id"))
        if conversation_id <= 0:
            return _api_error("conversation_id required", status=400)
        conv = Conversation.objects.filter(id=conversation_id).first()
        if not conv:
            return _api_error("conversation not found", status=404)
        out = {
            "id": conv.id,
            "project_id": conv.project_id,
            "title": _safe_text(conv.title),
        }
        try:
            conv.delete()
        except ProgrammingError as exc:
            return _db_write_error_response(exc, fallback="conversation delete failed")
        _log_activity(
            project_id=out["project_id"],
            action="conversation_delete",
            user_id=_current_user_id(request),
            target_type="conversation",
            target_id=out["id"],
            detail={"title": out["title"]},
        )
        return _safe_json_response({"ok": True, "deleted": out})

    return _api_error("method not allowed", status=405)

def api_digests(request):
    return _api_error("Not implemented")

@require_node("projectnotes", api=True)
def api_messages(request):
    if request.method != "GET":
        return _api_error("method not allowed", status=405)
    conversation_id = _to_int(request.GET.get("conversation_id"))
    if conversation_id <= 0:
        return _api_error("conversation_id required", status=400)

    conv = Conversation.objects.filter(id=conversation_id).select_related("project").first()
    if not conv:
        return _api_error("conversation not found", status=404)

    rows = (
        Message.objects
        .filter(conversation_id=conversation_id)
        .order_by("created_at", "id")
        .prefetch_related("citations__document_chunk__document_version__document__source")
    )
    messages = []
    for row in rows:
        citations = []
        for cite in row.citations.all():
            chunk = getattr(cite, "document_chunk", None)
            if not chunk or not getattr(chunk, "document_version", None):
                continue
            doc = chunk.document_version.document
            src = doc.source
            citations.append(
                {
                    "ref": _safe_text(cite.citation_text) or "C",
                    "source_id": src.id,
                    "source_title": _safe_text(src.name),
                    "chunk_index": chunk.chunk_index,
                    "excerpt": _safe_text(chunk.content),
                }
            )
        messages.append(
            {
                "id": row.id,
                "conversation_id": row.conversation_id,
                "sender_type": _safe_text(row.sender_type),
                "sender_id": _safe_text(row.sender_id),
                "content": _safe_text(row.content),
                "created_at": row.created_at.isoformat() if row.created_at else "",
                "citations": citations,
            }
        )
    return _safe_json_response({"ok": True, "conversation": {"id": conv.id, "project_id": conv.project_id, "title": _safe_text(conv.title)}, "messages": messages})

def audit_page(request):
    return _api_error("Not implemented")

def metrics_page(request):
    return _api_error("Not implemented")

def api_search(request):
    return _api_error("Not implemented")

def api_comments(request):
    return _api_error("Not implemented")

@require_node("projectnotes", api=True)
def api_audit_logs(request):
    if request.method != "GET":
        return _api_error("method not allowed", status=405)

    project_id = _to_int(request.GET.get("project_id"))
    user_id = _safe_text(request.GET.get("user_id"))
    limit = max(1, min(200, _to_int(request.GET.get("limit"), 50)))

    qs = ActivityLog.objects.all().select_related("project").order_by("-created_at", "-id")
    if project_id > 0:
        qs = qs.filter(project_id=project_id)
    if user_id:
        qs = qs.filter(user_id=user_id)

    rows = []
    for row in qs[:limit]:
        detail = row.detail_json if isinstance(row.detail_json, dict) else {}
        status = _safe_text(detail.get("status")) or "-"
        rows.append(
            {
                "id": row.id,
                "project_id": row.project_id,
                "project_name": _safe_text(getattr(row.project, "name", "")),
                "user_id": _safe_text(row.user_id),
                "action": _safe_text(row.action),
                "status": status,
                "target_type": _safe_text(row.target_type),
                "target_id": row.target_id,
                "conversation_id": detail.get("conversation_id"),
                "detail": detail,
                "created_at": row.created_at.isoformat() if row.created_at else "",
            }
        )

    return _safe_json_response({"ok": True, "rows": rows})

@require_node("projectnotes", api=True)
def api_metrics(request):
    if request.method != "GET":
        return _api_error("method not allowed", status=405)

    project_id = _to_int(request.GET.get("project_id"))
    days = max(1, min(90, _to_int(request.GET.get("days"), 7)))
    cutoff = time.time() - (days * 86400)

    qs = ActivityLog.objects.all().order_by("-created_at")
    if project_id > 0:
        qs = qs.filter(project_id=project_id)

    rows = []
    for row in qs[:5000]:
        created_ts = row.created_at.timestamp() if row.created_at else 0.0
        if created_ts < cutoff:
            continue
        rows.append(row)

    chat_rows = [r for r in rows if _safe_text(r.action) == "chat_query"]
    source_upload_rows = [r for r in rows if _safe_text(r.action) == "source_upload"]
    source_delete_rows = [r for r in rows if _safe_text(r.action) == "source_delete"]
    conv_create_rows = [r for r in rows if _safe_text(r.action) == "conversation_create"]
    citation_click_rows = [r for r in rows if _safe_text(r.action) == "citation_click"]

    insufficient = 0
    latency_vals: List[float] = []
    citation_total = 0
    for row in chat_rows:
        detail = row.detail_json if isinstance(row.detail_json, dict) else {}
        if _safe_text(detail.get("status")).lower() == "insufficient":
            insufficient += 1
        try:
            latency = float(detail.get("latency_ms"))
            if latency >= 0:
                latency_vals.append(latency)
        except Exception:
            pass
        citation_total += max(0, _to_int(detail.get("citation_count"), 0))

    metrics = {
        "usage_count": len(rows),
        "query_count": len(chat_rows),
        "insufficient_count": insufficient,
        "insufficient_rate": round((insufficient / len(chat_rows)), 4) if chat_rows else 0.0,
        "citation_click_count": len(citation_click_rows),
        "citation_click_rate": round((len(citation_click_rows) / len(chat_rows)), 4) if chat_rows else 0.0,
        "avg_latency_ms": round((sum(latency_vals) / len(latency_vals)), 2) if latency_vals else 0.0,
        "source_upload_count": len(source_upload_rows),
        "source_delete_count": len(source_delete_rows),
        "conversation_create_count": len(conv_create_rows),
        "avg_citation_count": round((citation_total / len(chat_rows)), 2) if chat_rows else 0.0,
    }

    return _safe_json_response({"ok": True, "days": days, "metrics": metrics})

@csrf_exempt
@require_node("projectnotes", api=True)
def api_source_versions(request):
    if request.method != "GET":
        return _api_error("method not allowed", status=405)
    project_id = _to_int(request.GET.get("project_id"))
    if project_id <= 0:
        return _api_error("project_id required", status=400)

    rows = Source.objects.filter(project_id=project_id).order_by("name", "created_at", "id")
    grouped: Dict[str, List[Source]] = {}
    for s in rows:
        key = _safe_text(s.name) or f"source_{s.id}"
        grouped.setdefault(key, []).append(s)

    items = []
    for title, srcs in grouped.items():
        versions = []
        for idx, s in enumerate(srcs, start=1):
            versions.append(
                {
                    "id": s.id,
                    "source_version": f"v{idx}",
                    "source_type": s.source_type,
                    "updated_at": s.updated_at.isoformat() if s.updated_at else "",
                }
            )
        latest_id = srcs[-1].id if srcs else 0
        items.append(
            {
                "title": title,
                "latest_id": latest_id,
                "versions": versions,
            }
        )

    return _safe_json_response({"ok": True, "items": items})

@require_node("projectnotes", api=True)
def api_source_content(request, source_id: int):
    if request.method != "GET":
        return _api_error("method not allowed", status=405)
    sid = _to_int(source_id, 0)
    if sid <= 0:
        return _api_error("source_id required", status=400)

    source = Source.objects.filter(id=sid).first()
    if not source:
        return _api_error("source not found", status=404)

    doc = Document.objects.filter(source_id=sid).order_by("-id").first()
    if not doc:
        return _safe_json_response(
            {
                "ok": True,
                "source": {"id": source.id, "title": source.name, "source_type": source.source_type},
                "chunks": [],
            }
        )

    doc_ver = DocumentVersion.objects.filter(document_id=doc.id).order_by("-version_number", "-id").first()
    if not doc_ver:
        return _safe_json_response(
            {
                "ok": True,
                "source": {"id": source.id, "title": source.name, "source_type": source.source_type},
                "chunks": [],
            }
        )

    rows = DocumentChunk.objects.filter(document_version_id=doc_ver.id).order_by("chunk_index")[:12]
    out_chunks = []
    for r in rows:
        out_chunks.append(
            {
                "chunk_index": r.chunk_index,
                "content": _safe_text(r.content),
            }
        )

    return _safe_json_response(
        {
            "ok": True,
            "source": {
                "id": source.id,
                "title": source.name,
                "source_type": source.source_type,
                "document_title": _safe_text(doc.title),
                "path": _safe_text(doc.path),
                "version": f"v{doc_ver.version_number}",
            },
            "chunks": out_chunks,
        }
    )

def api_source_toggle(request, source_id):
    return _api_error("Not implemented")

def api_source_resync(request, source_id):
    return _api_error("Not implemented")

@csrf_exempt
@require_node("projectnotes", api=True)
def api_citation_click(request):
    if request.method != "POST":
        return _api_error("method not allowed", status=405)
    body = _read_json_body(request)
    if _is_bad_utf8_request(request):
        return _api_error("request body must be UTF-8", error_code="invalid_encoding", status=400)

    project_id = _to_int(body.get("project_id"))
    source_id = _to_int(body.get("source_id"))
    chunk_index = _to_int(body.get("chunk_index"), -1)
    conversation_id = _to_int(body.get("conversation_id"))
    ref = _safe_text(body.get("ref"))
    excerpt = _safe_text(body.get("excerpt"))

    if project_id <= 0:
        return _api_error("project_id required", status=400)
    if source_id <= 0:
        return _api_error("source_id required", status=400)

    _log_activity(
        project_id=project_id,
        action="citation_click",
        user_id=_current_user_id(request),
        target_type="source",
        target_id=source_id,
        detail={
            "project_id": project_id,
            "conversation_id": conversation_id if conversation_id > 0 else None,
            "source_id": source_id,
            "chunk_index": chunk_index if chunk_index >= 0 else None,
            "ref": ref,
            "excerpt": excerpt[:300],
            "status": "ok",
        },
    )
    return _safe_json_response({"ok": True})

def api_citation_context(request):
    return _api_error("Not implemented")

@require_node("projectnotes", api=True)
def api_overview(request):
    if request.method != "GET":
        return _api_error("method not allowed", status=405)
    project_id = _to_int(request.GET.get("project_id"))
    if project_id <= 0:
        return _api_error("project_id required", status=400)
    if not Project.objects.filter(id=project_id).exists():
        return _api_error("project not found", status=404)
    return _safe_json_response({"ok": True, "overview": _build_project_overview(project_id)})

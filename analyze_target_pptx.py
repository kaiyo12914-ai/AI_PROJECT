from pptx import Presentation
import os

def analyze_pptx(path):
    if not os.path.exists(path):
        print(f"Error: File {path} not found.")
        return

    try:
        prs = Presentation(path)
        print(f"--- Analysis for: {os.path.basename(path)} ---")
        print(f"Total Layouts: {len(prs.slide_layouts)}")
        
        for i, layout in enumerate(prs.slide_layouts):
            print(f"\nLayout {i} Name: '{layout.name}'")
            placeholders = [p.name for p in layout.placeholders]
            print(f"  Placeholders found: {placeholders}")
            
            # Check for generic text frames that are not placeholders
            text_frames = 0
            for shape in layout.shapes:
                if getattr(shape, "has_text_frame", False):
                    text_frames += 1
            print(f"  Total Text Shapes: {text_frames}")
            
    except Exception as e:
        print(f"Error reading PPTX: {e}")

if __name__ == "__main__":
    target = r"H:\AI\Django\webapps\text2pptx\pptx_templates\Portfolio卡通.pptx"
    analyze_pptx(target)
